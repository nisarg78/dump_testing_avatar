"""
Lambda 02: Document Extractor
=============================
MECH Avatar - BMO Lab Account (939620275271)

Extracts text and structured elements from documents, creates semantic chunks.

Supported Formats:
    PDF, DOCX, DOC, XLSX, XLS, PPTX, PPT, TXT, MD, CSV, JSON, HTML, XML, XSD

S3 Structure:
    Input:  mechavatar-lab-cac1-s3-inbound/documents/{file}
    Output: mechavatar-lab-cac1-s3-inbound/chunks/{document_id}.jsonl

Environment Variables:
    AWS_DEFAULT_REGION              = ca-central-1
    AWS_ACCOUNT_ID                  = 939620275271
    BUCKET_NAME                     = mechavatar-lab-cac1-s3-inbound
    PIPELINE_TABLE                  = mechavatar-lab-cac1-mech-processing-pipeline
    HAIKU_INFERENCE_PROFILE_ARN     = anthropic.claude-3-haiku-20240307-v1:0
    CLAUDE_VISION_MODEL             = anthropic.claude-3-haiku-20240307-v1:0
    ENABLE_LLM_METADATA             = true
    ENABLE_CLAUDE_VISION_OCR        = true

Chunking Settings:
    SECTION_MAX_TOKENS  = 3500   (max tokens per section)
    CHUNK_TARGET_TOKENS = 1500   (target for semantic chunks)
    CHUNK_MAX_TOKENS    = 2500   (hard limit before split)
    CHUNK_OVERLAP       = 400    (overlap between chunks)

Lambda Configuration:
    Memory: 3008 MB | Timeout: 900s | Runtime: Python 3.11
    Layers: Layer2A (pymupdf, python-docx, openpyxl, etc.)
    Trigger: Step Functions after Lambda 01
"""

from __future__ import annotations

import json
import boto3
import os
import tempfile
import logging
import hashlib
import re
import subprocess
import zipfile
import base64
import uuid
import time
from datetime import datetime
from typing import Dict, List, Any, Tuple, Optional
from dataclasses import dataclass, asdict, field
from enum import Enum
from abc import ABC, abstractmethod
from pydantic import BaseModel, ConfigDict, ValidationError
from botocore.exceptions import BotoCoreError, ClientError
import io

# Optional: tiktoken for token counting
try:
    import tiktoken
except Exception:
    tiktoken = None

# Import centralized BMO configuration
import sys
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
try:
    from config import (
        AWS_REGION, AWS_ACCOUNT_ID, BUCKET_NAME, PIPELINE_TABLE,
        BEDROCK_HAIKU_MODEL
    )
except ImportError:
    AWS_REGION = os.environ.get('AWS_DEFAULT_REGION', 'ca-central-1')
    AWS_ACCOUNT_ID = os.environ.get('AWS_ACCOUNT_ID', '939620275271')
    BUCKET_NAME = os.environ.get('BUCKET_NAME', 'mechavatar-lab-cac1-s3-inbound')
    PIPELINE_TABLE = os.environ.get('PIPELINE_TABLE', 'mechavatar-lab-cac1-mech-processing-pipeline')
    BEDROCK_HAIKU_MODEL = os.environ.get('BEDROCK_HAIKU_MODEL', 'anthropic.claude-3-haiku-20240307-v1:0')

# ============================================================================
# LOGGING CONFIGURATION
# ============================================================================

logger = logging.getLogger()
logger.setLevel(logging.INFO)

_correlation_id: Optional[str] = None


def set_correlation_id(request_id: str) -> None:
    global _correlation_id
    _correlation_id = request_id


def log_event(event_type: str, message: str, **kwargs) -> None:
    log_data = {
        'correlation_id': _correlation_id,
        'event': event_type,
        'message': message,
        'timestamp': datetime.utcnow().isoformat(),
        **kwargs
    }
    logger.info(json.dumps(log_data))

# ============================================================================
# CHUNKING CONFIGURATION (env vars for tuning)
# ============================================================================

CHUNK_SIZE = int(os.environ.get('CHUNK_SIZE', 3500))
CHUNK_OVERLAP = int(os.environ.get('CHUNK_OVERLAP', 400))
MIN_CHUNK_SIZE = int(os.environ.get('MIN_CHUNK_SIZE', 150))

# Advanced chunking configuration (token-based)
ENABLE_ADVANCED_CHUNKING = os.environ.get('ENABLE_ADVANCED_CHUNKING', 'true').lower() == 'true'
TARGET_CHUNK_SIZE = int(os.environ.get('TARGET_CHUNK_SIZE', 1500))
MAX_CHUNK_SIZE = int(os.environ.get('MAX_CHUNK_SIZE', 2500))
MIN_TOKEN_CHUNK_SIZE = int(os.environ.get('MIN_TOKEN_CHUNK_SIZE', 200))
CHUNK_OVERLAP_TOKENS = int(os.environ.get('CHUNK_OVERLAP_TOKENS', 200))
TABLE_MAX_TOKENS = int(os.environ.get('TABLE_MAX_TOKENS', 2500))
TABLE_ROWS_PER_CHUNK = int(os.environ.get('TABLE_ROWS_PER_CHUNK', 50))
CODE_TARGET_LINES = int(os.environ.get('CODE_TARGET_LINES', 350))
CODE_OVERLAP_LINES = int(os.environ.get('CODE_OVERLAP_LINES', 50))

# Extractor configuration
OCR_ENABLED = os.environ.get('OCR_ENABLED', 'true').lower() == 'true'
TEXTRACT_ENABLED = os.environ.get('TEXTRACT_ENABLED', 'true').lower() == 'true'
TEMP_DIR = os.environ.get('TEMP_DIR', '/tmp/mech_avatar_extract')

# Processing limits
MAX_FILE_SIZE_MB = int(os.environ.get('MAX_FILE_SIZE_MB', 100))
MAX_PAGES = int(os.environ.get('MAX_PAGES', 500))


# ============================================================================
# MECH DOMAIN PATTERNS AND CRITICAL LITERALS (Layer 1 & 2 Integration)
# ============================================================================

# Mainframe-specific patterns for entity extraction
MECH_PATTERNS = {
    # BDD Table names (MYxxxx, BDDxxxx patterns)
    'bdd_tables': [
        r'\bMY[A-Z]{2,6}\d?\b',        # MYTP03, MYBFTA, MYASSH, MYRMKS
        r'\bBDD[A-Z0-9]{2,8}\b',       # BDDOPTR, BDDMSPA, BDDOPT6
        r'\bFSFL[A-Z]{2,4}\b',         # FSFLRTB, FSFLLAN
    ],
    
    # Field names (uppercase with numbers, 4-10 chars)
    'field_names': [
        r'\bDABI[A-Z]{2,6}\b',         # DABISEC, DABITAP, DABILGT
        r'\bOPTR[A-Z]{2,5}\b',         # OPTRTYP
        r'\bOPDM[A-Z]{2,5}\b',         # OPDMLGT
        r'\bSPTL[A-Z]{2,6}\b',         # SPTLNCT
        r'\bJIEN[A-Z]{2,4}\b',         # JIENAMT, JIENPOS, JIENGLA
    ],
    
    # SHF Segments
    'shf_segments': r'\bSHF0([134])\b',
    
    # XML Tags - for RL-03 questions
    'xml_tags': r'<([A-Za-z][A-Za-z0-9_-]*)(?:\s[^>]*)?>',
    
    # Hex data patterns
    'hex_data': [
        r"X'[0-9A-Fa-f]{2,}'",         # X'0A', X'C6'
        r'\b0x[0-9A-Fa-f]{2,}\b',      # 0x0A
    ],
    
    # JCL patterns
    'jcl': [
        r'^//[A-Z][A-Z0-9]*\s+(?:JOB|EXEC|DD|PROC)',
        r'\bDSNUTILB\b',
        r'\bIEBCOPY\b',
        r'\bDFSSORT\b',
    ],
    
    # Assembly/Macro patterns
    'macros': [
        r'\b[A-Z]{4,8}\s+(?:PROC|ENDP|CSECT|DSECT)\b',
        r'\bSTDPC\s+TYPE=[A-Z]+\b',
        r'\bBMOXMLB\s+TYPE=[A-Z]+\b',
    ],
    
    # Program names
    'programs': [
        r'\bM[A-Z]{4,7}\b',            # MSHFPA, MXDBMI, MPOSDA
        r'\bTLG\d{3}\b',               # TLG100
    ],
    
    # FSF Fields - for statement testing
    'fsf_fields': r'\bFSFL([A-Z]{2,4})\b',
    
    # OPTRTYP values - for back item mappings
    'optrtyp_values': r"(?:OPTRTYP|type)\s*[=:]\s*(?:C'([0-9A-F])'|X'([0-9A-F]{2})'|(\d+))",
    
    # Call chain patterns - generic STDPC macros
    'call_chain': r'STDPC\s+TYPE=(?:LINK|RETURN)\s+(\w+)',
}


class MechDomain(Enum):
    """MECH documentation domains for routing."""
    SHF_PROCESSING = "shf_processing"
    TXNLOG_OPERATIONS = "txnlog_operations"
    RL03_XML = "rl03_xml"
    BACK_ITEM_MAPPING = "back_item_mapping"
    STATEMENT_TESTING = "statement_testing"
    XXA_TRANSACTIONS = "xxa_transactions"
    GENERAL = "general"


# Domain classification patterns - generic regex-based detection
DOMAIN_PATTERNS = {
    MechDomain.SHF_PROCESSING: [
        r'\bSHF[-_]?\d{2}\b',           # SHF segments (SHF03, SHF04, SHF14, etc.)
        r'\bMSHFP[A-Z]\b',              # SHF processor macros
        r'\bformatted\s+segment\b',     # SHF terminology
    ],
    MechDomain.TXNLOG_OPERATIONS: [
        r'\bMECH_TXNLOG\b', r'\bTXNLOG\b',
        r'\bTLG\d{3}\b',                # Job names
        r'\bpartition\s+rotat',         # Partition rotation
    ],
    MechDomain.RL03_XML: [
        r'\bRL[-\s]?03\b',              # RL03 reporting
        r'\bRevenue\s+Qu[éeè]bec\b',    # Quebec revenue
        r'\bXML\s+(?:tag|header|schema)\b',
    ],
    MechDomain.BACK_ITEM_MAPPING: [
        r'\bDABI[A-Z]+\b',              # DABI fields
        r'\bOPTR(?:TYP)?\b',            # OPTR fields
        r'\bBDDOP[A-Z0-9]+\b',          # BDD mapping fields
    ],
    MechDomain.STATEMENT_TESTING: [
        r'\bFSFL[A-Z]+\b',              # FSF fields
        r'\bSTMTTST\b', 
        r'\b@[123][A-Z]+\b',            # JCL procedures
        r'\bGCHK\d*\b',                 # Checkpoint files
    ],
    MechDomain.XXA_TRANSACTIONS: [
        r'\bXXA\b',                      # XXA transactions
        r'\bSTDPC\s+TYPE=',             # STDPC macro calls
        r'\bprogram\s+call\s+chain\b',
    ],
}


class ExtractionMethod(Enum):
    """Extraction methods used."""
    PYMUPDF4LLM = "pymupdf4llm"
    PYMUPDF = "pymupdf"
    PYPDF = "pypdf"
    PYTHON_DOCX = "python-docx"
    OPENPYXL = "openpyxl"
    XLRD = "xlrd"
    PYTHON_PPTX = "python-pptx"
    NATIVE = "native"
    BEAUTIFULSOUP = "beautifulsoup"
    XML_PARSER = "xml_parser"
    IMAGE_METADATA = "image_metadata"
    FALLBACK = "fallback"


@dataclass
class ExtractionResult:
    """Structured extraction result."""
    success: bool
    text: str
    method: str
    pages: int = 0
    tables: int = 0
    images: int = 0
    sheets: int = 0
    slides: int = 0
    word_count: int = 0
    char_count: int = 0
    error: Optional[str] = None
    metadata: Optional[Dict[str, Any]] = None


@dataclass
class Chunk:
    """Document chunk with enhanced metadata for RAG."""
    document_id: str
    chunk_id: str
    chunk_number: int
    text: str
    char_count: int
    word_count: int
    page_number: Optional[int] = None
    
    # Enhanced context for RAG
    section_title: Optional[str] = None
    heading_hierarchy: Optional[List[str]] = None  # ["Chapter 1", "Section 1.2"]
    
    # Content type indicators (for filtering)
    has_table: bool = False
    has_image: bool = False
    has_list: bool = False
    has_code: bool = False
    
    # Adjacent context
    prev_chunk_preview: Optional[str] = None
    next_chunk_preview: Optional[str] = None
    
    # Document metadata
    document_title: Optional[str] = None
    document_author: Optional[str] = None
    document_type: Optional[str] = None
    
    # Additional metadata
    metadata: Optional[Dict[str, Any]] = None


# ============================================================================
# EMBEDDED EXTRACTORS (Standalone Lambda)
# ============================================================================

class BaseExtractor(ABC):
    """Abstract base class for document extractors."""

    def __init__(self, config):
        self.config = config

    @abstractmethod
    def extract(self, file_path: str, document_id: str) -> Dict[str, Any]:
        """Extract structured elements from a file."""
        raise NotImplementedError


# ------------------------------ PDF Extractor ------------------------------
try:
    import pymupdf4llm
    import pymupdf as fitz
except Exception:
    pymupdf4llm = None
    fitz = None


class PDFExtractor(BaseExtractor):
    """Extract content from PDF files using PyMuPDF4LLM."""

    def extract(self, file_path: str, document_id: str) -> Dict[str, Any]:
        if pymupdf4llm is None or fitz is None:
            raise ImportError("PyMuPDF4LLM not installed. Install with: pip install pymupdf4llm")

        logger.info(f"Extracting PDF: {file_path}")

        md_pages = pymupdf4llm.to_markdown(
            doc=file_path,
            page_chunks=True,
            write_images=True,
            image_path=os.path.join(self.config.TEMP_DIR, document_id),
            image_format="png",
            dpi=150
        )

        logger.info(f"Extracted {len(md_pages)} pages")

        doc = fitz.open(file_path)

        elements = []
        images_list = []
        tables_list = []

        for page_num, page_data in enumerate(md_pages):
            page_no = page_num + 1
            page_obj = doc[page_num]

            text_content = page_data.get('text', '')
            if text_content.strip():
                elements.append({
                    'type': 'text',
                    'content': text_content,
                    'metadata': {
                        'page_no': page_no,
                        'total_pages': len(md_pages),
                        'source': 'pymupdf4llm'
                    }
                })

            tables = page_obj.find_tables()
            for table_idx, table in enumerate(tables.tables):
                try:
                    table_df = table.to_pandas()
                    table_markdown = table_df.to_markdown(index=False)
                    table_html = table_df.to_html(index=False, classes=['table', 'table-striped'])
                    bbox = table.bbox

                    table_metadata = {
                        'page_no': page_no,
                        'table_index': table_idx,
                        'bbox': bbox,
                        'num_rows': len(table_df),
                        'num_cols': len(table_df.columns),
                        'columns': list(table_df.columns),
                        'html': table_html
                    }

                    elements.append({
                        'type': 'table',
                        'content': table_markdown,
                        'metadata': table_metadata
                    })

                    tables_list.append({
                        'html': table_html,
                        'markdown': table_markdown,
                        'metadata': table_metadata
                    })
                except Exception as e:
                    logger.warning(f"Error extracting table on page {page_no}: {str(e)}")
                    continue

            image_list = page_obj.get_images()
            for img_idx, img in enumerate(image_list):
                try:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]

                    image_filename = f"page_{page_no}_img_{img_idx}.{image_ext}"
                    image_path = os.path.join(self.config.TEMP_DIR, document_id, image_filename)
                    os.makedirs(os.path.dirname(image_path), exist_ok=True)
                    with open(image_path, 'wb') as img_file:
                        img_file.write(image_bytes)

                    img_rects = page_obj.get_image_rects(xref)
                    bbox = img_rects[0] if img_rects else None

                    # ⚡ CLAUDE VISION OCR - Extract text from embedded PDF image
                    ocr_text = extract_text_from_image_bytes(
                        image_bytes, 
                        image_ext,
                        source_info=f"PDF page {page_no} image {img_idx}"
                    )

                    image_metadata = {
                        'page_no': page_no,
                        'image_index': img_idx,
                        'bbox': bbox,
                        'ext': image_ext,
                        'size': len(image_bytes),
                        'local_path': image_path,
                        'has_ocr_text': bool(ocr_text),
                        'ocr_text_length': len(ocr_text) if ocr_text else 0
                    }

                    # Add OCR text as separate text element if found
                    if ocr_text:
                        elements.append({
                            'type': 'text',
                            'content': f"[Image OCR - Page {page_no}]\n{ocr_text}",
                            'metadata': {
                                'page_no': page_no,
                                'source': 'claude_vision_ocr',
                                'image_index': img_idx,
                                'ocr_source': 'embedded_pdf_image'
                            }
                        })

                    elements.append({
                        'type': 'image',
                        'content': ocr_text if ocr_text else '',
                        'metadata': image_metadata
                    })

                    images_list.append({
                        'local_path': image_path,
                        'ext': image_ext,
                        'ocr_text': ocr_text,
                        'metadata': image_metadata
                    })
                except Exception as e:
                    logger.warning(f"Error extracting image on page {page_no}: {str(e)}")
                    continue

        doc.close()

        return {
            'elements': elements,
            'images': images_list,
            'tables': tables_list,
            'metadata': {
                'total_pages': len(md_pages),
                'total_elements': len(elements),
                'total_images': len(images_list),
                'total_tables': len(tables_list)
            }
        }


# ------------------------------ DOCX Extractors ----------------------------
try:
    from unstructured.partition.docx import partition_docx
    from unstructured.partition.pptx import partition_pptx
except Exception:
    partition_docx = None
    partition_pptx = None

try:
    import docx2txt
    DOCX2TXT_AVAILABLE = True
except Exception:
    DOCX2TXT_AVAILABLE = False

try:
    import textract
    TEXTRACT_AVAILABLE = True
except Exception:
    TEXTRACT_AVAILABLE = False


def _safe_metadata_dict(element) -> dict:
    """Safely extract metadata dictionary from Unstructured element."""
    try:
        if element.metadata is None:
            return {}
        if hasattr(element.metadata, 'to_dict'):
            return element.metadata.to_dict()
        if isinstance(element.metadata, dict):
            return element.metadata
    except Exception as e:
        logger.warning(f"Could not extract metadata: {e}")
    return {}


class DOCXExtractor(BaseExtractor):
    """Extract content from DOCX and legacy DOC files with multiple fallbacks."""

    def extract(self, file_path: str, document_id: str) -> Dict[str, Any]:
        if partition_docx is None:
            raise ImportError("Unstructured library not installed.")

        logger.info(f"Extracting Word document: {file_path}")
        file_ext = os.path.splitext(file_path.lower())[1]

        if file_ext == '.doc':
            logger.info("Detected legacy .doc format - using fallback extraction methods")
            return self._extract_legacy_doc(file_path, document_id)

        elements = partition_docx(
            filename=file_path,
            include_page_breaks=True,
            infer_table_structure=True
        )

        structured_elements = []
        tables_list = []
        images_list = []

        for idx, element in enumerate(elements):
            element_type = element.category

            if element_type in ['Title', 'NarrativeText', 'Text']:
                structured_elements.append({
                    'type': 'text',
                    'content': str(element),
                    'metadata': {
                        'element_type': element_type,
                        'element_index': idx,
                        **_safe_metadata_dict(element)
                    }
                })

            elif element_type == 'Table':
                table_text = str(element)
                try:
                    table_html = element.metadata.text_as_html if hasattr(element.metadata, 'text_as_html') else None
                    table_metadata = {
                        'element_type': 'table',
                        'element_index': idx,
                        **_safe_metadata_dict(element)
                    }
                    structured_elements.append({
                        'type': 'table',
                        'content': table_text,
                        'metadata': table_metadata
                    })
                    tables_list.append({
                        'html': table_html or f"<pre>{table_text}</pre>",
                        'markdown': table_text,
                        'metadata': table_metadata
                    })
                except Exception as e:
                    logger.warning(f"Error parsing table: {str(e)}")
                    structured_elements.append({
                        'type': 'text',
                        'content': table_text,
                        'metadata': {
                            'element_type': 'table_as_text',
                            'element_index': idx
                        }
                    })

            elif element_type == 'Image':
                image_metadata = {
                    'element_type': 'image',
                    'element_index': idx,
                    **_safe_metadata_dict(element)
                }
                structured_elements.append({
                    'type': 'image',
                    'content': '[Image - see images list for binary data]',
                    'metadata': image_metadata
                })

            elif element_type == 'ListItem':
                structured_elements.append({
                    'type': 'text',
                    'content': str(element),
                    'metadata': {
                        'element_type': 'list_item',
                        'element_index': idx,
                        **_safe_metadata_dict(element)
                    }
                })
            else:
                structured_elements.append({
                    'type': 'text',
                    'content': str(element),
                    'metadata': {
                        'element_type': element_type.lower(),
                        'element_index': idx,
                        **_safe_metadata_dict(element)
                    }
                })

        images_list = self._extract_images_from_package(file_path, document_id)

        # Add OCR text from images as searchable text elements
        for img_idx, img in enumerate(images_list):
            ocr_text = img.get('ocr_text', '')
            if ocr_text:
                structured_elements.append({
                    'type': 'text',
                    'content': f"[Image OCR - {img.get('filename', 'unknown')}]\n{ocr_text}",
                    'metadata': {
                        'element_type': 'image_ocr_text',
                        'element_index': len(structured_elements),
                        'image_filename': img.get('filename'),
                        'image_index': img_idx,
                        'ocr_source': 'claude_vision',
                        'source': 'embedded_docx_image'
                    }
                })

        return {
            'elements': structured_elements,
            'images': images_list,
            'tables': tables_list,
            'metadata': {
                'total_elements': len(structured_elements),
                'total_images': len(images_list),
                'total_tables': len(tables_list),
                'images_with_ocr': sum(1 for img in images_list if img.get('has_ocr_text')),
                'extraction_method': 'unstructured'
            }
        }

    def _extract_images_from_package(self, file_path: str, document_id: str) -> List[Dict[str, Any]]:
        images = []
        try:
            with zipfile.ZipFile(file_path, 'r') as docx_zip:
                media_files = [f for f in docx_zip.namelist() if f.startswith('word/media/')]
                for idx, media_path in enumerate(media_files):
                    try:
                        image_bytes = docx_zip.read(media_path)
                        file_name = os.path.basename(media_path)
                        ext = os.path.splitext(file_name)[1].lower()
                        mime_map = {
                            '.png': 'image/png',
                            '.jpg': 'image/jpeg',
                            '.jpeg': 'image/jpeg',
                            '.gif': 'image/gif',
                            '.bmp': 'image/bmp',
                            '.tiff': 'image/tiff',
                            '.tif': 'image/tiff',
                            '.emf': 'image/x-emf',
                            '.wmf': 'image/x-wmf'
                        }
                        mime_type = mime_map.get(ext, 'application/octet-stream')
                        image_hash = hashlib.md5(image_bytes).hexdigest()
                        
                        # ⚡ CLAUDE VISION OCR - Extract text from embedded DOCX image
                        ocr_text = extract_text_from_image_bytes(
                            image_bytes,
                            ext.replace('.', ''),
                            source_info=f"DOCX image {file_name}"
                        )
                        
                        images.append({
                            'image_id': f"{document_id}_img_{idx}",
                            'filename': file_name,
                            'format': ext.replace('.', ''),
                            'mime_type': mime_type,
                            'size_bytes': len(image_bytes),
                            'hash': image_hash,
                            'base64_data': base64.b64encode(image_bytes).decode('utf-8'),
                            'ocr_text': ocr_text,  # Claude Vision OCR text
                            'has_ocr_text': bool(ocr_text),
                            'metadata': {
                                'source_path': media_path,
                                'document_id': document_id,
                                'image_index': idx,
                                'ocr_text_length': len(ocr_text) if ocr_text else 0
                            }
                        })
                    except Exception as e:
                        logger.warning(f"Failed to extract image {media_path}: {e}")
                        continue
        except Exception as e:
            logger.warning(f"Error extracting images from package: {e}")
        return images

    def _extract_legacy_doc(self, file_path: str, document_id: str) -> Dict[str, Any]:
        logger.info(f"Processing legacy .doc file: {os.path.basename(file_path)}")

        try:
            converted_path = self._convert_doc_to_docx_libreoffice(file_path)
            if converted_path and os.path.exists(converted_path):
                logger.info("Converted .doc to .docx using LibreOffice")
                result = self.extract(converted_path, document_id)
                result['metadata']['converted_from_legacy_doc'] = True
                result['metadata']['conversion_method'] = 'libreoffice'
                try:
                    os.remove(converted_path)
                except Exception:
                    pass
                return result
        except Exception as e:
            logger.warning(f"LibreOffice conversion failed: {str(e)}")

        if TEXTRACT_AVAILABLE:
            try:
                text_content = textract.process(file_path).decode('utf-8', errors='ignore')
                logger.info(f"Extracted {len(text_content)} chars using textract")
                return self._create_text_only_result(text_content, 'textract')
            except Exception as e:
                logger.warning(f"Textract extraction failed: {str(e)}")

        if DOCX2TXT_AVAILABLE:
            try:
                text_content = docx2txt.process(file_path)
                logger.info(f"Extracted {len(text_content)} chars using docx2txt")
                return self._create_text_only_result(text_content, 'docx2txt')
            except Exception as e:
                logger.warning(f"docx2txt extraction failed: {str(e)}")

        try:
            text_content = self._extract_with_antiword(file_path)
            if text_content:
                logger.info(f"Extracted {len(text_content)} chars using antiword")
                return self._create_text_only_result(text_content, 'antiword')
        except Exception as e:
            logger.warning(f"Antiword extraction failed: {str(e)}")

        try:
            with open(file_path, 'rb') as f:
                raw_bytes = f.read()
                text_content = raw_bytes.decode('latin-1', errors='ignore')
                text_content = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text_content)
                text_content = ' '.join(text_content.split())
                if len(text_content) > 100:
                    logger.warning("Using raw text extraction (low quality)")
                    return self._create_text_only_result(text_content, 'raw_binary')
        except Exception as e:
            logger.error(f"Raw text extraction failed: {str(e)}")

        error_msg = (
            f"Failed to extract content from legacy .doc file: {os.path.basename(file_path)}. "
            f"All extraction methods failed. Install LibreOffice or textract for better support."
        )
        logger.error(error_msg)
        raise Exception(error_msg)

    def _convert_doc_to_docx_libreoffice(self, file_path: str) -> Optional[str]:
        output_dir = tempfile.gettempdir()
        libreoffice_commands = [
            'libreoffice',
            'soffice',
            '/usr/bin/libreoffice',
            '/opt/libreoffice/program/soffice',
            '/Applications/LibreOffice.app/Contents/MacOS/soffice'
        ]
        for cmd in libreoffice_commands:
            try:
                subprocess.run(
                    [cmd, '--headless', '--convert-to', 'docx', '--outdir', output_dir, file_path],
                    capture_output=True,
                    timeout=30,
                    check=True
                )
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                converted_path = os.path.join(output_dir, f"{base_name}.docx")
                if os.path.exists(converted_path):
                    logger.info(f"Converted using: {cmd}")
                    return converted_path
            except Exception:
                continue
        return None

    def _extract_with_antiword(self, file_path: str) -> Optional[str]:
        try:
            result = subprocess.run(
                ['antiword', file_path],
                capture_output=True,
                timeout=30,
                check=True
            )
            return result.stdout.decode('utf-8', errors='ignore')
        except Exception:
            return None

    def _create_text_only_result(self, text_content: str, method: str) -> Dict[str, Any]:
        paragraphs = [p.strip() for p in text_content.split('\n\n') if p.strip()]
        if len(paragraphs) <= 1:
            paragraphs = [p.strip() for p in text_content.split('\n') if p.strip() and len(p.strip()) > 20]

        structured_elements = []
        for idx, para in enumerate(paragraphs):
            if para:
                structured_elements.append({
                    'type': 'text',
                    'content': para,
                    'metadata': {
                        'element_type': 'paragraph',
                        'element_index': idx,
                        'extraction_method': method,
                        'is_legacy_doc': True
                    }
                })

        return {
            'elements': structured_elements,
            'images': [],
            'tables': [],
            'metadata': {
                'total_elements': len(structured_elements),
                'total_images': 0,
                'total_tables': 0,
                'extraction_method': method,
                'is_legacy_doc': True,
                'format': 'doc',
                'warning': 'Legacy .doc format: tables, images, and formatting may be lost'
            }
        }


# Enhanced DOCX extractor (python-docx)
try:
    from docx import Document
    from docx.oxml.ns import qn
    from docx.table import Table as DocxTable
    from docx.text.paragraph import Paragraph
    PYTHON_DOCX_AVAILABLE = True
except Exception:
    PYTHON_DOCX_AVAILABLE = False

try:
    from PIL import Image
    PIL_AVAILABLE = True
except Exception:
    PIL_AVAILABLE = False


class EnhancedDOCXExtractor(BaseExtractor):
    """Enhanced DOCX extractor using python-docx."""

    def __init__(self, config):
        super().__init__(config)
        if not PYTHON_DOCX_AVAILABLE:
            raise ImportError("python-docx not installed. Install with: pip install python-docx")

    def extract(self, file_path: str, document_id: str) -> Dict[str, Any]:
        logger.info(f"Enhanced DOCX extraction: {file_path}")

        try:
            doc = Document(file_path)
            elements = []
            images_list = []
            tables_list = []
            element_counter = 0
            paragraph_counter = 0
            total_paragraphs = self._count_total_paragraphs(doc)
            doc_metadata = self._extract_doc_metadata(doc)
            heading_stack = []

            for block in self._iter_block_items(doc):
                relative_pos = paragraph_counter / max(total_paragraphs, 1)

                if isinstance(block, Paragraph):
                    para_data = self._process_paragraph(
                        block, element_counter, paragraph_counter,
                        heading_stack, relative_pos, document_id
                    )
                    if para_data:
                        elements.append(para_data['element'])
                        para_images = self._extract_paragraph_images(
                            block, document_id, element_counter,
                            paragraph_counter, elements
                        )
                        for img in para_images:
                            elements.append(img['element'])
                            images_list.append(img['image_data'])
                            element_counter += 1
                        element_counter += 1
                    paragraph_counter += 1

                elif isinstance(block, DocxTable):
                    table_data = self._process_table(
                        block, element_counter, paragraph_counter,
                        heading_stack, relative_pos
                    )
                    if table_data:
                        elements.append(table_data['element'])
                        tables_list.append(table_data['table_data'])
                        element_counter += 1
                    paragraph_counter += 1

            media_images = self._extract_media_images(doc, document_id, element_counter)
            for img in media_images:
                if not self._image_already_extracted(img, images_list):
                    images_list.append(img)
                    
                    # Add OCR text from image as searchable text element
                    ocr_text = img.get('ocr_text', '')
                    if ocr_text:
                        elements.append({
                            'type': 'text',
                            'content': f"[Image OCR - Media]\n{ocr_text}",
                            'metadata': {
                                'element_type': 'image_ocr_text',
                                'element_index': element_counter,
                                'image_id': img.get('image_id'),
                                'ocr_source': 'claude_vision',
                                'source': 'embedded_docx_media'
                            }
                        })
                        element_counter += 1

            logger.info(f"Extracted: {len(elements)} elements, {len(images_list)} images, {len(tables_list)} tables")

            return {
                'elements': elements,
                'images': images_list,
                'tables': tables_list,
                'metadata': {
                    **doc_metadata,
                    'total_elements': len(elements),
                    'total_images': len(images_list),
                    'total_tables': len(tables_list),
                    'images_with_ocr': sum(1 for img in images_list if img.get('has_ocr_text')),
                    'extraction_method': 'enhanced_python_docx'
                }
            }
        except Exception as e:
            logger.error(f"Enhanced DOCX extraction failed: {e}", exc_info=True)
            logger.info("Falling back to basic Unstructured extraction")
            basic_extractor = DOCXExtractor(self.config)
            return basic_extractor.extract(file_path, document_id)

    def _iter_block_items(self, doc):
        parent = doc.element.body
        for child in parent.iterchildren():
            if child.tag == qn('w:p'):
                for para in doc.paragraphs:
                    if para._element is child:
                        yield para
                        break
            elif child.tag == qn('w:tbl'):
                for table in doc.tables:
                    if table._element is child:
                        yield table
                        break

    def _count_total_paragraphs(self, doc) -> int:
        count = 0
        for item in doc.element.body:
            if item.tag in [qn('w:p'), qn('w:tbl')]:
                count += 1
        return count

    def _extract_doc_metadata(self, doc) -> Dict[str, Any]:
        metadata = {}
        try:
            core_props = doc.core_properties
            metadata['title'] = core_props.title or ''
            metadata['author'] = core_props.author or ''
            metadata['subject'] = core_props.subject or ''
            metadata['keywords'] = core_props.keywords or ''
            metadata['created'] = str(core_props.created) if core_props.created else ''
            metadata['modified'] = str(core_props.modified) if core_props.modified else ''
        except Exception:
            pass
        return metadata

    def _process_paragraph(self, para: Paragraph, elem_idx: int, para_idx: int,
                           heading_stack: List[str], rel_pos: float,
                           document_id: str) -> Optional[Dict]:
        text = para.text.strip()
        if not text:
            return None

        style_name = para.style.name if para.style else 'Normal'
        element_type = 'text'
        heading_level = 0

        if style_name.startswith('Heading'):
            try:
                heading_level = int(style_name.replace('Heading ', '').strip())
                element_type = 'heading'
                heading_stack[:] = heading_stack[:heading_level - 1] + [text[:100]]
            except ValueError:
                pass
        elif style_name == 'Title':
            element_type = 'heading'
            heading_level = 1
            heading_stack[:] = [text[:100]]

        if para._element.xpath('.//w:numPr'):
            element_type = 'list_item'

        formatting = self._extract_paragraph_formatting(para)

        return {
            'element': {
                'type': element_type,
                'content': text,
                'metadata': {
                    'element_type': style_name,
                    'element_index': elem_idx,
                    'paragraph_index': para_idx,
                    'heading_level': heading_level,
                    'headings': heading_stack.copy(),
                    'relative_position': rel_pos,
                    'formatting': formatting,
                    'has_images': len(para._element.xpath(
                        './/a:blip', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                    )) > 0
                }
            }
        }

    def _extract_paragraph_formatting(self, para: Paragraph) -> Dict[str, Any]:
        formatting = {
            'style': para.style.name if para.style else 'Normal',
            'alignment': str(para.alignment) if para.alignment else 'LEFT',
        }
        if para.runs:
            first_run = para.runs[0]
            formatting['bold'] = first_run.bold or False
            formatting['italic'] = first_run.italic or False
            formatting['underline'] = first_run.underline or False
            if first_run.font.size:
                formatting['font_size_pt'] = first_run.font.size.pt
        return formatting

    def _extract_paragraph_images(self, para: Paragraph, document_id: str,
                                   elem_idx: int, para_idx: int,
                                   elements: List[Dict]) -> List[Dict]:
        images = []
        drawing_elements = para._element.xpath('.//w:drawing')
        for draw_idx, drawing in enumerate(drawing_elements):
            try:
                extent = drawing.xpath(
                    './/wp:extent', namespaces={'wp': 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing'}
                )
                width_emu = 0
                height_emu = 0
                if extent:
                    width_emu = int(extent[0].get('cx', 0))
                    height_emu = int(extent[0].get('cy', 0))

                width_inches = width_emu / 914400.0
                height_inches = height_emu / 914400.0

                blip = drawing.xpath(
                    './/a:blip', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                )
                if blip:
                    embed_attr = blip[0].get(qn('r:embed'))
                    if embed_attr:
                        alt_text = ''
                        desc_attr = drawing.xpath(
                            './/wp:docPr', namespaces={'wp': 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing'}
                        )
                        if desc_attr:
                            alt_text = desc_attr[0].get('descr', '') or desc_attr[0].get('title', '')

                        context_before = self._get_context_before(elements, 500)

                        image_element = {
                            'element': {
                                'type': 'image',
                                'content': f"[Image: {alt_text or 'Inline image'}]",
                                'metadata': {
                                    'element_type': 'image',
                                    'element_index': elem_idx + draw_idx + 1,
                                    'paragraph_index': para_idx,
                                    'image_id': embed_attr,
                                    'width_inches': width_inches,
                                    'height_inches': height_inches,
                                    'width_emu': width_emu,
                                    'height_emu': height_emu,
                                    'alt_text': alt_text,
                                    'context_before': context_before[:500],
                                    'is_inline': True,
                                    'position': {
                                        'type': 'inline',
                                        'paragraph_index': para_idx,
                                        'relative_position': para_idx
                                    }
                                }
                            },
                            'image_data': {
                                'image_id': embed_attr,
                                'width_inches': width_inches,
                                'height_inches': height_inches,
                                'alt_text': alt_text,
                                'format': 'embedded',
                                'paragraph_index': para_idx,
                                'metadata': {
                                    'width_emu': width_emu,
                                    'height_emu': height_emu
                                }
                            }
                        }
                        images.append(image_element)
            except Exception as e:
                logger.warning(f"Error extracting image from paragraph: {e}")
                continue
        return images

    def _get_context_before(self, elements: List[Dict], max_chars: int) -> str:
        context_parts = []
        chars_collected = 0
        for elem in reversed(elements[-10:]):
            if elem.get('type') == 'text':
                content = elem.get('content', '')
                context_parts.insert(0, content)
                chars_collected += len(content)
                if chars_collected >= max_chars:
                    break
        return ' '.join(context_parts)[-max_chars:]

    def _process_table(self, table: DocxTable, elem_idx: int, para_idx: int,
                       heading_stack: List[str], rel_pos: float) -> Optional[Dict]:
        try:
            rows = []
            headers = []
            for row_idx, row in enumerate(table.rows):
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    row_data.append(cell_text)
                if row_idx == 0:
                    headers = row_data
                rows.append(row_data)

            markdown_lines = []
            if rows:
                markdown_lines.append('| ' + ' | '.join(rows[0]) + ' |')
                markdown_lines.append('|' + '---|' * len(rows[0]))
                for row in rows[1:]:
                    markdown_lines.append('| ' + ' | '.join(row) + ' |')

            markdown_content = '\n'.join(markdown_lines)

            return {
                'element': {
                    'type': 'table',
                    'content': markdown_content,
                    'metadata': {
                        'element_type': 'table',
                        'element_index': elem_idx,
                        'paragraph_index': para_idx,
                        'num_rows': len(rows),
                        'num_cols': len(rows[0]) if rows else 0,
                        'headers': headers,
                        'headings': heading_stack.copy(),
                        'relative_position': rel_pos,
                        'has_header_row': True,
                    }
                },
                'table_data': {
                    'markdown': markdown_content,
                    'rows': rows,
                    'headers': headers,
                    'num_rows': len(rows),
                    'num_cols': len(rows[0]) if rows else 0,
                    'element_index': elem_idx,
                }
            }
        except Exception as e:
            logger.warning(f"Error processing table: {e}")
            return None

    def _extract_media_images(self, doc, document_id: str, start_idx: int) -> List[Dict]:
        images = []
        try:
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        image_part = rel.target_part
                        image_bytes = image_part.blob
                        content_type = image_part.content_type
                        if 'png' in content_type:
                            img_format = 'png'
                        elif 'jpeg' in content_type or 'jpg' in content_type:
                            img_format = 'jpeg'
                        elif 'gif' in content_type:
                            img_format = 'gif'
                        else:
                            img_format = content_type.split('/')[-1]

                        width = 0
                        height = 0
                        if PIL_AVAILABLE:
                            try:
                                img = Image.open(io.BytesIO(image_bytes))
                                width, height = img.size
                            except Exception:
                                pass

                        img_id = hashlib.md5(image_bytes[:1000]).hexdigest()[:12]
                        img_filename = f"{document_id}_img_{img_id}.{img_format}"
                        img_path = os.path.join(self.config.TEMP_DIR, document_id, img_filename)
                        os.makedirs(os.path.dirname(img_path), exist_ok=True)
                        with open(img_path, 'wb') as f:
                            f.write(image_bytes)

                        # ⚡ CLAUDE VISION OCR - Extract text from embedded DOCX image
                        ocr_text = extract_text_from_image_bytes(
                            image_bytes,
                            img_format,
                            source_info=f"DOCX media image {img_filename}"
                        )

                        images.append({
                            'image_id': img_id,
                            'local_path': img_path,
                            'format': img_format,
                            'width_px': width,
                            'height_px': height,
                            'size_bytes': len(image_bytes),
                            'ocr_text': ocr_text,  # Claude Vision OCR text
                            'has_ocr_text': bool(ocr_text),
                            'metadata': {
                                'source': 'document_media',
                                'rel_id': rel.rId,
                                'ocr_text_length': len(ocr_text) if ocr_text else 0
                            }
                        })
                    except Exception as e:
                        logger.warning(f"Error extracting media image: {e}")
                        continue
        except Exception as e:
            logger.warning(f"Error accessing document media: {e}")
        return images

    def _image_already_extracted(self, image: Dict, existing: List[Dict]) -> bool:
        img_id = image.get('image_id', '')
        for existing_img in existing:
            if existing_img.get('image_id', '') == img_id:
                return True
        return False


# ------------------------------ XLSX Extractor -----------------------------
try:
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter
    OPENPYXL_AVAILABLE = True
except Exception:
    OPENPYXL_AVAILABLE = False


class XLSXExtractor(BaseExtractor):
    """Extract content from XLSX files using openpyxl (optimized - no pandas dependency)."""

    def extract(self, file_path: str, document_id: str) -> Dict[str, Any]:
        if not OPENPYXL_AVAILABLE:
            raise ImportError("openpyxl not installed. Install with: pip install openpyxl")

        logger.info(f"Extracting XLSX: {file_path}")
        
        try:
            # Load workbook in read-only mode for better performance
            wb = load_workbook(file_path, data_only=True, read_only=True)
            sheet_names = wb.sheetnames
        except Exception as e:
            logger.error(f"Failed to load workbook {file_path}: {e}")
            raise

        structured_elements = []
        tables_list = []

        for sheet_idx, sheet_name in enumerate(sheet_names):
            try:
                sheet = wb[sheet_name]
                
                # Extract all rows
                rows = list(sheet.iter_rows(values_only=True))
                if not rows:
                    logger.info(f"Skipping empty sheet: {sheet_name}")
                    continue
                
                # Get header row (first row)
                headers = [str(cell) if cell is not None else f"Column{i+1}" 
                          for i, cell in enumerate(rows[0])]
                
                # Build markdown table
                markdown_lines = []
                markdown_lines.append("| " + " | ".join(headers) + " |")
                markdown_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                
                # Add data rows
                data_rows = rows[1:]
                for row in data_rows:
                    # Skip completely empty rows
                    if all(cell is None or str(cell).strip() == '' for cell in row):
                        continue
                    row_values = [str(cell) if cell is not None else '' for cell in row]
                    # Pad row to match header length
                    while len(row_values) < len(headers):
                        row_values.append('')
                    markdown_lines.append("| " + " | ".join(row_values[:len(headers)]) + " |")
                
                table_markdown = "\n".join(markdown_lines)
                
                # Build HTML table
                html_lines = ['<table class="table table-striped">']
                html_lines.append('  <thead>')
                html_lines.append('    <tr>')
                for header in headers:
                    html_lines.append(f'      <th>{header}</th>')
                html_lines.append('    </tr>')
                html_lines.append('  </thead>')
                html_lines.append('  <tbody>')
                
                for row in data_rows:
                    if all(cell is None or str(cell).strip() == '' for cell in row):
                        continue
                    html_lines.append('    <tr>')
                    for i, cell in enumerate(row):
                        if i >= len(headers):
                            break
                        cell_value = str(cell) if cell is not None else ''
                        html_lines.append(f'      <td>{cell_value}</td>')
                    # Pad with empty cells if needed
                    for i in range(len(row), len(headers)):
                        html_lines.append('      <td></td>')
                    html_lines.append('    </tr>')
                
                html_lines.append('  </tbody>')
                html_lines.append('</table>')
                table_html = '\n'.join(html_lines)

                metadata = {
                    'element_type': 'table',
                    'sheet_name': sheet_name,
                    'sheet_index': sheet_idx,
                    'num_rows': len(data_rows),
                    'num_cols': len(headers),
                    'columns': headers
                }

                structured_elements.append({
                    'type': 'table',
                    'content': table_markdown,
                    'metadata': metadata
                })

                tables_list.append({
                    'html': table_html,
                    'markdown': table_markdown,
                    'metadata': metadata
                })
            except Exception as e:
                logger.warning(f"Error reading sheet '{sheet_name}': {str(e)}")
                continue

        wb.close()

        return {
            'elements': structured_elements,
            'images': [],
            'tables': tables_list,
            'metadata': {
                'total_sheets': len(sheet_names),
                'total_tables': len(tables_list)
            }
        }


# ------------------------------ PPTX Extractor -----------------------------
class PPTXExtractor(BaseExtractor):
    """Extract content from PPTX files using Unstructured."""

    def extract(self, file_path: str, document_id: str) -> Dict[str, Any]:
        if partition_pptx is None:
            raise ImportError("Unstructured library not installed. Install with: pip install unstructured[pptx]")

        logger.info(f"Extracting PPTX: {file_path}")
        elements = partition_pptx(
            filename=file_path,
            include_page_breaks=True,
            infer_table_structure=True
        )

        structured_elements = []
        tables_list = []
        current_slide = 1

        for idx, element in enumerate(elements):
            element_type = element.category

            if element_type == 'PageBreak':
                current_slide += 1
                continue

            metadata = {
                'element_type': element_type,
                'element_index': idx,
                'slide_number': current_slide,
                **_safe_metadata_dict(element)
            }

            if element_type in ['Title', 'NarrativeText', 'Text']:
                structured_elements.append({
                    'type': 'text',
                    'content': str(element),
                    'metadata': metadata
                })
            elif element_type == 'Table':
                table_text = str(element)
                structured_elements.append({
                    'type': 'table',
                    'content': table_text,
                    'metadata': metadata
                })
                tables_list.append({
                    'html': f"<pre>{table_text}</pre>",
                    'markdown': table_text,
                    'metadata': metadata
                })
            elif element_type == 'ListItem':
                structured_elements.append({
                    'type': 'text',
                    'content': f"- {str(element)}",
                    'metadata': metadata
                })
            else:
                structured_elements.append({
                    'type': 'text',
                    'content': str(element),
                    'metadata': metadata
                })

        return {
            'elements': structured_elements,
            'images': [],
            'tables': tables_list,
            'metadata': {
                'total_slides': current_slide,
                'total_elements': len(structured_elements),
                'total_tables': len(tables_list)
            }
        }


# ------------------------------ Text Extractor -----------------------------
class TextExtractor(BaseExtractor):
    """Extract content from plain text files (.txt, .md, etc.)."""

    def extract(self, file_path: str, document_id: str) -> Dict[str, Any]:
        logger.info(f"Extracting text file: {file_path}")
        try:
            text_content = None
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        text_content = f.read()
                    logger.info(f"Read file with {encoding} encoding")
                    break
                except UnicodeDecodeError:
                    continue

            if text_content is None:
                with open(file_path, 'rb') as f:
                    raw_bytes = f.read()
                text_content = raw_bytes.decode('utf-8', errors='ignore')
                logger.warning("Used fallback binary read with error handling")

            structured_elements = self._parse_text_to_elements(text_content, file_path)
            logger.info(f"Extracted {len(structured_elements)} elements from text file")
            return {
                'elements': structured_elements,
                'images': [],
                'tables': [],
                'metadata': {
                    'total_elements': len(structured_elements),
                    'total_images': 0,
                    'total_tables': 0,
                    'file_size': os.path.getsize(file_path),
                    'extraction_method': 'text_file',
                    'format': os.path.splitext(file_path)[1].lower()
                }
            }
        except Exception as e:
            logger.error(f"Error extracting text file: {str(e)}", exc_info=True)
            raise

    def _parse_text_to_elements(self, content: str, file_path: str) -> List[Dict]:
        elements = []
        content = content.replace('\r\n', '\n').replace('\r', '\n')

        is_markdown = (
            file_path.lower().endswith('.md') or
            re.search(r'^#+\s+\S', content, re.MULTILINE) is not None or
            re.search(r'^```', content, re.MULTILINE) is not None
        )

        if is_markdown:
            return self._parse_markdown(content)
        return self._parse_plain_text(content)

    def _parse_markdown(self, content: str) -> List[Dict]:
        elements = []
        lines = content.split('\n')
        current_text = []
        in_code_block = False
        code_lang = ''
        element_index = 0
        i = 0
        while i < len(lines):
            line = lines[i]
            if line.startswith('```'):
                if in_code_block:
                    code_content = '\n'.join(current_text)
                    if code_content.strip():
                        elements.append({
                            'type': 'code',
                            'content': code_content,
                            'metadata': {
                                'element_type': 'code_block',
                                'element_index': element_index,
                                'language': code_lang
                            }
                        })
                        element_index += 1
                    current_text = []
                    in_code_block = False
                else:
                    if current_text:
                        text_content = '\n'.join(current_text).strip()
                        if text_content:
                            elements.append({
                                'type': 'text',
                                'content': text_content,
                                'metadata': {
                                    'element_type': 'paragraph',
                                    'element_index': element_index
                                }
                            })
                            element_index += 1
                        current_text = []
                    in_code_block = True
                    code_lang = line[3:].strip()
                i += 1
                continue

            if in_code_block:
                current_text.append(line)
                i += 1
                continue

            heading_match = re.match(r'^(#{1,6})\s+(.+)$', line)
            if heading_match:
                if current_text:
                    text_content = '\n'.join(current_text).strip()
                    if text_content:
                        elements.append({
                            'type': 'text',
                            'content': text_content,
                            'metadata': {
                                'element_type': 'paragraph',
                                'element_index': element_index
                            }
                        })
                        element_index += 1
                    current_text = []
                level = len(heading_match.group(1))
                heading_text = heading_match.group(2).strip()
                elements.append({
                    'type': 'text',
                    'content': heading_text,
                    'metadata': {
                        'element_type': 'heading',
                        'element_index': element_index,
                        'heading_level': level
                    }
                })
                element_index += 1
                i += 1
                continue

            if i + 1 < len(lines):
                next_line = lines[i + 1]
                if re.match(r'^={3,}$', next_line):
                    if current_text:
                        text_content = '\n'.join(current_text).strip()
                        if text_content:
                            elements.append({
                                'type': 'text',
                                'content': text_content,
                                'metadata': {
                                    'element_type': 'paragraph',
                                    'element_index': element_index
                                }
                            })
                            element_index += 1
                        current_text = []

                    elements.append({
                        'type': 'text',
                        'content': line.strip(),
                        'metadata': {
                            'element_type': 'heading',
                            'element_index': element_index,
                            'heading_level': 1
                        }
                    })
                    element_index += 1
                    i += 2
                    continue
                if re.match(r'^-{3,}$', next_line):
                    if current_text:
                        text_content = '\n'.join(current_text).strip()
                        if text_content:
                            elements.append({
                                'type': 'text',
                                'content': text_content,
                                'metadata': {
                                    'element_type': 'paragraph',
                                    'element_index': element_index
                                }
                            })
                            element_index += 1
                        current_text = []

                    elements.append({
                        'type': 'text',
                        'content': line.strip(),
                        'metadata': {
                            'element_type': 'heading',
                            'element_index': element_index,
                            'heading_level': 2
                        }
                    })
                    element_index += 1
                    i += 2
                    continue

            if not line.strip():
                if current_text:
                    text_content = '\n'.join(current_text).strip()
                    if text_content:
                        elements.append({
                            'type': 'text',
                            'content': text_content,
                            'metadata': {
                                'element_type': 'paragraph',
                                'element_index': element_index
                            }
                        })
                        element_index += 1
                    current_text = []
                i += 1
                continue

            current_text.append(line)
            i += 1

        if current_text:
            text_content = '\n'.join(current_text).strip()
            if text_content:
                if in_code_block:
                    elements.append({
                        'type': 'code',
                        'content': text_content,
                        'metadata': {
                            'element_type': 'code_block',
                            'element_index': element_index,
                            'language': code_lang
                        }
                    })
                else:
                    elements.append({
                        'type': 'text',
                        'content': text_content,
                        'metadata': {
                            'element_type': 'paragraph',
                            'element_index': element_index
                        }
                    })
        return elements

    def _parse_plain_text(self, content: str) -> List[Dict]:
        elements = []
        element_index = 0
        paragraphs = re.split(r'\n\s*\n', content)
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            lines = para.split('\n')
            if len(lines) == 1 and len(para) < 100:
                if para.isupper() or (para.endswith(':') and len(para) < 60):
                    elements.append({
                        'type': 'text',
                        'content': para,
                        'metadata': {
                            'element_type': 'heading',
                            'element_index': element_index,
                            'heading_level': 2
                        }
                    })
                    element_index += 1
                    continue
            elements.append({
                'type': 'text',
                'content': para,
                'metadata': {
                    'element_type': 'paragraph',
                    'element_index': element_index
                }
            })
            element_index += 1
        return elements


# ------------------------------ Image Extractor ----------------------------
# Uses Claude Vision (Bedrock) for OCR - zero additional package size
# Fallback: AWS Textract (paid ~$1.50/1000 pages)

import base64

CLAUDE_VISION_MODEL = os.environ.get(
    'CLAUDE_VISION_MODEL',
    os.environ.get('HAIKU_INFERENCE_PROFILE_ARN', 'arn:aws:bedrock:us-east-1:966245098051:application-inference-profile/wt15xad4wqki')
)
ENABLE_CLAUDE_VISION_OCR = os.environ.get('ENABLE_CLAUDE_VISION_OCR', 'true').lower() == 'true'


# ============================================================================
# GLOBAL IMAGE OCR HELPER - Used by ALL extractors (PDF, DOCX, XLSX, PPTX)
# ============================================================================

_bedrock_client_global = None
_textract_client_global = None

def get_bedrock_client():
    """Get or create global Bedrock client for Claude Vision OCR."""
    global _bedrock_client_global
    if _bedrock_client_global is None:
        try:
            _bedrock_client_global = boto3.client('bedrock-runtime', region_name=AWS_REGION)
            logger.info("Initialized Bedrock client for Claude Vision OCR")
        except Exception as e:
            logger.warning(f"Failed to initialize Bedrock client: {e}")
            _bedrock_client_global = None
    return _bedrock_client_global


def get_textract_client():
    """Get or create global Textract client for OCR fallback."""
    global _textract_client_global
    if _textract_client_global is None and TEXTRACT_ENABLED:
        try:
            _textract_client_global = boto3.client('textract', region_name=AWS_REGION)
            logger.info("Initialized Textract client for OCR fallback")
        except Exception as e:
            logger.warning(f"Failed to initialize Textract client: {e}")
            _textract_client_global = None
    return _textract_client_global


def extract_text_from_image_bytes(image_bytes: bytes, image_format: str = 'png', source_info: str = '') -> str:
    """
    Extract text from image bytes using Claude Vision OCR.
    
    This function is called by ALL document extractors (PDF, DOCX, XLSX, PPTX)
    to OCR embedded images.
    
    Args:
        image_bytes: Raw image data
        image_format: Image format (png, jpg, jpeg, etc.)
        source_info: Description of image source for logging
    
    Returns:
        Extracted text from image, or empty string if no text found
    """
    if not image_bytes:
        return ""
    
    ocr_text = ""
    ocr_source = "none"
    
    try:
        # Try Claude Vision first (if enabled)
        if ENABLE_CLAUDE_VISION_OCR and OCR_ENABLED:
            bedrock_client = get_bedrock_client()
            if bedrock_client:
                try:
                    # Map format to MIME type
                    mime_type_map = {
                        'png': 'image/png',
                        'jpg': 'image/jpeg',
                        'jpeg': 'image/jpeg',
                        'gif': 'image/gif',
                        'bmp': 'image/bmp',
                        'webp': 'image/webp',
                        'tiff': 'image/tiff',
                        'tif': 'image/tiff'
                    }
                    mime_type = mime_type_map.get(image_format.lower().replace('.', ''), 'image/png')
                    
                    response = bedrock_client.invoke_model(
                        modelId=CLAUDE_VISION_MODEL,
                        body=json.dumps({
                            'anthropic_version': 'bedrock-2023-05-31',
                            'max_tokens': 4096,
                            'messages': [{
                                'role': 'user',
                                'content': [
                                    {
                                        'type': 'image',
                                        'source': {
                                            'type': 'base64',
                                            'media_type': mime_type,
                                            'data': base64.standard_b64encode(image_bytes).decode('utf-8')
                                        }
                                    },
                                    {
                                        'type': 'text',
                                        'text': 'Extract ALL text from this image including labels, titles, tables, code, and diagrams. Format as clean readable text. If no text is visible, respond "NO_TEXT".'
                                    }
                                ]
                            }]
                        })
                    )
                    
                    result = json.loads(response['body'].read())
                    text = result.get('content', [{}])[0].get('text', '')
                    
                    if text and 'NO_TEXT' not in text.upper():
                        ocr_text = text.strip()
                        ocr_source = 'claude_vision'
                        logger.info(f"Claude Vision OCR: extracted {len(ocr_text)} chars from {source_info}")
                    
                except Exception as e:
                    logger.warning(f"Claude Vision OCR failed for {source_info}: {e}")
        
        # Fallback to Textract if Claude didn't work
        if not ocr_text and TEXTRACT_ENABLED and OCR_ENABLED:
            textract_client = get_textract_client()
            if textract_client:
                try:
                    response = textract_client.detect_document_text(Document={'Bytes': image_bytes})
                    lines = [block['Text'] for block in response.get('Blocks', []) if block['BlockType'] == 'LINE']
                    ocr_text = '\n'.join(lines)
                    ocr_source = 'textract'
                    logger.info(f"Textract OCR: extracted {len(ocr_text)} chars from {source_info}")
                except Exception as e:
                    logger.warning(f"Textract OCR failed for {source_info}: {e}")
    
    except Exception as e:
        logger.error(f"Image OCR failed for {source_info}: {e}")
    
    return ocr_text


class ImageExtractor(BaseExtractor):
    """Extract text from images using Claude Vision (Bedrock)."""

    def __init__(self, config):
        super().__init__(config)
        self._textract_client = None
        self._bedrock_client = None

    @property
    def bedrock_client(self):
        if self._bedrock_client is None:
            try:
                self._bedrock_client = boto3.client('bedrock-runtime', region_name=AWS_REGION)
            except Exception as e:
                logger.warning(f"Could not initialize Bedrock client: {e}")
                self._bedrock_client = False
        return self._bedrock_client if self._bedrock_client else None

    @property
    def textract_client(self):
        if self._textract_client is None and self.config.TEXTRACT_ENABLED:
            try:
                self._textract_client = boto3.client('textract')
            except Exception as e:
                logger.warning(f"Could not initialize Textract client: {e}")
                self._textract_client = False
        return self._textract_client if self._textract_client else None

    def extract(self, file_path: str, document_id: str) -> Dict[str, Any]:
        logger.info(f"Extracting image: {file_path}")
        try:
            with open(file_path, 'rb') as img_file:
                image_bytes = img_file.read()

            ocr_text = ""
            ocr_source = "none"
            
            if self.config.OCR_ENABLED:
                # Try Claude Vision first, then Textract as fallback
                if ENABLE_CLAUDE_VISION_OCR and self.bedrock_client:
                    ocr_text = self._extract_with_claude_vision(image_bytes, file_path)
                    if ocr_text:
                        ocr_source = 'claude_vision'
                
                if not ocr_text and self.textract_client:
                    ocr_text = self._extract_with_textract(image_bytes)
                    if ocr_text:
                        ocr_source = 'textract'
                
                if ocr_text:
                    logger.info(f"Extracted {len(ocr_text)} chars via {ocr_source}")

            image_ext = os.path.splitext(file_path)[1][1:]
            elements = []

            if ocr_text.strip():
                elements.append({
                    'type': 'text',
                    'content': ocr_text,
                    'metadata': {
                        'element_type': 'ocr_text',
                        'source': ocr_source,
                        'image_file': os.path.basename(file_path)
                    }
                })

            elements.append({
                'type': 'image',
                'content': '',
                'metadata': {
                    'element_type': 'image',
                    'ext': image_ext,
                    'size': len(image_bytes),
                    'has_ocr': len(ocr_text) > 0
                }
            })

            return {
                'elements': elements,
                'images': [{'local_path': file_path, 'ext': image_ext, 'metadata': {'size': len(image_bytes)}}],
                'tables': [],
                'metadata': {'ocr_performed': self.config.OCR_ENABLED, 'ocr_text_length': len(ocr_text)}
            }
        except Exception as e:
            logger.error(f"Error extracting image: {e}")
            raise

    def _extract_with_claude_vision(self, image_bytes: bytes, file_path: str) -> str:
        """Extract text using Claude Vision via Bedrock."""
        try:
            ext = os.path.splitext(file_path)[1].lower()
            media_types = {'.png': 'image/png', '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg', 
                          '.gif': 'image/gif', '.webp': 'image/webp', '.bmp': 'image/bmp'}
            media_type = media_types.get(ext, 'image/png')
            
            response = self.bedrock_client.invoke_model(
                modelId=CLAUDE_VISION_MODEL,
                body=json.dumps({
                    'anthropic_version': 'bedrock-2023-05-31',
                    'max_tokens': 4096,
                    'messages': [{
                        'role': 'user',
                        'content': [
                            {'type': 'image', 'source': {'type': 'base64', 'media_type': media_type, 
                             'data': base64.standard_b64encode(image_bytes).decode('utf-8')}},
                            {'type': 'text', 'text': 'Extract ALL text from this image including labels, titles, tables, and code. Format as clean readable text. If no text, respond "NO_TEXT".'}
                        ]
                    }]
                })
            )
            
            text = json.loads(response['body'].read()).get('content', [{}])[0].get('text', '')
            return "" if 'NO_TEXT' in text else text.strip()
        except Exception as e:
            logger.error(f"Claude Vision error: {e}")
            return ""

    def _extract_with_textract(self, image_bytes: bytes) -> str:
        """Extract text using AWS Textract (fallback)."""
        try:
            response = self.textract_client.detect_document_text(Document={'Bytes': image_bytes})
            return '\n'.join([b['Text'] for b in response.get('Blocks', []) if b['BlockType'] == 'LINE'])
        except Exception as e:
            logger.error(f"Textract error: {e}")
            return ""


# ============================================================================
# PARENT-CHILD CHUNKING (Small-to-Big Retrieval)
# ============================================================================

# Feature flag for parent-child chunking
ENABLE_PARENT_CHILD_CHUNKING = os.environ.get('ENABLE_PARENT_CHILD_CHUNKING', 'true').lower() == 'true'
CHILD_CHUNK_SENTENCES = int(os.environ.get('CHILD_CHUNK_SENTENCES', '3'))
CHILD_CHUNK_MIN_TOKENS = int(os.environ.get('CHILD_CHUNK_MIN_TOKENS', '50'))


class ParentChildChunker:
    """
    Parent-Child Chunking for Small-to-Big Retrieval.
    
    Research-backed approach (FLARE, Sentence Window):
    1. Index SMALL chunks (sentences) for HIGH PRECISION retrieval
    2. Link each child to PARENT (paragraph/section) for FULL CONTEXT
    3. At retrieval: search children → fetch parents → LLM gets full context
    
    Benefits:
    - Search precision: Small chunks match specific queries better
    - Generation context: LLM receives full surrounding context
    - Best of both worlds without embedding larger chunks
    """
    
    def __init__(self, 
                 sentences_per_child: int = None,
                 min_child_tokens: int = None,
                 config = None):
        self.sentences_per_child = sentences_per_child or CHILD_CHUNK_SENTENCES
        self.min_child_tokens = min_child_tokens or CHILD_CHUNK_MIN_TOKENS
        self.config = config
        
        # Initialize tokenizer if available
        if tiktoken is not None:
            self.tokenizer = tiktoken.get_encoding("cl100k_base")
        else:
            self.tokenizer = None
    
    def _count_tokens(self, text: str) -> int:
        """Count tokens in text."""
        if self.tokenizer:
            return len(self.tokenizer.encode(text))
        return len(text.split())  # Fallback to word count
    
    def create_parent_child_chunks(self, 
                                   parent_chunks: List[Dict], 
                                   document_id: str) -> Tuple[List[Dict], List[Dict]]:
        """
        Create child chunks from parent chunks with linking.
        
        Args:
            parent_chunks: List of standard chunks (the "parents")
            document_id: Document identifier
        
        Returns:
            Tuple of (parent_chunks, child_chunks) where:
            - parent_chunks: Original chunks marked as parents
            - child_chunks: Small sentence-level chunks linked to parents
        """
        if not ENABLE_PARENT_CHILD_CHUNKING:
            return parent_chunks, []
        
        parents = []
        children = []
        child_counter = 0
        
        for parent in parent_chunks:
            parent_id = parent.get('chunk_id')
            parent_content = parent.get('content', '')
            parent_type = parent.get('type', 'text')
            
            # Mark parent as a parent chunk
            parent_copy = parent.copy()
            parent_copy['is_parent'] = True
            parent_copy['child_ids'] = []
            parents.append(parent_copy)
            
            # Skip creating children for code/table chunks (keep as-is)
            if parent_type in ['code', 'table', 'image']:
                continue
            
            # Split parent into sentences
            sentences = self._split_into_sentences(parent_content)
            
            if len(sentences) < 2:
                # Too small for children
                continue
            
            # Create child chunks (groups of sentences)
            for i in range(0, len(sentences), self.sentences_per_child):
                child_sentences = sentences[i:i + self.sentences_per_child]
                child_content = ' '.join(child_sentences)
                
                # Skip if too small
                if self._count_tokens(child_content) < self.min_child_tokens:
                    continue
                
                child_counter += 1
                child_id = f"{document_id}-child-{child_counter:04d}"
                
                child_chunk = {
                    'chunk_id': child_id,
                    'document_id': document_id,
                    'parent_chunk_id': parent_id,  # KEY: Link to parent
                    'type': 'child',
                    'content': child_content,
                    'tokens': self._count_tokens(child_content),
                    'is_child': True,
                    'child_index': i // self.sentences_per_child,
                    'metadata': {
                        **parent.get('metadata', {}),
                        'parent_chunk_id': parent_id,
                        'sentence_range': f"{i+1}-{min(i + self.sentences_per_child, len(sentences))}",
                        'is_child': True
                    }
                }
                
                children.append(child_chunk)
                parent_copy['child_ids'].append(child_id)
        
        logger.info(f"Parent-Child chunking: {len(parents)} parents → {len(children)} children")
        
        return parents, children
    
    def _split_into_sentences(self, text: str) -> List[str]:
        """Split text into sentences using punctuation and newlines."""
        import re
        
        # Split on sentence boundaries
        # Handle: periods, exclamation, question marks, but preserve abbreviations
        sentence_pattern = r'(?<=[.!?])\s+(?=[A-Z])|(?<=\n)\s*(?=\S)'
        
        raw_sentences = re.split(sentence_pattern, text)
        
        # Clean and filter
        sentences = []
        for s in raw_sentences:
            s = s.strip()
            if s and len(s) > 10:  # Skip very short fragments
                sentences.append(s)
        
        return sentences
    
    def fetch_parent_for_child(self, 
                               child_chunk: Dict, 
                               all_parents: Dict[str, Dict]) -> Optional[Dict]:
        """
        Given a child chunk, retrieve its parent for full context.
        
        Args:
            child_chunk: The matched child chunk from search
            all_parents: Dictionary of parent_id → parent_chunk
        
        Returns:
            Parent chunk if found, None otherwise
        """
        parent_id = child_chunk.get('parent_chunk_id')
        if parent_id and parent_id in all_parents:
            return all_parents[parent_id]
        return None
    
    def expand_children_to_parents(self, 
                                   child_results: List[Dict],
                                   opensearch_client = None,
                                   index_name: str = None) -> List[Dict]:
        """
        Expand retrieved children to their parent chunks.
        
        This is the key retrieval step:
        1. Search returns small child chunks (high precision)
        2. We fetch the parent chunks for full context
        3. Deduplicate parents (multiple children may have same parent)
        
        Args:
            child_results: Search results (child chunks)
            opensearch_client: OpenSearch client for fetching parents
            index_name: OpenSearch index name
        
        Returns:
            Parent chunks with full context
        """
        if not child_results:
            return []
        
        # Collect unique parent IDs
        parent_ids = set()
        child_to_parent_score = {}  # Track best child score per parent
        
        for child in child_results:
            parent_id = child.get('parent_chunk_id') or child.get('metadata', {}).get('parent_chunk_id')
            
            if parent_id:
                parent_ids.add(parent_id)
                # Keep highest scoring child for each parent
                child_score = child.get('score', 0)
                if parent_id not in child_to_parent_score or child_score > child_to_parent_score[parent_id]:
                    child_to_parent_score[parent_id] = child_score
        
        if not parent_ids:
            # No parent IDs found, return children as-is
            return child_results
        
        # Fetch parent chunks from OpenSearch
        if opensearch_client and index_name:
            try:
                query = {
                    "query": {
                        "ids": {
                            "values": list(parent_ids)
                        }
                    },
                    "_source": {"excludes": ["embedding"]},
                    "size": len(parent_ids)
                }
                
                response = opensearch_client.search(index=index_name, body=query)
                hits = response.get('hits', {}).get('hits', [])
                
                parents = []
                for hit in hits:
                    source = hit['_source']
                    parent_id = source.get('chunk_id')
                    
                    parent_chunk = {
                        'chunk_id': parent_id,
                        'document_id': source.get('document_id'),
                        'text': source.get('text', source.get('content', '')),
                        'content': source.get('text', source.get('content', '')),
                        'type': source.get('type', 'text'),
                        'score': child_to_parent_score.get(parent_id, 0),  # Use child's score
                        'file_name': source.get('file_name'),
                        'file_type': source.get('file_type'),
                        'page_number': source.get('page_number'),
                        'section': source.get('section'),
                        'section_title': source.get('section_title'),
                        'metadata': source.get('metadata', {}),
                        'expanded_from_child': True
                    }
                    parents.append(parent_chunk)
                
                # Sort by score
                parents.sort(key=lambda x: x.get('score', 0), reverse=True)
                
                logger.info(f"Expanded {len(child_results)} children → {len(parents)} parent chunks")
                return parents
                
            except Exception as e:
                logger.error(f"Error fetching parent chunks: {e}")
                return child_results
        
        # Fallback: return children if can't fetch parents
        return child_results


# Initialize global parent-child chunker
parent_child_chunker = ParentChildChunker() if ENABLE_PARENT_CHILD_CHUNKING else None


# ============================================================================
# HYBRID CHUNKER (Recursive + Semantic + Parent-Child)
# ============================================================================

# Configuration for hybrid chunking
SECTION_MAX_TOKENS = int(os.environ.get('SECTION_MAX_TOKENS', '3500'))
CHUNK_TARGET_TOKENS = int(os.environ.get('CHUNK_TARGET_TOKENS', '1500'))
CHUNK_MAX_TOKENS = int(os.environ.get('CHUNK_MAX_TOKENS', '2000'))
SEMANTIC_THRESHOLD = float(os.environ.get('SEMANTIC_THRESHOLD', '0.7'))
HYBRID_CHUNK_OVERLAP = int(os.environ.get('HYBRID_CHUNK_OVERLAP', '400'))


@dataclass
class Section:
    """Represents a document section with its elements."""
    heading: str
    heading_level: int
    heading_path: List[str]
    elements: List[Dict]
    tables: List[Dict]
    images: List[Dict]
    page_range: List[int]


@dataclass
class ChunkHierarchy:
    """Container for multi-level chunk hierarchy."""
    section_chunks: List[Dict] = field(default_factory=list)
    semantic_chunks: List[Dict] = field(default_factory=list)
    child_chunks: List[Dict] = field(default_factory=list)
    table_chunks: List[Dict] = field(default_factory=list)
    image_chunks: List[Dict] = field(default_factory=list)
    code_chunks: List[Dict] = field(default_factory=list)
    
    def all_chunks(self) -> List[Dict]:
        """Return all chunks flattened."""
        return (
            self.section_chunks + 
            self.semantic_chunks + 
            self.child_chunks + 
            self.table_chunks + 
            self.image_chunks +
            self.code_chunks
        )


class HybridChunker:
    """
    Multi-strategy chunker combining:
    1. RECURSIVE: Respect document hierarchy (H1→H2→H3→paragraph)
    2. SEMANTIC: Detect topic shifts within sections
    3. PARENT-CHILD: Create searchable children linked to context-rich parents
    
    This chunker produces a 3-level hierarchy:
    - SECTION chunks: Full sections for broad context
    - SEMANTIC chunks: Topic-based chunks for retrieval  
    - CHILD chunks: Sentence-level for precision search
    """
    
    def __init__(self, config=None):
        self.section_max_tokens = config.SECTION_MAX_TOKENS if config else SECTION_MAX_TOKENS
        self.chunk_target_tokens = config.CHUNK_TARGET_TOKENS if config else CHUNK_TARGET_TOKENS
        self.chunk_max_tokens = config.CHUNK_MAX_TOKENS if config else CHUNK_MAX_TOKENS
        self.semantic_threshold = config.SEMANTIC_THRESHOLD if config else SEMANTIC_THRESHOLD
        self.chunk_overlap = config.HYBRID_CHUNK_OVERLAP if config else HYBRID_CHUNK_OVERLAP
        self.child_sentences = CHILD_CHUNK_SENTENCES
        
        # Initialize tokenizer
        if tiktoken is not None:
            self.tokenizer = tiktoken.get_encoding("cl100k_base")
        else:
            self.tokenizer = None
        
        self.chunk_counter = 0
        self.section_counter = 0
        self.child_counter = 0
    
    def _count_tokens(self, text: str) -> int:
        """Count tokens in text."""
        if not text:
            return 0
        if self.tokenizer:
            return len(self.tokenizer.encode(text))
        return len(text.split())
    
    def chunk_document(self, 
                       document_id: str,
                       elements: List[Dict], 
                       file_type: str,
                       document_metadata: Dict = None) -> ChunkHierarchy:
        """
        Main entry point - produces 3-level chunk hierarchy.
        
        Args:
            document_id: Unique document identifier
            elements: Extracted document elements
            file_type: File extension (pdf, docx, xlsx, etc.)
            document_metadata: Additional document metadata
        
        Returns:
            ChunkHierarchy with section, semantic, child, table, image chunks
        """
        document_metadata = document_metadata or {}
        self.chunk_counter = 0
        self.section_counter = 0
        self.child_counter = 0
        
        logger.info(f"HybridChunker: Processing {len(elements)} elements from {document_id} ({file_type})")
        
        # STEP 1: Group elements by heading hierarchy
        sections = self._group_by_sections(elements, file_type)
        
        all_chunks = ChunkHierarchy()
        
        for section in sections:
            # STEP 2: Create SECTION CHUNK (Level 1 - Parent)
            section_chunk = self._create_section_chunk(section, document_id, document_metadata)
            all_chunks.section_chunks.append(section_chunk)
            
            # STEP 3: Apply RECURSIVE + SEMANTIC chunking (Level 2)
            semantic_chunks = self._recursive_semantic_chunk(
                section.elements,
                parent_id=section_chunk['chunk_id'],
                section_id=section_chunk['chunk_id'],
                heading_path=section.heading_path,
                document_id=document_id,
                document_metadata=document_metadata
            )
            all_chunks.semantic_chunks.extend(semantic_chunks)
            
            # Update section with child references
            section_chunk['child_ids'] = [c['chunk_id'] for c in semantic_chunks]
            
            # STEP 4: Create CHILD chunks (Level 3 - highest precision)
            for semantic_chunk in semantic_chunks:
                children = self._create_child_chunks(
                    semantic_chunk,
                    document_id=document_id
                )
                all_chunks.child_chunks.extend(children)
                semantic_chunk['child_ids'] = [c['chunk_id'] for c in children]
            
            # STEP 5: Handle TABLES separately
            for table in section.tables:
                table_chunk = self._chunk_table(
                    table, 
                    parent_id=section_chunk['chunk_id'],
                    heading_path=section.heading_path,
                    document_id=document_id,
                    document_metadata=document_metadata
                )
                all_chunks.table_chunks.append(table_chunk)
            
            # STEP 6: Handle IMAGES separately
            for image in section.images:
                image_chunk = self._chunk_image(
                    image,
                    parent_id=section_chunk['chunk_id'],
                    heading_path=section.heading_path,
                    document_id=document_id,
                    document_metadata=document_metadata
                )
                all_chunks.image_chunks.append(image_chunk)
        
        logger.info(f"HybridChunker: Created {len(all_chunks.section_chunks)} sections, "
                    f"{len(all_chunks.semantic_chunks)} semantic, "
                    f"{len(all_chunks.child_chunks)} children, "
                    f"{len(all_chunks.table_chunks)} tables, "
                    f"{len(all_chunks.image_chunks)} images")
        
        return all_chunks
    
    def _group_by_sections(self, elements: List[Dict], file_type: str) -> List[Section]:
        """Group elements by heading hierarchy."""
        sections = []
        current_section = None
        heading_stack = []
        
        for element in elements:
            elem_type = element.get('type', 'text')
            elem_meta = element.get('metadata', {})
            content = element.get('content', '')
            
            # Check if this is a heading
            element_type = elem_meta.get('element_type', '')
            is_heading = element_type in ['Title', 'title', 'heading', 'Heading']
            heading_level = elem_meta.get('heading_level', 1) if is_heading else 0
            
            if is_heading or element_type == 'Title':
                # Start new section
                if current_section:
                    sections.append(current_section)
                
                # Update heading stack
                if heading_level == 1 or element_type == 'Title':
                    heading_stack = [content.strip()[:150]]
                else:
                    heading_stack = heading_stack[:heading_level - 1] + [content.strip()[:150]]
                
                current_section = Section(
                    heading=content.strip()[:150],
                    heading_level=heading_level or 1,
                    heading_path=heading_stack.copy(),
                    elements=[],
                    tables=[],
                    images=[],
                    page_range=[]
                )
            else:
                # Add element to current section
                if current_section is None:
                    # Create default section for content before first heading
                    current_section = Section(
                        heading="Document Content",
                        heading_level=1,
                        heading_path=["Document Content"],
                        elements=[],
                        tables=[],
                        images=[],
                        page_range=[]
                    )
                
                if elem_type == 'table':
                    current_section.tables.append(element)
                elif elem_type == 'image':
                    current_section.images.append(element)
                else:
                    current_section.elements.append(element)
                
                # Track page numbers
                page = elem_meta.get('page_no') or elem_meta.get('page_number')
                if page and page not in current_section.page_range:
                    current_section.page_range.append(page)
        
        # Don't forget last section
        if current_section:
            sections.append(current_section)
        
        return sections
    
    def _create_section_chunk(self, section: Section, document_id: str, document_metadata: Dict) -> Dict:
        """Create a section-level chunk (full section content)."""
        self.section_counter += 1
        chunk_id = f"{document_id}-section-{self.section_counter:04d}"
        
        # Combine all element content
        content_parts = [section.heading]
        for elem in section.elements:
            content_parts.append(elem.get('content', ''))
        
        # Add table summaries
        for table in section.tables:
            caption = table.get('metadata', {}).get('caption', 'Table')
            content_parts.append(f"[Table: {caption}]")
        
        # Add image descriptions
        for image in section.images:
            desc = image.get('metadata', {}).get('description', 'Image')
            content_parts.append(f"[Image: {desc}]")
        
        content = '\n\n'.join(content_parts)
        
        return {
            'chunk_id': chunk_id,
            'document_id': document_id,
            'type': 'section',
            'content': content,
            'tokens': self._count_tokens(content),
            'heading': section.heading,
            'heading_path': section.heading_path,
            'heading_level': section.heading_level,
            'page_numbers': sorted(section.page_range),
            'child_ids': [],
            'has_tables': len(section.tables) > 0,
            'has_images': len(section.images) > 0,
            'is_section': True,
            'metadata': {
                **document_metadata,
                'section_heading': section.heading,
                'heading_path': section.heading_path
            }
        }
    
    def _recursive_semantic_chunk(self, 
                                    elements: List[Dict],
                                    parent_id: str,
                                    section_id: str,
                                    heading_path: List[str],
                                    document_id: str,
                                    document_metadata: Dict) -> List[Dict]:
        """Apply recursive splitting first, then semantic within each part."""
        if not elements:
            return []
        
        combined_text = '\n'.join([e.get('content', '') for e in elements])
        token_count = self._count_tokens(combined_text)
        
        # If small enough, create single chunk
        if token_count <= self.chunk_max_tokens:
            return [self._create_semantic_chunk(
                combined_text, parent_id, section_id, heading_path, 
                document_id, document_metadata, elements
            )]
        
        # Split by semantic breaks (topic shifts)
        return self._semantic_split_with_overlap(
            elements, parent_id, section_id, heading_path, 
            document_id, document_metadata
        )
    
    def _semantic_split_with_overlap(self,
                                      elements: List[Dict],
                                      parent_id: str,
                                      section_id: str,
                                      heading_path: List[str],
                                      document_id: str,
                                      document_metadata: Dict) -> List[Dict]:
        """Split by detecting semantic breaks between elements."""
        chunks = []
        current_elements = []
        current_tokens = 0
        prev_element = None
        
        for element in elements:
            content = element.get('content', '')
            elem_tokens = self._count_tokens(content)
            
            # Check if adding this element would exceed max
            if current_tokens + elem_tokens > self.chunk_max_tokens and current_elements:
                # Create chunk from current elements
                chunk_text = '\n'.join([e.get('content', '') for e in current_elements])
                chunks.append(self._create_semantic_chunk(
                    chunk_text, parent_id, section_id, heading_path,
                    document_id, document_metadata, current_elements
                ))
                
                # Start new chunk with overlap (last element for context continuity)
                if current_elements:
                    current_elements = [current_elements[-1]]
                    current_tokens = self._count_tokens(current_elements[0].get('content', ''))
                else:
                    current_elements = []
                    current_tokens = 0
            
            # Check for semantic break (heading, page break, format change, etc.)
            is_break = self._is_semantic_break(element, prev_element)
            
            if is_break and current_tokens > self.chunk_target_tokens and current_elements:
                # Create chunk at semantic break
                chunk_text = '\n'.join([e.get('content', '') for e in current_elements])
                chunks.append(self._create_semantic_chunk(
                    chunk_text, parent_id, section_id, heading_path,
                    document_id, document_metadata, current_elements
                ))
                current_elements = []
                current_tokens = 0
            
            current_elements.append(element)
            current_tokens += elem_tokens
            prev_element = element
        
        # Don't forget last chunk
        if current_elements:
            chunk_text = '\n'.join([e.get('content', '') for e in current_elements])
            chunks.append(self._create_semantic_chunk(
                chunk_text, parent_id, section_id, heading_path,
                document_id, document_metadata, current_elements
            ))
        
        return chunks
    
    def _is_semantic_break(self, element: Dict, prev_element: Dict = None) -> bool:
        """
        Enhanced semantic break detection for better topic identification.
        
        Detects breaks based on:
        1. Element type (heading, title, page break)
        2. Format changes (bold, different style)
        3. Content patterns (numbered sections, keywords)
        4. Structural cues (blank lines, horizontal rules)
        """
        elem_meta = element.get('metadata', {})
        content = element.get('content', '').strip()
        element_type = elem_meta.get('element_type', '')
        
        # Type-based breaks (strongest signal)
        break_types = {
            'Title', 'title', 'heading', 'Heading', 
            'PageBreak', 'page_break',
            'Header', 'header',
            'SectionHeading', 'section_heading'
        }
        if element_type in break_types:
            return True
        
        # List items can be breaks if they start a new list
        if element_type in ['ListItem', 'list_item']:
            if prev_element:
                prev_type = prev_element.get('metadata', {}).get('element_type', '')
                if prev_type not in ['ListItem', 'list_item']:
                    return True  # First item in a new list
        
        # Content pattern detection
        if content:
            # Numbered section headers: "1.", "1.1", "Section 1", "Chapter X"
            if re.match(r'^(?:\d+\.)+\s*\w|^(?:Section|Chapter|Part|Article)\s+\d+', content, re.IGNORECASE):
                return True
            
            # All-caps headers (common in mainframe docs)
            if len(content) < 100 and content.isupper() and len(content.split()) <= 10:
                return True
            
            # Colon-ending short lines (definition/label style)
            if len(content) < 80 and content.endswith(':') and '\n' not in content:
                return True
            
            # Lines starting with common section markers
            section_markers = [
                r'^={3,}',  # === dividers
                r'^-{3,}',  # --- dividers
                r'^\*{3,}',  # *** dividers
                r'^#{1,6}\s',  # Markdown headers
                r'^(?:NOTE|WARNING|CAUTION|IMPORTANT|TODO|FIXME)\s*:',
                r'^(?:Purpose|Description|Parameters|Returns|Example|Usage)\s*:',
            ]
            for marker in section_markers:
                if re.match(marker, content, re.IGNORECASE):
                    return True
        
        # Format-based breaks
        is_bold = elem_meta.get('is_bold', False) or elem_meta.get('bold', False)
        prev_bold = prev_element.get('metadata', {}).get('is_bold', False) if prev_element else False
        
        # Bold text after non-bold can indicate a new topic
        if is_bold and not prev_bold and len(content) < 150:
            return True
        
        return False
    
    def _create_semantic_chunk(self,
                                content: str,
                                parent_id: str,
                                section_id: str,
                                heading_path: List[str],
                                document_id: str,
                                document_metadata: Dict,
                                elements: List[Dict] = None) -> Dict:
        """Create a semantic-level chunk."""
        self.chunk_counter += 1
        chunk_id = f"{document_id}-chunk-{self.chunk_counter:04d}"
        
        # Extract page numbers from elements
        page_numbers = set()
        if elements:
            for elem in elements:
                page = elem.get('metadata', {}).get('page_no') or elem.get('metadata', {}).get('page_number')
                if page:
                    page_numbers.add(page)
        
        return {
            'chunk_id': chunk_id,
            'document_id': document_id,
            'type': 'semantic',
            'content': content,
            'tokens': self._count_tokens(content),
            'parent_chunk_id': parent_id,
            'section_id': section_id,
            'heading_path': heading_path,
            'page_numbers': sorted(list(page_numbers)),
            'child_ids': [],
            'is_semantic': True,
            'metadata': {
                **document_metadata,
                'heading_path': heading_path,
                'parent_chunk_id': parent_id,
                'section_id': section_id
            }
        }
    
    def _create_child_chunks(self, semantic_chunk: Dict, document_id: str) -> List[Dict]:
        """Create sentence-level children for high-precision search."""
        content = semantic_chunk.get('content', '')
        sentences = self._split_sentences(content)
        
        if len(sentences) < 2:
            return []  # Too small to create children
        
        children = []
        for i in range(0, len(sentences), self.child_sentences):
            child_sentences = sentences[i:i + self.child_sentences]
            child_text = ' '.join(child_sentences)
            
            # Skip very small fragments
            if self._count_tokens(child_text) < CHILD_CHUNK_MIN_TOKENS:
                continue
            
            self.child_counter += 1
            child_id = f"{document_id}-child-{self.child_counter:04d}"
            
            children.append({
                'chunk_id': child_id,
                'document_id': document_id,
                'type': 'child',
                'content': child_text,
                'tokens': self._count_tokens(child_text),
                'parent_chunk_id': semantic_chunk['chunk_id'],
                'section_id': semantic_chunk.get('section_id'),
                'heading_path': semantic_chunk.get('heading_path', []),
                'is_child': True,
                'position_in_parent': i // self.child_sentences,
                'metadata': {
                    'parent_chunk_id': semantic_chunk['chunk_id'],
                    'section_id': semantic_chunk.get('section_id'),
                    'heading_path': semantic_chunk.get('heading_path', []),
                    'is_child': True
                }
            })
        
        return children
    
    def _split_sentences(self, text: str) -> List[str]:
        """Split text into sentences."""
        import re
        # Split on sentence boundaries
        sentences = re.split(r'(?<=[.!?])\s+(?=[A-Z])', text)
        return [s.strip() for s in sentences if s.strip() and len(s.strip()) > 10]
    
    def _chunk_table(self, table: Dict, parent_id: str, heading_path: List[str],
                     document_id: str, document_metadata: Dict) -> Dict:
        """Create a table chunk with searchable content."""
        self.chunk_counter += 1
        chunk_id = f"{document_id}-table-{self.chunk_counter:04d}"
        
        content = table.get('content', '')
        table_meta = table.get('metadata', {})
        caption = table_meta.get('caption', table_meta.get('title', 'Table'))
        
        # Create searchable summary
        searchable = f"Table: {caption}\n{content}"
        
        return {
            'chunk_id': chunk_id,
            'document_id': document_id,
            'type': 'table',
            'content': searchable,
            'tokens': self._count_tokens(searchable),
            'parent_chunk_id': parent_id,
            'heading_path': heading_path,
            'table_caption': caption,
            'num_rows': table_meta.get('num_rows', 0),
            'num_cols': table_meta.get('num_cols', 0),
            'is_table': True,
            'metadata': {
                **document_metadata,
                'parent_chunk_id': parent_id,
                'heading_path': heading_path,
                'is_table': True,
                'table_caption': caption
            }
        }
    
    def _chunk_image(self, image: Dict, parent_id: str, heading_path: List[str],
                     document_id: str, document_metadata: Dict) -> Dict:
        """Create an image chunk with description."""
        self.chunk_counter += 1
        chunk_id = f"{document_id}-image-{self.chunk_counter:04d}"
        
        image_meta = image.get('metadata', {})
        description = image_meta.get('description', '')
        caption = image_meta.get('caption', 'Image')
        ocr_text = image_meta.get('ocr_text', '')
        
        # Create searchable content
        content_parts = [f"Image: {caption}"]
        if description:
            content_parts.append(f"Description: {description}")
        if ocr_text:
            content_parts.append(f"Text in image: {ocr_text}")
        
        content = '\n'.join(content_parts)
        
        return {
            'chunk_id': chunk_id,
            'document_id': document_id,
            'type': 'image',
            'content': content,
            'tokens': self._count_tokens(content),
            'parent_chunk_id': parent_id,
            'heading_path': heading_path,
            'image_caption': caption,
            'is_image': True,
            'metadata': {
                **document_metadata,
                'parent_chunk_id': parent_id,
                'heading_path': heading_path,
                'is_image': True,
                'image_caption': caption,
                'has_ocr': bool(ocr_text)
            }
        }


# Initialize global hybrid chunker
hybrid_chunker = HybridChunker()


# ============================================================================
# MAINFRAME CODE CHUNKER (Assembler, JCL, COBOL, Copybooks)
# ============================================================================

class MainframeCodeChunker:
    """
    Specialized chunker for mainframe code files with rich metadata extraction.
    
    Supports:
    - Assembler (.asm) - procedure/CSECT boundaries
    - JCL (.jcl) - JOB/STEP/PROC boundaries  
    - COBOL (.cbl, .cob) - DIVISION/SECTION/PARAGRAPH
    - Copybooks (.cpy) - data structure definitions
    - PROC/CLIST - procedure libraries
    - REXX (.rexx, .rex) - procedure calls
    
    Extracts:
    - Procedure names and boundaries
    - Labels and symbols
    - Table references (DSECT, USING)
    - CALL graph (external calls)
    - Register usage patterns
    - Data structure definitions
    """
    
    def __init__(self, target_lines: int = 150, overlap_lines: int = 20):
        self.target_lines = target_lines
        self.overlap_lines = overlap_lines
        
        # Language-specific patterns
        self.patterns = {
            'assembler': {
                'procedure': re.compile(r'^(\w+)\s+(PROC|CSECT|DSECT)\b', re.IGNORECASE | re.MULTILINE),
                'label': re.compile(r'^([A-Z_][A-Z0-9_]*)\s', re.IGNORECASE | re.MULTILINE),
                'table_ref': re.compile(r'\bUSING\s+(\w+)', re.IGNORECASE),
                'call': re.compile(r'\b(CALL|BAL|BALR|BAS|BASR|BASSM)\s+(\w+)', re.IGNORECASE),
                'macro': re.compile(r'^\s+(\w+)\s+MACRO\b', re.IGNORECASE | re.MULTILINE),
                'end': re.compile(r'^\w*\s+END\b', re.IGNORECASE | re.MULTILINE)
            },
            'jcl': {
                'job': re.compile(r'^//(\w+)\s+JOB\b', re.IGNORECASE | re.MULTILINE),
                'step': re.compile(r'^//(\w+)\s+EXEC\s+(\w+)', re.IGNORECASE | re.MULTILINE),
                'proc': re.compile(r'^//(\w+)\s+PROC\b', re.IGNORECASE | re.MULTILINE),
                'dd': re.compile(r'^//(\w+)\s+DD\b', re.IGNORECASE | re.MULTILINE),
                'dataset': re.compile(r"DSN=(['\w\.]+)", re.IGNORECASE)
            },
            'cobol': {
                'division': re.compile(r'^\s*(\w+)\s+DIVISION\b', re.IGNORECASE | re.MULTILINE),
                'section': re.compile(r'^\s*(\w+)\s+SECTION\b', re.IGNORECASE | re.MULTILINE),
                'paragraph': re.compile(r'^(\s{0,3})([A-Z0-9][\w-]*)\.\s*$', re.IGNORECASE | re.MULTILINE),
                'copy': re.compile(r'\bCOPY\s+(\w+)', re.IGNORECASE),
                'call': re.compile(r'\bCALL\s+[\'"]?(\w+)[\'"]?', re.IGNORECASE)
            },
            'copybook': {
                'level': re.compile(r'^\s*(\d{1,2})\s+(\w+)', re.MULTILINE),
                'pic': re.compile(r'PIC\s+(\S+)', re.IGNORECASE),
                'redefines': re.compile(r'REDEFINES\s+(\w+)', re.IGNORECASE)
            }
        }
    
    def detect_language(self, content: str, file_extension: str) -> str:
        """Detect mainframe language from content and extension."""
        ext = file_extension.lower().lstrip('.')
        
        if ext in ['asm', 's', 'mac']:
            return 'assembler'
        elif ext in ['jcl', 'job', 'prc']:
            return 'jcl'
        elif ext in ['cbl', 'cob', 'cobol']:
            return 'cobol'
        elif ext in ['cpy', 'copy']:
            return 'copybook'
        elif ext in ['rexx', 'rex', 'exec']:
            return 'rexx'
        elif ext in ['clist', 'proc']:
            return 'clist'
        
        # Content-based detection
        content_upper = content.upper()[:2000]
        if 'CSECT' in content_upper or 'DSECT' in content_upper:
            return 'assembler'
        elif content_upper.lstrip().startswith('//'):
            return 'jcl'
        elif 'IDENTIFICATION DIVISION' in content_upper or 'PROCEDURE DIVISION' in content_upper:
            return 'cobol'
        elif content.lstrip().startswith('/*'):
            return 'rexx'
        
        return 'assembler'  # Default
    
    def chunk_code(self, 
                   content: str, 
                   document_id: str,
                   file_extension: str,
                   file_name: str,
                   document_metadata: Dict = None) -> List[Dict]:
        """
        Chunk mainframe code with procedure-boundary awareness and rich metadata.
        """
        document_metadata = document_metadata or {}
        language = self.detect_language(content, file_extension)
        lines = content.split('\n')
        total_lines = len(lines)
        
        logger.info(f"MainframeCodeChunker: {total_lines} lines, language={language}")
        
        # Extract code structure
        if language == 'assembler':
            procedures = self._extract_assembler_procedures(lines)
            metadata = self._extract_assembler_metadata(content, lines)
        elif language == 'jcl':
            procedures = self._extract_jcl_steps(lines)
            metadata = self._extract_jcl_metadata(content, lines)
        elif language == 'cobol':
            procedures = self._extract_cobol_sections(lines)
            metadata = self._extract_cobol_metadata(content, lines)
        elif language == 'copybook':
            procedures = []
            metadata = self._extract_copybook_metadata(content, lines)
        else:
            procedures = self._extract_generic_procedures(lines)
            metadata = {}
        
        # Create chunks respecting procedure boundaries
        chunks = self._create_procedure_aware_chunks(
            lines=lines,
            procedures=procedures,
            document_id=document_id,
            language=language,
            file_name=file_name,
            code_metadata=metadata,
            document_metadata=document_metadata
        )
        
        return chunks
    
    def _extract_assembler_procedures(self, lines: List[str]) -> List[Dict]:
        """Extract PROC/CSECT/DSECT boundaries from assembler code."""
        procedures = []
        current_proc = None
        
        for i, line in enumerate(lines):
            line_upper = line.upper()
            
            # Check for PROC/CSECT/DSECT start
            for marker in ['CSECT', 'DSECT', 'PROC']:
                if f' {marker}' in line_upper or line_upper.endswith(marker):
                    if current_proc:
                        current_proc['end_line'] = i - 1
                        procedures.append(current_proc)
                    
                    # Extract name (first non-empty token)
                    parts = line.split()
                    name = parts[0] if parts and not parts[0].startswith('*') else f'{marker}_{i}'
                    
                    current_proc = {
                        'name': name,
                        'type': marker,
                        'start_line': i,
                        'end_line': None
                    }
                    break
            
            # Check for END
            if current_proc and (' END ' in line_upper or line_upper.strip() == 'END'):
                current_proc['end_line'] = i
                procedures.append(current_proc)
                current_proc = None
        
        if current_proc:
            current_proc['end_line'] = len(lines) - 1
            procedures.append(current_proc)
        
        return procedures
    
    def _extract_assembler_metadata(self, content: str, lines: List[str]) -> Dict:
        """Extract rich metadata from assembler code."""
        patterns = self.patterns['assembler']
        
        # Labels
        labels = []
        for match in patterns['label'].finditer(content):
            label = match.group(1)
            if label not in labels and len(label) <= 8:  # Assembler label max length
                labels.append(label)
        
        # External calls
        calls = []
        for match in patterns['call'].finditer(content):
            call_type = match.group(1)
            target = match.group(2)
            calls.append({'type': call_type, 'target': target})
        
        # Table references (USING)
        table_refs = [m.group(1) for m in patterns['table_ref'].finditer(content)]
        
        # Register usage pattern (R0-R15)
        registers_used = set()
        reg_pattern = re.compile(r'\bR(\d{1,2})\b|\b(\d{1,2}),')
        for match in reg_pattern.finditer(content):
            reg = match.group(1) or match.group(2)
            if reg and int(reg) <= 15:
                registers_used.add(f'R{reg}')
        
        return {
            'labels': labels[:100],
            'calls': calls[:50],
            'table_refs': list(set(table_refs))[:30],
            'registers_used': sorted(list(registers_used)),
            'has_macros': bool(patterns['macro'].search(content))
        }
    
    def _extract_jcl_steps(self, lines: List[str]) -> List[Dict]:
        """Extract JOB/EXEC step boundaries from JCL."""
        procedures = []
        current_step = None
        
        for i, line in enumerate(lines):
            if not line.startswith('//'):
                continue
            
            line_upper = line.upper()
            
            # JOB card
            if ' JOB ' in line_upper:
                if current_step:
                    current_step['end_line'] = i - 1
                    procedures.append(current_step)
                
                parts = line.split()
                name = parts[0][2:] if len(parts) > 0 else f'JOB_{i}'
                current_step = {'name': name, 'type': 'JOB', 'start_line': i, 'end_line': None}
            
            # EXEC step
            elif ' EXEC ' in line_upper:
                if current_step and current_step['type'] != 'JOB':
                    current_step['end_line'] = i - 1
                    procedures.append(current_step)
                
                parts = line.split()
                name = parts[0][2:] if len(parts) > 0 else f'STEP_{i}'
                
                # Extract program/proc name
                exec_match = re.search(r'EXEC\s+(?:PGM=)?(\w+)', line_upper)
                pgm = exec_match.group(1) if exec_match else 'UNKNOWN'
                
                current_step = {
                    'name': name,
                    'type': 'STEP',
                    'program': pgm,
                    'start_line': i,
                    'end_line': None
                }
        
        if current_step:
            current_step['end_line'] = len(lines) - 1
            procedures.append(current_step)
        
        return procedures
    
    def _extract_jcl_metadata(self, content: str, lines: List[str]) -> Dict:
        """Extract metadata from JCL."""
        patterns = self.patterns['jcl']
        
        # Datasets referenced
        datasets = [m.group(1).strip("'") for m in patterns['dataset'].finditer(content)]
        
        # DD names
        dd_names = [m.group(1) for m in patterns['dd'].finditer(content)]
        
        # Programs executed
        programs = []
        for m in patterns['step'].finditer(content):
            pgm = m.group(2)
            if pgm not in programs:
                programs.append(pgm)
        
        return {
            'datasets': list(set(datasets))[:50],
            'dd_names': list(set(dd_names))[:50],
            'programs': programs[:30]
        }
    
    def _extract_cobol_sections(self, lines: List[str]) -> List[Dict]:
        """Extract DIVISION/SECTION boundaries from COBOL."""
        procedures = []
        current_section = None
        
        content = '\n'.join(lines)
        patterns = self.patterns['cobol']
        
        # Find divisions and sections
        for match in patterns['division'].finditer(content):
            name = match.group(1)
            start = content[:match.start()].count('\n')
            
            if current_section:
                current_section['end_line'] = start - 1
                procedures.append(current_section)
            
            current_section = {
                'name': f'{name} DIVISION',
                'type': 'DIVISION',
                'start_line': start,
                'end_line': None
            }
        
        for match in patterns['section'].finditer(content):
            name = match.group(1)
            start = content[:match.start()].count('\n')
            
            if current_section:
                current_section['end_line'] = start - 1
                procedures.append(current_section)
            
            current_section = {
                'name': f'{name} SECTION',
                'type': 'SECTION',
                'start_line': start,
                'end_line': None
            }
        
        if current_section:
            current_section['end_line'] = len(lines) - 1
            procedures.append(current_section)
        
        return procedures
    
    def _extract_cobol_metadata(self, content: str, lines: List[str]) -> Dict:
        """Extract metadata from COBOL."""
        patterns = self.patterns['cobol']
        
        # COPY statements
        copybooks = [m.group(1) for m in patterns['copy'].finditer(content)]
        
        # CALL statements
        calls = [m.group(1) for m in patterns['call'].finditer(content)]
        
        # Paragraphs
        paragraphs = [m.group(2) for m in patterns['paragraph'].finditer(content)]
        
        return {
            'copybooks': list(set(copybooks)),
            'calls': list(set(calls)),
            'paragraphs': paragraphs[:50]
        }
    
    def _extract_copybook_metadata(self, content: str, lines: List[str]) -> Dict:
        """Extract data structure info from copybook."""
        patterns = self.patterns['copybook']
        
        # Data levels and names
        data_items = []
        for match in patterns['level'].finditer(content):
            level = match.group(1)
            name = match.group(2)
            data_items.append({'level': int(level), 'name': name})
        
        # PIC clauses
        pics = [m.group(1) for m in patterns['pic'].finditer(content)]
        
        # REDEFINES
        redefines = [m.group(1) for m in patterns['redefines'].finditer(content)]
        
        # Find 01 level items (top-level structures)
        top_level = [d['name'] for d in data_items if d['level'] == 1]
        
        return {
            'data_items': data_items[:100],
            'top_level_structures': top_level,
            'redefines': list(set(redefines)),
            'has_numeric': any('9' in p for p in pics),
            'has_alpha': any('X' in p for p in pics)
        }
    
    def _extract_generic_procedures(self, lines: List[str]) -> List[Dict]:
        """Generic procedure extraction for other code types."""
        return []
    
    def _create_procedure_aware_chunks(self,
                                        lines: List[str],
                                        procedures: List[Dict],
                                        document_id: str,
                                        language: str,
                                        file_name: str,
                                        code_metadata: Dict,
                                        document_metadata: Dict) -> List[Dict]:
        """Create chunks respecting procedure boundaries."""
        chunks = []
        chunk_counter = 0
        total_lines = len(lines)
        
        if not procedures:
            # No procedures found - chunk by line count
            return self._chunk_by_lines(
                lines, document_id, language, file_name,
                code_metadata, document_metadata
            )
        
        # Sort procedures by start line
        procedures = sorted(procedures, key=lambda p: p['start_line'])
        
        for proc in procedures:
            start = proc['start_line']
            end = proc.get('end_line', total_lines - 1)
            proc_lines = lines[start:end + 1]
            proc_content = '\n'.join(proc_lines)
            
            # If procedure is small enough, single chunk
            if len(proc_lines) <= self.target_lines:
                chunk_counter += 1
                chunks.append({
                    'chunk_id': f"{document_id}-code-{chunk_counter:04d}",
                    'document_id': document_id,
                    'type': 'code',
                    'content': proc_content,
                    'tokens': len(proc_content.split()),  # Approximate
                    'metadata': {
                        'language': language,
                        'procedure_name': proc['name'],
                        'procedure_type': proc['type'],
                        'start_line': start + 1,
                        'end_line': end + 1,
                        'line_count': len(proc_lines),
                        'is_complete_procedure': True,
                        'file_name': file_name,
                        **code_metadata,
                        **document_metadata
                    }
                })
            else:
                # Split large procedure
                sub_chunks = self._split_large_procedure(
                    proc_lines, proc, document_id, language,
                    file_name, code_metadata, document_metadata,
                    chunk_counter, start
                )
                chunk_counter += len(sub_chunks)
                chunks.extend(sub_chunks)
        
        return chunks
    
    def _split_large_procedure(self,
                                lines: List[str],
                                proc: Dict,
                                document_id: str,
                                language: str,
                                file_name: str,
                                code_metadata: Dict,
                                document_metadata: Dict,
                                base_counter: int,
                                global_start: int) -> List[Dict]:
        """Split a large procedure into smaller chunks with overlap."""
        chunks = []
        counter = base_counter
        i = 0
        total_parts = (len(lines) + self.target_lines - 1) // self.target_lines
        part = 0
        
        while i < len(lines):
            end = min(i + self.target_lines, len(lines))
            chunk_lines = lines[i:end]
            
            counter += 1
            part += 1
            
            chunks.append({
                'chunk_id': f"{document_id}-code-{counter:04d}",
                'document_id': document_id,
                'type': 'code',
                'content': '\n'.join(chunk_lines),
                'tokens': len('\n'.join(chunk_lines).split()),
                'metadata': {
                    'language': language,
                    'procedure_name': proc['name'],
                    'procedure_type': proc['type'],
                    'start_line': global_start + i + 1,
                    'end_line': global_start + end,
                    'line_count': len(chunk_lines),
                    'is_complete_procedure': False,
                    'part': part,
                    'total_parts': total_parts,
                    'file_name': file_name,
                    **code_metadata,
                    **document_metadata
                }
            })
            
            i = end - self.overlap_lines
        
        return chunks
    
    def _chunk_by_lines(self,
                        lines: List[str],
                        document_id: str,
                        language: str,
                        file_name: str,
                        code_metadata: Dict,
                        document_metadata: Dict) -> List[Dict]:
        """Fallback: chunk by line count."""
        chunks = []
        counter = 0
        i = 0
        
        while i < len(lines):
            end = min(i + self.target_lines, len(lines))
            chunk_lines = lines[i:end]
            
            counter += 1
            chunks.append({
                'chunk_id': f"{document_id}-code-{counter:04d}",
                'document_id': document_id,
                'type': 'code',
                'content': '\n'.join(chunk_lines),
                'tokens': len('\n'.join(chunk_lines).split()),
                'metadata': {
                    'language': language,
                    'start_line': i + 1,
                    'end_line': end,
                    'line_count': len(chunk_lines),
                    'file_name': file_name,
                    **code_metadata,
                    **document_metadata
                }
            })
            
            i = end - self.overlap_lines
        
        return chunks


# Initialize global mainframe code chunker
mainframe_code_chunker = MainframeCodeChunker()


# ============================================================================
# EMBEDDED PRODUCTION CHUNKER (Advanced Chunking)
# ============================================================================

@dataclass
class ChunkMetadata:
    """Structured metadata for each chunk."""
    document_id: str
    chunk_id: str
    chunk_index: int
    primary_type: str
    has_tables: bool = False
    has_images: bool = False
    has_code: bool = False
    headings: List[str] = field(default_factory=list)
    page_numbers: List[int] = field(default_factory=list)
    element_count: int = 0
    token_count: int = 0
    part: Optional[int] = None
    total_parts: Optional[int] = None
    parent_id: Optional[str] = None
    procedures: List[str] = field(default_factory=list)
    labels: List[str] = field(default_factory=list)
    start_line: Optional[int] = None
    end_line: Optional[int] = None
    table_id: Optional[str] = None
    num_rows: Optional[int] = None
    num_cols: Optional[int] = None
    table_caption: Optional[str] = None
    source_file: Optional[str] = None
    section: Optional[str] = None


class ProductionChunker:
    """Production-ready semantic chunker."""

    def __init__(self, config):
        self.config = config
        if tiktoken is None:
            raise ImportError("tiktoken not installed. Install with: pip install tiktoken")

        self.tokenizer = tiktoken.get_encoding("cl100k_base")
        self.text_target_tokens = config.TARGET_CHUNK_SIZE
        self.text_overlap_tokens = config.CHUNK_OVERLAP
        self.text_max_tokens = config.MAX_CHUNK_SIZE
        self.text_min_tokens = config.MIN_CHUNK_SIZE
        self.table_max_tokens = config.TABLE_MAX_TOKENS
        self.table_rows_per_chunk = getattr(config, 'TABLE_ROWS_PER_CHUNK', 50)
        self.code_target_lines = config.CODE_TARGET_LINES
        self.code_overlap_lines = config.CODE_OVERLAP_LINES
        self.procedure_markers = getattr(config, 'ASSEMBLY_PROCEDURE_MARKERS',
                                         ['PROC', 'ENDP', 'CSECT', 'DSECT'])

    def chunk(self, elements: List[Dict], document_id: str, document_metadata: Dict) -> List[Dict]:
        logger.info(f"Production chunking: {len(elements)} elements from {document_id}")
        text_elements = []
        table_elements = []
        code_elements = []
        image_elements = []

        for element in elements:
            elem_type = element.get('type', 'text')
            if elem_type == 'table':
                table_elements.append(element)
            elif elem_type == 'code':
                code_elements.append(element)
            elif elem_type == 'image':
                image_elements.append(element)
            else:
                text_elements.append(element)

        all_chunks = []
        chunk_counter = [0]
        all_chunks.extend(self._chunk_text_elements(text_elements, document_id, document_metadata, chunk_counter))
        all_chunks.extend(self._chunk_table_elements(table_elements, document_id, document_metadata, chunk_counter))
        all_chunks.extend(self._chunk_code_elements(code_elements, document_id, document_metadata, chunk_counter))
        all_chunks.extend(self._chunk_image_elements(image_elements, text_elements, document_id, document_metadata, chunk_counter))

        validated_chunks = self._validate_and_enrich(all_chunks, document_id)
        logger.info(f"Created {len(validated_chunks)} chunks (avg tokens: {self._avg_tokens(validated_chunks):.0f})")
        return validated_chunks

    def _chunk_text_elements(self, elements: List[Dict], document_id: str,
                             doc_metadata: Dict, chunk_counter: List[int]) -> List[Dict]:
        if not elements:
            return []

        chunks = []
        current_elements = []
        current_tokens = 0
        heading_stack = []

        for element in elements:
            content = element.get('content', '')
            elem_meta = element.get('metadata', {})
            tokens = self._count_tokens(content)

            elem_type = elem_meta.get('element_type', '')
            if elem_type in ['Title', 'title']:
                heading_stack = [content.strip()[:100]]
            elif elem_type == 'heading':
                level = elem_meta.get('heading_level', 1)
                heading_stack = heading_stack[:level - 1] + [content.strip()[:100]]

            if current_tokens + tokens > self.text_max_tokens:
                if current_elements:
                    chunk = self._create_text_chunk(
                        current_elements, document_id, heading_stack.copy(),
                        doc_metadata, chunk_counter
                    )
                    chunks.append(chunk)

                overlap_elements, overlap_tokens = self._get_overlap_elements(
                    current_elements, self.text_overlap_tokens
                )
                current_elements = overlap_elements + [element]
                current_tokens = overlap_tokens + tokens
            elif current_tokens + tokens > self.text_target_tokens:
                if self._is_semantic_break(element):
                    if current_elements:
                        chunk = self._create_text_chunk(
                            current_elements, document_id, heading_stack.copy(),
                            doc_metadata, chunk_counter
                        )
                        chunks.append(chunk)
                    current_elements = [element]
                    current_tokens = tokens
                else:
                    current_elements.append(element)
                    current_tokens += tokens
            else:
                current_elements.append(element)
                current_tokens += tokens

        if current_elements:
            chunk = self._create_text_chunk(
                current_elements, document_id, heading_stack.copy(),
                doc_metadata, chunk_counter
            )
            chunks.append(chunk)

        return chunks

    def _create_text_chunk(self, elements: List[Dict], document_id: str,
                           headings: List[str], doc_metadata: Dict,
                           chunk_counter: List[int]) -> Dict:
        chunk_counter[0] += 1
        chunk_id = f"{document_id}-chunk-{chunk_counter[0]:04d}"
        content_parts = []
        page_numbers = set()
        has_tables = False
        has_images = False

        for elem in elements:
            content = elem.get('content', '')
            elem_meta = elem.get('metadata', {})
            elem_type = elem.get('type', 'text')
            if elem_type == 'table':
                content_parts.append(f"\n{content}\n")
                has_tables = True
            elif elem_type == 'image':
                has_images = True
                caption = elem_meta.get('caption', 'Image')
                ocr = elem_meta.get('ocr_text', '')
                desc = elem_meta.get('description', '')
                img_text = f"[{caption}]"
                if ocr:
                    img_text += f"\nText in image: {ocr}"
                if desc:
                    img_text += f"\nDescription: {desc}"
                content_parts.append(img_text)
            else:
                content_parts.append(content)

            page = elem_meta.get('page_no') or elem_meta.get('page_number')
            if page:
                page_numbers.add(page)

        combined_content = '\n'.join(content_parts)
        tokens = self._count_tokens(combined_content)
        section = headings[-1] if headings else None

        return {
            'chunk_id': chunk_id,
            'document_id': document_id,
            'type': 'text',
            'content': combined_content,
            'tokens': tokens,
            'metadata': {
                'headings': headings,
                'page_numbers': sorted(list(page_numbers)) if page_numbers else [],
                'element_count': len(elements),
                'has_tables': has_tables,
                'has_images': has_images,
                'has_code': False,
                'section': section,
                'chunk_index': chunk_counter[0],
                **doc_metadata
            }
        }

    def _chunk_table_elements(self, elements: List[Dict], document_id: str,
                               doc_metadata: Dict, chunk_counter: List[int]) -> List[Dict]:
        chunks = []
        for elem in elements:
            content = elem.get('content', '')
            elem_meta = elem.get('metadata', {})
            tokens = self._count_tokens(content)
            caption = elem_meta.get('caption', elem_meta.get('title', 'Table'))
            num_rows = elem_meta.get('num_rows', 0)
            num_cols = elem_meta.get('num_cols', 0)
            table_html = elem_meta.get('html', '')
            page_no = elem_meta.get('page_no')
            table_id = f"table-{uuid.uuid4().hex[:8]}"

            if tokens <= self.table_max_tokens:
                chunk_counter[0] += 1
                chunk_id = f"{document_id}-chunk-{chunk_counter[0]:04d}"
                table_desc = f"Table: {caption}\n"
                table_desc += f"This table has {num_rows} rows and {num_cols} columns.\n\n"
                full_content = table_desc + content
                chunks.append({
                    'chunk_id': chunk_id,
                    'document_id': document_id,
                    'type': 'table',
                    'content': full_content,
                    'tokens': self._count_tokens(full_content),
                    'metadata': {
                        'table_id': table_id,
                        'table_caption': caption,
                        'num_rows': num_rows,
                        'num_cols': num_cols,
                        'table_html': table_html,
                        'page_numbers': [page_no] if page_no else [],
                        'has_tables': True,
                        'has_images': False,
                        'has_code': False,
                        'is_complete': True,
                        'chunk_index': chunk_counter[0],
                        **doc_metadata
                    }
                })
            else:
                table_chunks = self._split_large_table(
                    content, elem_meta, document_id, table_id,
                    doc_metadata, chunk_counter
                )
                chunks.extend(table_chunks)
        return chunks

    def _split_large_table(self, content: str, elem_meta: Dict, document_id: str,
                            table_id: str, doc_metadata: Dict,
                            chunk_counter: List[int]) -> List[Dict]:
        lines = content.strip().split('\n')
        if len(lines) < 2:
            return []

        header_lines = []
        data_start = 0
        for i, line in enumerate(lines):
            if re.match(r'^[\s|:-]+$', line):
                header_lines = lines[:i + 1]
                data_start = i + 1
                break
            if i == 0:
                header_lines = [line]

        if not header_lines:
            header_lines = [lines[0]]
            data_start = 1

        header = '\n'.join(header_lines)
        data_lines = lines[data_start:]
        if not data_lines:
            return []

        rows_per_chunk = self.table_rows_per_chunk
        total_parts = (len(data_lines) + rows_per_chunk - 1) // rows_per_chunk

        caption = elem_meta.get('caption', 'Table')
        num_cols = elem_meta.get('num_cols', 0)
        page_no = elem_meta.get('page_no')
        table_html = elem_meta.get('html', '')

        chunks = []
        for part in range(total_parts):
            start_row = part * rows_per_chunk
            end_row = min(start_row + rows_per_chunk, len(data_lines))
            part_rows = data_lines[start_row:end_row]
            part_content = header + '\n' + '\n'.join(part_rows)

            chunk_counter[0] += 1
            chunk_id = f"{document_id}-chunk-{chunk_counter[0]:04d}"

            desc = f"Table: {caption} (Part {part + 1} of {total_parts})\n"
            desc += f"Rows {start_row + 1} to {end_row} of {len(data_lines)} total rows.\n\n"
            full_content = desc + part_content

            chunks.append({
                'chunk_id': chunk_id,
                'document_id': document_id,
                'type': 'table',
                'content': full_content,
                'tokens': self._count_tokens(full_content),
                'metadata': {
                    'table_id': table_id,
                    'table_caption': caption,
                    'num_rows': len(part_rows),
                    'num_cols': num_cols,
                    'part': part + 1,
                    'total_parts': total_parts,
                    'rows_range': f"{start_row + 1}-{end_row}",
                    'page_numbers': [page_no] if page_no else [],
                    'has_tables': True,
                    'has_images': False,
                    'has_code': False,
                    'is_complete': False,
                    'table_html': table_html if part == 0 else '',
                    'chunk_index': chunk_counter[0],
                    **doc_metadata
                }
            })
        return chunks

    def _chunk_code_elements(self, elements: List[Dict], document_id: str,
                              doc_metadata: Dict, chunk_counter: List[int]) -> List[Dict]:
        chunks = []
        for elem in elements:
            content = elem.get('content', '')
            elem_meta = elem.get('metadata', {})
            lines = content.split('\n')
            num_lines = len(lines)
            procedures = self._find_procedures(lines)

            if num_lines <= self.code_target_lines:
                chunk_counter[0] += 1
                chunk_id = f"{document_id}-chunk-{chunk_counter[0]:04d}"
                labels = self._extract_labels(lines)
                proc_names = [p['name'] for p in procedures]
                chunks.append({
                    'chunk_id': chunk_id,
                    'document_id': document_id,
                    'type': 'code',
                    'content': content,
                    'tokens': self._count_tokens(content),
                    'metadata': {
                        'has_code': True,
                        'has_tables': False,
                        'has_images': False,
                        'line_count': num_lines,
                        'start_line': 1,
                        'end_line': num_lines,
                        'procedures': proc_names,
                        'labels': labels[:50],
                        'is_complete': True,
                        'language': elem_meta.get('language', 'assembly'),
                        'chunk_index': chunk_counter[0],
                        **doc_metadata
                    }
                })
            else:
                code_chunks = self._split_large_code(
                    lines, procedures, elem_meta, document_id,
                    doc_metadata, chunk_counter
                )
                chunks.extend(code_chunks)
        return chunks

    def _find_procedures(self, lines: List[str]) -> List[Dict]:
        procedures = []
        current_proc = None
        for i, line in enumerate(lines):
            line_upper = line.upper()
            for marker in ['PROC', 'CSECT']:
                if marker in line_upper:
                    if current_proc:
                        current_proc['end_line'] = i - 1
                        procedures.append(current_proc)
                    parts = line.split()
                    name = parts[0] if parts else f"PROC_{i}"
                    current_proc = {'name': name, 'start_line': i, 'type': marker}
                    break
            if 'ENDP' in line_upper or 'END' in line_upper:
                if current_proc:
                    current_proc['end_line'] = i
                    procedures.append(current_proc)
                    current_proc = None
        if current_proc:
            current_proc['end_line'] = len(lines) - 1
            procedures.append(current_proc)
        return procedures

    def _extract_labels(self, lines: List[str]) -> List[str]:
        labels = []
        label_pattern = re.compile(r'^([A-Z_][A-Z0-9_]*)\s*:', re.IGNORECASE)
        for line in lines:
            match = label_pattern.match(line)
            if match:
                labels.append(match.group(1))
        return labels

    def _split_large_code(self, lines: List[str], procedures: List[Dict],
                           elem_meta: Dict, document_id: str,
                           doc_metadata: Dict, chunk_counter: List[int]) -> List[Dict]:
        chunks = []
        num_lines = len(lines)
        start_line = 0
        while start_line < num_lines:
            end_line = min(start_line + self.code_target_lines, num_lines)
            for proc in procedures:
                proc_end = proc['end_line']
                if start_line < proc_end <= end_line + 20:
                    end_line = proc_end + 1
                    break
            chunk_lines = lines[start_line:end_line]
            chunk_content = '\n'.join(chunk_lines)
            chunk_counter[0] += 1
            chunk_id = f"{document_id}-chunk-{chunk_counter[0]:04d}"
            chunk_procs = [
                p['name'] for p in procedures
                if p['start_line'] >= start_line and p['start_line'] < end_line
            ]
            labels = self._extract_labels(chunk_lines)
            chunks.append({
                'chunk_id': chunk_id,
                'document_id': document_id,
                'type': 'code',
                'content': chunk_content,
                'tokens': self._count_tokens(chunk_content),
                'metadata': {
                    'has_code': True,
                    'has_tables': False,
                    'has_images': False,
                    'line_count': len(chunk_lines),
                    'start_line': start_line + 1,
                    'end_line': end_line,
                    'procedures': chunk_procs,
                    'labels': labels[:30],
                    'is_complete': False,
                    'language': elem_meta.get('language', 'assembly'),
                    'chunk_index': chunk_counter[0],
                    **doc_metadata
                }
            })
            start_line = end_line - self.code_overlap_lines
        return chunks

    def _chunk_image_elements(self, image_elements: List[Dict],
                               text_elements: List[Dict], document_id: str,
                               doc_metadata: Dict, chunk_counter: List[int]) -> List[Dict]:
        chunks = []
        for elem in image_elements:
            elem_meta = elem.get('metadata', {})
            caption = elem_meta.get('caption', 'Image')
            ocr_text = elem_meta.get('ocr_text', '')
            description = elem_meta.get('description', '')
            preceding = elem_meta.get('preceding_text', '')
            following = elem_meta.get('following_text', '')
            page_no = elem_meta.get('page_no')
            image_path = elem_meta.get('local_path', '')

            content_parts = []
            content_parts.append(f"Image: {caption}")
            if preceding:
                content_parts.append(f"\nContext: {preceding[:300]}")
            if ocr_text:
                content_parts.append(f"\nText in image: {ocr_text[:500]}")
            if description:
                content_parts.append(f"\nDescription: {description[:600]}")
            if following:
                content_parts.append(f"\nExplanation: {following[:300]}")

            full_content = '\n'.join(content_parts)
            chunk_counter[0] += 1
            chunk_id = f"{document_id}-chunk-{chunk_counter[0]:04d}"

            chunks.append({
                'chunk_id': chunk_id,
                'document_id': document_id,
                'type': 'image',
                'content': full_content,
                'tokens': self._count_tokens(full_content),
                'metadata': {
                    'has_images': True,
                    'has_tables': False,
                    'has_code': False,
                    'image_path': image_path,
                    'image_caption': caption,
                    'has_ocr': bool(ocr_text),
                    'has_description': bool(description),
                    'page_numbers': [page_no] if page_no else [],
                    'chunk_index': chunk_counter[0],
                    **doc_metadata
                }
            })

        return chunks

    def _validate_and_enrich(self, chunks: List[Dict], document_id: str) -> List[Dict]:
        validated = []
        for chunk in chunks:
            tokens = chunk.get('tokens', 0)
            if tokens < self.text_min_tokens:
                chunk['metadata']['quality_flag'] = 'undersized'
            elif tokens > self.text_max_tokens:
                chunk['metadata']['quality_flag'] = 'oversized'
            else:
                chunk['metadata']['quality_flag'] = 'valid'

            chunk['metadata'].setdefault('headings', [])
            chunk['metadata'].setdefault('page_numbers', [])
            chunk['metadata'].setdefault('has_tables', False)
            chunk['metadata'].setdefault('has_images', False)
            chunk['metadata'].setdefault('has_code', False)
            chunk['metadata']['created_at'] = datetime.utcnow().isoformat() + 'Z'
            validated.append(chunk)
        return validated

    def _count_tokens(self, text: str) -> int:
        if not text:
            return 0
        try:
            return len(self.tokenizer.encode(text))
        except Exception:
            return int(len(text.split()) * 1.3)

    def _is_semantic_break(self, element: Dict) -> bool:
        elem_type = element.get('metadata', {}).get('element_type', '')
        return elem_type in ['Title', 'title', 'heading', 'PageBreak', 'section_break']

    def _get_overlap_elements(self, elements: List[Dict],
                               target_tokens: int) -> Tuple[List[Dict], int]:
        if not elements:
            return [], 0
        overlap_elements = []
        tokens_accumulated = 0
        for elem in reversed(elements):
            elem_tokens = self._count_tokens(elem.get('content', ''))
            if tokens_accumulated + elem_tokens <= target_tokens * 1.2:
                overlap_elements.insert(0, elem)
                tokens_accumulated += elem_tokens
            else:
                break
        return overlap_elements, tokens_accumulated

    def _avg_tokens(self, chunks: List[Dict]) -> float:
        if not chunks:
            return 0
        return sum(c.get('tokens', 0) for c in chunks) / len(chunks)


# ============================================================================
# AWS CLIENTS
# ============================================================================

_clients: Dict[str, Any] = {}
_circuit_breakers: Dict[str, Any] = {}


class CircuitBreaker:
    """Simple circuit breaker for downstream dependencies."""

    def __init__(self, threshold: int = 5, recovery_seconds: int = 30):
        self.threshold = threshold
        self.recovery_seconds = recovery_seconds
        self.failures = 0
        self.open_until = 0.0

    def allow(self) -> bool:
        if self.open_until == 0.0:
            return True
        return time.time() >= self.open_until

    def record_success(self) -> None:
        self.failures = 0
        self.open_until = 0.0

    def record_failure(self) -> None:
        self.failures += 1
        if self.failures >= self.threshold:
            self.open_until = time.time() + self.recovery_seconds


def _get_circuit_breaker(name: str) -> CircuitBreaker:
    if name not in _circuit_breakers:
        _circuit_breakers[name] = CircuitBreaker()
    return _circuit_breakers[name]


def get_client(service: str):
    """Get or create AWS client."""
    if service not in _clients:
        _clients[service] = boto3.client(service, region_name=AWS_REGION)
    return _clients[service]


def get_resource(service: str):
    """Get or create AWS resource."""
    key = f"{service}_resource"
    if key not in _clients:
        _clients[key] = boto3.resource(service, region_name=AWS_REGION)
    return _clients[key]


def call_with_circuit_breaker(name: str, func, *args, **kwargs):
    breaker = _get_circuit_breaker(name)
    if not breaker.allow():
        raise RuntimeError(f"Circuit breaker open for {name}")
    try:
        result = func(*args, **kwargs)
        breaker.record_success()
        return result
    except Exception:
        breaker.record_failure()
        raise


class DocumentExtractEvent(BaseModel):
    model_config = ConfigDict(extra='allow')
    document_id: str
    s3_key: str
    file_type: str
    s3_bucket: Optional[str] = None
    file_name: Optional[str] = None


def get_pipeline_status(document_id: str) -> Optional[str]:
    """Read current pipeline status for idempotency checks."""
    try:
        table = get_resource('dynamodb').Table(PIPELINE_TABLE)
        response = table.get_item(
            Key={'document_id': document_id},
            ProjectionExpression='#s',
            ExpressionAttributeNames={'#s': 'status'}
        )
        return response.get('Item', {}).get('status')
    except Exception as e:
        logger.warning(f"Failed to read pipeline status: {e}")
        return None


# ============================================================================
# MECH ENCODING NORMALIZATION
# ============================================================================

def normalize_mainframe_encoding(text: str) -> str:
    """
    Normalize mainframe document encodings for better searchability.
    
    Addresses failure mode: "encoding mangled" causing missed searches.
    
    Fixes:
    - Ligatures (fi, fl, etc.)
    - Non-breaking spaces
    - Smart quotes
    - Hex escape variations
    - OCR errors in numeric contexts
    """
    # Ligature fixes
    ligature_map = {
        '\ufb01': 'fi', '\ufb02': 'fl', '\ufb00': 'ff',
        '\ufb03': 'ffi', '\ufb04': 'ffl',
    }
    
    # Quote and dash fixes
    quote_map = {
        '\u2018': "'", '\u2019': "'", '\u201c': '"', '\u201d': '"',
        '\u2013': '-', '\u2014': '-', '\u00a0': ' ',
    }
    
    # EBCDIC remnant fixes
    ebcdic_fixes = {
        '\x15': '\n', '\x25': '\n', '\x0d': '\n', '\x05': '\t',
    }
    
    for lig, replacement in ligature_map.items():
        text = text.replace(lig, replacement)
    
    for char, replacement in quote_map.items():
        text = text.replace(char, replacement)
    
    for char, replacement in ebcdic_fixes.items():
        text = text.replace(char, replacement)
    
    # Normalize hex formats: 0x0A -> X'0A'
    text = re.sub(r'0x([0-9A-Fa-f]{2})', lambda m: f"X'{m.group(1).upper()}'", text)
    
    # Fix JCL comment markers that get mangled
    text = re.sub(r'/\*\s*\*/', '/* */', text)
    
    return text


def classify_mech_domain(content: str) -> str:
    """Classify content into MECH domain for routing."""
    scores = {}
    
    for domain, patterns in DOMAIN_PATTERNS.items():
        score = 0
        for pattern in patterns:
            matches = re.findall(pattern, content, re.IGNORECASE)
            score += len(matches)
        scores[domain] = score
    
    if max(scores.values()) > 0:
        best_domain = max(scores, key=scores.get)
        return best_domain.value
    return MechDomain.GENERAL.value


def extract_mech_metadata_from_chunk(content: str) -> Dict[str, Any]:
    """
    Extract MECH-specific metadata from chunk content.
    
    Extracts:
    - OPTRTYP/BDD field mappings
    - SHF segment identification (SHF03, SHF04, etc.)
    - Checkpoint file references (GCHK*)
    - Program call chains (STDPC macros)
    """
    metadata = {
        'domain': classify_mech_domain(content),
        'bdd_tables': [],
        'field_names': [],
        'shf_segment': None,
        'xml_tags': [],
        'program_chain': [],
        'fsf_fields': [],
        'contains_txnlog': False,
        'contains_rl03': False,
        'contains_xxa': False,
        'contains_hex': False,
        'contains_jcl': False,
        'contains_macro': False,
        'is_mapping': False,
        'is_procedure': False,
    }
    
    # Extract BDD tables
    for pattern in MECH_PATTERNS['bdd_tables']:
        matches = re.findall(pattern, content)
        metadata['bdd_tables'].extend(matches)
    metadata['bdd_tables'] = list(set(metadata['bdd_tables']))
    
    # Extract field names
    for pattern in MECH_PATTERNS['field_names']:
        matches = re.findall(pattern, content)
        metadata['field_names'].extend(matches)
    metadata['field_names'] = list(set(metadata['field_names']))
    
    # Extract SHF segment
    shf_match = re.search(MECH_PATTERNS['shf_segments'], content)
    if shf_match:
        metadata['shf_segment'] = f"SHF0{shf_match.group(1)}"
    
    # Extract XML tags
    xml_matches = re.findall(MECH_PATTERNS['xml_tags'], content)
    metadata['xml_tags'] = list(set(xml_matches))
    
    # Extract programs
    for pattern in MECH_PATTERNS['programs']:
        matches = re.findall(pattern, content)
        metadata['program_chain'].extend(matches)
    metadata['program_chain'] = list(set(metadata['program_chain']))
    
    # Extract FSF fields
    fsf_matches = re.findall(MECH_PATTERNS['fsf_fields'], content)
    metadata['fsf_fields'] = list(set(f'FSFL{m}' for m in fsf_matches))
    
    # Boolean flags
    metadata['contains_txnlog'] = bool(re.search(r'\b(?:TXNLOG|MECH_TXNLOG)\b', content, re.I))
    metadata['contains_rl03'] = bool(re.search(r'\bRL[-\s]?03\b', content, re.I))
    metadata['contains_xxa'] = bool(re.search(r'\bXXA\b', content))
    metadata['contains_hex'] = any(re.search(p, content) for p in MECH_PATTERNS['hex_data'])
    metadata['contains_jcl'] = any(re.search(p, content, re.M) for p in MECH_PATTERNS['jcl'])
    metadata['contains_macro'] = any(re.search(p, content) for p in MECH_PATTERNS['macros'])
    
    # Quality indicators
    metadata['is_mapping'] = bool(re.search(r'\b(?:map(?:ping|ped)?|correspond|contains)\b', content, re.I))
    metadata['is_procedure'] = bool(re.search(r'\b(?:step\s*\d|then|finally|first)\b', content, re.I))
    
    # Count technical identifier density (instead of checking against hardcoded list)
    all_identifiers = (
        metadata.get('bdd_tables', []) + 
        metadata.get('field_names', []) + 
        metadata.get('programs', []) + 
        metadata.get('shf_segments', [])
    )
    metadata['identifier_count'] = len(all_identifiers)
    metadata['identifier_density'] = len(all_identifiers) / max(len(content.split()), 1)
    
    return metadata


# ============================================================================
# GENERIC METADATA EXTRACTION (LLM-Powered)
# ============================================================================

# Feature flag for LLM-based metadata extraction
ENABLE_LLM_METADATA = os.environ.get('ENABLE_LLM_METADATA', 'true').lower() == 'true'
HAIKU_MODEL = os.environ.get(
    'HAIKU_INFERENCE_PROFILE_ARN',
    'arn:aws:bedrock:us-east-1:966245098051:application-inference-profile/wt15xad4wqki'
)

# Metadata extraction settings
METADATA_SAMPLE_SIZE = int(os.environ.get('METADATA_SAMPLE_SIZE', 2000))
METADATA_BATCH_PAUSE = float(os.environ.get('METADATA_BATCH_PAUSE', 0.2))


def extract_document_metadata(text_sample: str, file_name: str = None) -> Dict[str, Any]:
    """
    GENERIC LLM-based document metadata extraction.
    
    NO HARDCODED DOCUMENT TYPES - fully dynamic extraction based on content.
    This function analyzes the document and extracts semantic metadata that
    can be used for retrieval boosting without predefined categories.
    
    Extracts:
        - doc_summary: 1-2 sentence summary of what the document is about
        - primary_topics: List of main topics covered (dynamically extracted)
        - key_terms: Important technical terms, acronyms, system names
        - doc_type_inferred: Inferred document type (procedure, reference, specification, etc.)
        - domain_indicators: Domain-specific markers for retrieval matching
    
    Returns:
        Dict with dynamically extracted metadata
    """
    if not ENABLE_LLM_METADATA:
        return _get_default_metadata()
    
    try:
        bedrock = get_client('bedrock-runtime')
        
        # Use larger sample for document-level analysis
        sample = text_sample[:METADATA_SAMPLE_SIZE] if len(text_sample) > METADATA_SAMPLE_SIZE else text_sample
        
        prompt = f"""<task>Extract metadata from this technical document for search/retrieval optimization.</task>

<document_info>
Filename: {file_name or 'Unknown'}
</document_info>

<document_excerpt>
{sample}
</document_excerpt>

<extraction_instructions>
Analyze the actual content and extract the following fields:

1. **doc_summary**: 1-2 concise sentences describing what this document covers.
   - Be specific about the system, process, or topic
   - Example: "Describes the SHF Switch process for MECH mainframe, including trigger conditions and recovery procedures."

2. **primary_topics**: 3-5 main topics/subjects (extract from content, not generic terms)
   - Examples: "SHF_processing", "batch_job_configuration", "error_recovery", "TSS_authorization"

3. **key_terms**: 5-15 important technical terms, acronyms, system names FOUND in the text
   - Include: program names (MYPTLA), file names (SHF03), acronyms (CICS, TSS), codes (S0C7)
   - Extract EXACTLY as they appear in the document

4. **doc_type_inferred**: Single document type based on content style:
   - procedure: step-by-step instructions
   - reference: lookup tables, parameter lists
   - specification: technical details, field layouts
   - guide: conceptual explanations
   - troubleshooting: error resolution, diagnostics
   - configuration: setup, parameters
   - code: source code, JCL, scripts
   - overview: high-level descriptions

5. **domain_indicators**: 3-8 terms that would help match user queries to this document
   - Think: what would someone search to find this document?
   - Include synonyms and related terms
</extraction_instructions>

<rules>
- Extract ACTUAL terms from the document
- Do NOT use generic placeholders
- If a field cannot be determined, use empty string or empty array
</rules>

<output>Return ONLY valid JSON:
{{"doc_summary": "...", "primary_topics": [...], "key_terms": [...], "doc_type_inferred": "...", "domain_indicators": [...]}}</output>"""

        payload = {
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 500,
            "temperature": 0,
            "messages": [{"role": "user", "content": prompt}]
        }
        
        response = call_with_circuit_breaker(
            'bedrock',
            bedrock.invoke_model,
            modelId=HAIKU_MODEL,
            body=json.dumps(payload)
        )
        
        result = json.loads(response['body'].read())
        metadata_text = result['content'][0]['text'].strip()
        
        try:
            metadata = json.loads(metadata_text)
            
            # Ensure all expected fields exist
            metadata.setdefault('doc_summary', '')
            metadata.setdefault('primary_topics', [])
            metadata.setdefault('key_terms', [])
            metadata.setdefault('doc_type_inferred', 'unknown')
            metadata.setdefault('domain_indicators', [])
            
            # Add system defaults
            metadata['sensitivity'] = 'internal'
            metadata['department'] = 'mainframe'
            metadata['upload_date'] = datetime.now().isoformat()
            metadata['source'] = 's3'
            metadata['extraction_method'] = 'llm_v2'
            
            logger.info(f"Extracted metadata: topics={len(metadata['primary_topics'])}, "
                       f"terms={len(metadata['key_terms'])}, type={metadata['doc_type_inferred']}")
            
            return metadata
            
        except json.JSONDecodeError:
            logger.warning(f"Failed to parse document metadata JSON: {metadata_text[:100]}")
            return _get_default_metadata()
        
    except Exception as e:
        logger.warning(f"Document metadata extraction failed (non-blocking): {e}")
        return _get_default_metadata()


def _get_default_metadata() -> Dict[str, Any]:
    """Return default metadata when LLM extraction fails."""
    return {
        'doc_summary': '',
        'primary_topics': [],
        'key_terms': [],
        'doc_type_inferred': 'unknown',
        'domain_indicators': [],
        'sensitivity': 'internal',
        'department': 'mainframe',
        'upload_date': datetime.now().isoformat(),
        'source': 's3',
        'extraction_method': 'default'
    }


def extract_rich_metadata(chunk_text: str, file_name: str = None) -> Dict[str, Any]:
    """
    Extract rich semantic metadata from a single chunk using Claude Haiku.
    
    FULLY DYNAMIC - No hardcoded categories or document types.
    Extracts whatever is relevant from the actual content.
    
    Extracts:
        - entities: Technical terms, acronyms, system names found in this chunk
        - topics: High-level topics covered in this chunk
        - content_type: Type of content (procedure, troubleshooting, definition, code, overview, etc.)
        - section_title: Inferred section title if identifiable
        - key_concepts: Main concepts discussed for retrieval matching
    
    Returns:
        Dict with metadata fields
    """
    if not ENABLE_LLM_METADATA:
        return {
            'entities': [],
            'topics': [],
            'content_type': 'unknown',
            'section_title': None,
            'sensitivity': 'internal',
            'department': 'mainframe'
        }
    
    try:
        bedrock = get_client('bedrock-runtime')
        
        # Truncate to first 1500 chars for cost efficiency
        text_sample = chunk_text[:1500] if len(chunk_text) > 1500 else chunk_text
        
        prompt = f"""<task>Extract semantic metadata from this documentation chunk for search optimization.</task>

<source_file>{file_name or 'Unknown'}</source_file>

<chunk_text>
{text_sample}
</chunk_text>

<extraction_fields>
1. **entities**: Technical identifiers found in this chunk
   - System names: MECH, CADS, SAGE, IMS, CICS
   - Program names: MYPTLA, MYNBTA, MSHFPA
   - File/dataset names: SHF03, SHF04, CIF, CAF$
   - Error codes: S0C7, ABEND, RC=12
   - Security terms: TSS, RACF, ACF2
   - Extract EXACTLY as written (5-20 items)

2. **topics**: 1-3 high-level topics this chunk covers
   - Choose from: batch_processing, error_handling, security, configuration, troubleshooting, 
     data_processing, file_operations, transaction_processing, reporting, maintenance
   - Or create domain-specific topic if needed

3. **content_type**: Single type that best describes this chunk
   | Type | Indicators |
   |------|------------|
   | procedure | numbered steps, "do this", action verbs |
   | troubleshooting | error codes, "if...then", problem/solution |
   | definition | "is defined as", explanatory text |
   | code | JCL, COBOL, assembly, scripts |
   | overview | conceptual, introduces topics |
   | reference | tables, parameter lists, lookups |
   | configuration | settings, parameters, setup |

4. **section_title**: Inferred title for this chunk (or null if unclear)
   - Look for: headers, bold text, topic sentences
   - Example: "SHF Switch Activation Procedure"
</extraction_fields>

<output>Return ONLY valid JSON:
{{"entities": [...], "topics": [...], "content_type": "...", "section_title": "...or null"}}</output>"""

        payload = {
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 600,  # Increased from 300 to handle entity-rich content
            "temperature": 0,
            "messages": [{"role": "user", "content": prompt}]
        }
        
        response = call_with_circuit_breaker(
            'bedrock',
            bedrock.invoke_model,
            modelId=HAIKU_MODEL,
            body=json.dumps(payload)
        )
        
        result = json.loads(response['body'].read())
        metadata_text = result['content'][0]['text'].strip()
        
        try:
            # Try to extract JSON from response (handles LLM adding extra text)
            # Look for {...} pattern
            import re
            json_match = re.search(r'\{.*\}', metadata_text, re.DOTALL)
            if json_match:
                metadata_text = json_match.group(0)
            
            metadata = json.loads(metadata_text)
            
            # Normalize content_type
            valid_types = ['procedure', 'troubleshooting', 'definition', 'code', 'overview', 'reference', 'configuration']
            if metadata.get('content_type') not in valid_types:
                metadata['content_type'] = 'overview'
            
            # Add default fields
            metadata['sensitivity'] = 'internal'
            metadata['department'] = 'mainframe'
            metadata['upload_date'] = datetime.now().isoformat()
            metadata['source'] = 's3'
            
            return metadata
            
        except json.JSONDecodeError as je:
            logger.warning(f"Failed to parse metadata JSON (truncated at {len(metadata_text)} chars): {str(je)}")
            return {
                'entities': [],
                'topics': [],
                'content_type': 'unknown',
                'section_title': None,
                'sensitivity': 'internal',
                'department': 'mainframe'
            }
        
    except Exception as e:
        logger.warning(f"Metadata extraction failed (non-blocking): {e}")
        return {
            'entities': [],
            'topics': [],
            'content_type': 'unknown',
            'section_title': None,
            'sensitivity': 'internal',
            'department': 'mainframe'
        }


def enrich_chunks_with_metadata(chunks: List[Chunk], file_name: str = None,
                                 batch_size: int = 5) -> List[Chunk]:
    """
    PHASE 1.5: Enrich all chunks with LLM-extracted metadata.
    
    Processes in batches to avoid throttling.
    """
    if not ENABLE_LLM_METADATA or not chunks:
        return chunks
    
    logger.info(f"Enriching {len(chunks)} chunks with LLM metadata...")
    
    enriched_count = 0
    
    for i, chunk in enumerate(chunks):
        try:
            # Extract metadata
            rich_metadata = extract_rich_metadata(chunk.text, file_name)
            
            # Update chunk metadata
            if chunk.metadata is None:
                chunk.metadata = {}
            
            chunk.metadata.update(rich_metadata)
            
            # Also update section_title if we got a better one
            if rich_metadata.get('section_title') and not chunk.section_title:
                chunk.section_title = rich_metadata['section_title']
            
            enriched_count += 1
            
            # Rate limiting - small delay every batch_size chunks
            if (i + 1) % batch_size == 0:
                import time
                time.sleep(0.2)  # 200ms pause to avoid throttling
                
        except Exception as e:
            logger.warning(f"Failed to enrich chunk {i}: {e}")
            continue
    
    logger.info(f"Enriched {enriched_count}/{len(chunks)} chunks with metadata")
    return chunks


# ============================================================================
# EXTRACTION ERROR HANDLING
# ============================================================================

class ExtractionError(Exception):
    """Custom exception for extraction failures."""
    def __init__(self, message: str, method: str = None, recoverable: bool = True):
        super().__init__(message)
        self.method = method
        self.recoverable = recoverable


# ============================================================================
# PDF EXTRACTION (Multiple fallback methods)
# ============================================================================

def extract_pdf(file_path: str) -> ExtractionResult:
    """
    Extract text from PDF with multiple fallback methods.
    
    Priority:
        1. pymupdf4llm (best for LLM, includes tables as markdown)
        2. PyMuPDF (fitz) direct extraction
        3. pypdf (basic fallback)
    """
    errors = []
    
    # Method 1: pymupdf4llm (optimal for LLM processing)
    try:
        import pymupdf4llm
        import fitz
        
        doc = fitz.open(file_path)
        page_count = len(doc)
        
        # Count images and tables
        image_count = 0
        table_count = 0
        for page in doc:
            image_count += len(page.get_images(full=True))
            tables = page.find_tables()
            if tables and tables.tables:
                table_count += len(tables.tables)
        doc.close()
        
        # Extract with pymupdf4llm
        md_text = pymupdf4llm.to_markdown(
            file_path,
            page_chunks=False,
            write_images=False,
            show_progress=False
        )
        
        # Add page markers
        text = _add_page_markers(md_text, file_path)
        
        return ExtractionResult(
            success=True,
            text=text,
            method=ExtractionMethod.PYMUPDF4LLM.value,
            pages=page_count,
            tables=table_count,
            images=image_count,
            word_count=len(text.split()),
            char_count=len(text)
        )
        
    except ImportError:
        errors.append("pymupdf4llm not available")
    except Exception as e:
        errors.append(f"pymupdf4llm failed: {str(e)}")
    
    # Method 2: PyMuPDF direct
    try:
        import fitz
        
        doc = fitz.open(file_path)
        page_count = len(doc)
        text_parts = []
        image_count = 0
        
        for page_num, page in enumerate(doc):
            page_text = page.get_text("text")
            if page_text.strip():
                text_parts.append(f"\n[Page {page_num + 1}]\n{page_text}")
            image_count += len(page.get_images(full=True))
        
        doc.close()
        text = '\n'.join(text_parts)
        
        return ExtractionResult(
            success=True,
            text=text,
            method=ExtractionMethod.PYMUPDF.value,
            pages=page_count,
            images=image_count,
            word_count=len(text.split()),
            char_count=len(text)
        )
        
    except ImportError:
        errors.append("PyMuPDF not available")
    except Exception as e:
        errors.append(f"PyMuPDF failed: {str(e)}")
    
    # Method 3: pypdf (basic fallback)
    try:
        from pypdf import PdfReader
        
        reader = PdfReader(file_path)
        page_count = len(reader.pages)
        text_parts = []
        
        for page_num, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text_parts.append(f"\n[Page {page_num + 1}]\n{page_text}")
        
        text = '\n'.join(text_parts)
        
        return ExtractionResult(
            success=True,
            text=text,
            method=ExtractionMethod.PYPDF.value,
            pages=page_count,
            word_count=len(text.split()),
            char_count=len(text)
        )
        
    except ImportError:
        errors.append("pypdf not available")
    except Exception as e:
        errors.append(f"pypdf failed: {str(e)}")
    
    # All methods failed
    raise ExtractionError(
        f"PDF extraction failed. Errors: {'; '.join(errors)}",
        method="pdf",
        recoverable=False
    )


def _add_page_markers(md_text: str, file_path: str) -> str:
    """Add page markers to pymupdf4llm output."""
    try:
        import fitz
        doc = fitz.open(file_path)
        page_count = len(doc)
        doc.close()
        
        # Split by page break markers
        parts = re.split(r'\n-{3,}\n|\n\n---\n\n', md_text)
        
        if len(parts) == 1:
            return f"[Page 1]\n{md_text}"
        
        result_parts = []
        for i, part in enumerate(parts):
            if part.strip():
                result_parts.append(f"\n[Page {i + 1}]\n{part.strip()}")
        
        return '\n'.join(result_parts)
    except Exception:
        return md_text


# ============================================================================
# WORD DOCUMENT EXTRACTION
# ============================================================================

def extract_docx(file_path: str) -> ExtractionResult:
    """
    Extract text from DOCX with full formatting preservation.
    
    Features:
        - Paragraphs with style detection
        - Tables as markdown
        - Headers/footers
        - Document properties
    """
    try:
        from docx import Document
        from docx.table import Table
        from docx.oxml.text.paragraph import CT_P
        from docx.oxml.table import CT_Tbl
        
        doc = Document(file_path)
        text_parts = []
        table_count = 0
        
        # Extract document properties
        props = doc.core_properties
        if props.title:
            text_parts.append(f"# {props.title}\n")
        if props.subject:
            text_parts.append(f"**Subject:** {props.subject}\n")
        if props.author:
            text_parts.append(f"**Author:** {props.author}\n")
        
        text_parts.append("\n---\n")
        
        # Process body elements in order (paragraphs and tables)
        for element in doc.element.body:
            if isinstance(element, CT_P):
                para = None
                for p in doc.paragraphs:
                    if p._element is element:
                        para = p
                        break
                if para and para.text.strip():
                    # Detect heading style
                    style_name = para.style.name if para.style else ""
                    if "Heading 1" in style_name:
                        text_parts.append(f"\n# {para.text}\n")
                    elif "Heading 2" in style_name:
                        text_parts.append(f"\n## {para.text}\n")
                    elif "Heading 3" in style_name:
                        text_parts.append(f"\n### {para.text}\n")
                    elif "Title" in style_name:
                        text_parts.append(f"\n# {para.text}\n")
                    else:
                        text_parts.append(f"{para.text}\n")
                        
            elif isinstance(element, CT_Tbl):
                # Find matching table
                for tbl in doc.tables:
                    if tbl._element is element:
                        table_md = _table_to_markdown(tbl)
                        text_parts.append(f"\n{table_md}\n")
                        table_count += 1
                        break
        
        # Extract headers and footers
        for section in doc.sections:
            header = section.header
            if header and header.paragraphs:
                header_text = ' '.join(p.text for p in header.paragraphs if p.text.strip())
                if header_text:
                    text_parts.insert(0, f"[Header: {header_text}]\n")
        
        text = '\n'.join(text_parts)
        
        return ExtractionResult(
            success=True,
            text=text,
            method=ExtractionMethod.PYTHON_DOCX.value,
            tables=table_count,
            word_count=len(text.split()),
            char_count=len(text),
            metadata={
                'title': props.title,
                'author': props.author,
                'subject': props.subject,
                'created': str(props.created) if props.created else None,
                'modified': str(props.modified) if props.modified else None
            }
        )
        
    except ImportError:
        raise ExtractionError("python-docx library not available", "docx")
    except Exception as e:
        raise ExtractionError(f"DOCX extraction failed: {str(e)}", "docx")


def _table_to_markdown(table) -> str:
    """Convert DOCX table to markdown format."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
        rows.append('| ' + ' | '.join(cells) + ' |')
    
    if len(rows) >= 1:
        # Add header separator
        header_sep = '| ' + ' | '.join(['---'] * len(table.rows[0].cells)) + ' |'
        rows.insert(1, header_sep)
    
    return '\n'.join(rows)


def extract_doc_legacy(file_path: str) -> ExtractionResult:
    """
    Extract text from legacy .doc files using OLE parsing.
    
    .doc files are OLE (Object Linking and Embedding) compound documents.
    Uses olefile for proper parsing with multiple fallback strategies.
    """
    errors = []
    
    # Strategy 1: Use olefile for proper OLE parsing
    try:
        import olefile
        
        ole = olefile.OleFileIO(file_path)
        text_parts = []
        
        # Get document metadata
        try:
            meta = ole.get_metadata()
            if meta:
                if meta.title:
                    text_parts.append(f"# {meta.title}")
                if meta.subject:
                    text_parts.append(f"**Subject:** {meta.subject}")
                if meta.author:
                    text_parts.append(f"**Author:** {meta.author}")
                if text_parts:
                    text_parts.append("\n---\n")
        except Exception:
            pass
        
        # Extract from WordDocument stream
        if ole.exists('WordDocument'):
            try:
                wd_stream = ole.openstream('WordDocument')
                wd_data = wd_stream.read()
                
                # Parse Word binary format
                extracted = _extract_word_text(wd_data)
                if extracted and len(extracted.strip()) > 50:
                    text_parts.append(extracted)
            except Exception as e:
                errors.append(f"WordDocument stream: {e}")
        
        # Try other text-containing streams
        for stream_path in ole.listdir():
            stream_name = '/'.join(stream_path)
            if any(x in stream_name.lower() for x in ['text', 'content', 'data']):
                try:
                    stream = ole.openstream(stream_path)
                    data = stream.read()
                    extracted = _extract_text_from_ole_stream(data)
                    if extracted and len(extracted.strip()) > 30:
                        text_parts.append(extracted)
                except Exception:
                    continue
        
        ole.close()
        
        if text_parts:
            text = '\n\n'.join(text_parts)
            # Clean up
            text = re.sub(r'\s{3,}', '\n\n', text)
            text = re.sub(r'[^\x20-\x7E\n\r\t\u00A0-\u00FF]+', ' ', text)
            
            return ExtractionResult(
                success=True,
                text=text,
                method="olefile",
                word_count=len(text.split()),
                char_count=len(text),
                metadata={'format': 'doc_legacy', 'extraction': 'ole_parsing'}
            )
            
    except ImportError:
        errors.append("olefile not available - install with: pip install olefile")
    except Exception as e:
        errors.append(f"OLE parsing failed: {e}")
    
    # Strategy 2: Enhanced binary extraction with encoding detection
    try:
        with open(file_path, 'rb') as f:
            content = f.read()
        
        text = None
        
        # Try UTF-16LE (common in Word docs)
        try:
            decoded = content.decode('utf-16-le', errors='ignore')
            text = ''.join(c for c in decoded if c.isprintable() or c in '\n\r\t')
            if len(text.strip()) < 100:
                text = None
        except Exception:
            pass
        
        # Try CP1252 (Windows default encoding)
        if not text:
            try:
                text_parts = []
                current = []
                
                for byte in content:
                    if 32 <= byte <= 126 or byte in (9, 10, 13):
                        current.append(chr(byte))
                    elif 128 <= byte <= 255:
                        # Windows-1252 extended characters
                        try:
                            current.append(bytes([byte]).decode('cp1252'))
                        except Exception:
                            if len(current) > 15:
                                text_parts.append(''.join(current))
                            current = []
                    else:
                        if len(current) > 15:
                            text_parts.append(''.join(current))
                        current = []
                
                if current and len(current) > 15:
                    text_parts.append(''.join(current))
                
                text = '\n'.join(text_parts)
            except Exception:
                pass
        
        if text and len(text.strip()) > 50:
            # Clean and structure the text
            text = re.sub(r'(.)\1{5,}', r'\1', text)  # Remove repeated chars
            text = re.sub(r'\s+', ' ', text)
            
            # Try to split into paragraphs
            sentences = re.split(r'(?<=[.!?])\s+', text)
            paragraphs = [' '.join(sentences[i:i+5]) for i in range(0, len(sentences), 5)]
            text = '\n\n'.join(paragraphs)
            
            return ExtractionResult(
                success=True,
                text=f"[Legacy DOC - Binary Extraction]\n\n{text}",
                method=ExtractionMethod.FALLBACK.value,
                word_count=len(text.split()),
                char_count=len(text),
                metadata={'warning': 'Binary extraction - some content may be missing', 'errors': errors}
            )
        
    except Exception as e:
        errors.append(f"Binary extraction failed: {e}")
    
    # All strategies failed
    raise ExtractionError(
        f"Legacy .doc extraction failed. Errors: {'; '.join(errors)}. Consider converting to .docx",
        "doc_legacy"
    )


def _extract_word_text(data: bytes) -> str:
    """Extract text from Word binary document stream."""
    import struct
    
    text_parts = []
    
    # Method 1: Look for Unicode text runs (UTF-16LE)
    try:
        i = 0
        current_run = []
        
        while i < len(data) - 1:
            char_code = struct.unpack('<H', data[i:i+2])[0]
            
            if 32 <= char_code <= 126 or char_code in (9, 10, 13):
                current_run.append(chr(char_code))
            elif 0x00A0 <= char_code <= 0x00FF:  # Latin-1 supplement
                current_run.append(chr(char_code))
            elif 0x0100 <= char_code <= 0x017F:  # Latin Extended-A
                current_run.append(chr(char_code))
            else:
                if len(current_run) > 5:
                    text_parts.append(''.join(current_run))
                current_run = []
            
            i += 2
        
        if current_run and len(current_run) > 5:
            text_parts.append(''.join(current_run))
            
    except Exception:
        pass
    
    # Method 2: Extract CP1252 text
    try:
        cp_parts = []
        current = []
        
        for byte in data:
            if 32 <= byte <= 126 or byte in (9, 10, 13):
                current.append(chr(byte))
            elif 128 <= byte <= 255:
                try:
                    current.append(bytes([byte]).decode('cp1252'))
                except Exception:
                    if len(current) > 20:
                        cp_parts.append(''.join(current))
                    current = []
            else:
                if len(current) > 20:
                    cp_parts.append(''.join(current))
                current = []
        
        if current and len(current) > 20:
            cp_parts.append(''.join(current))
        
        # Use whichever extraction got more text
        if cp_parts:
            cp_text = '\n'.join(cp_parts)
            if len(cp_text) > len(' '.join(text_parts)):
                text_parts = cp_parts
                
    except Exception:
        pass
    
    # Deduplicate and clean
    seen = set()
    unique_parts = []
    for part in text_parts:
        cleaned = re.sub(r'\s+', ' ', part).strip()
        if cleaned and cleaned not in seen and len(cleaned) > 3:
            seen.add(cleaned)
            unique_parts.append(cleaned)
    
    return '\n\n'.join(unique_parts)


def _extract_text_from_ole_stream(data: bytes) -> str:
    """Extract readable text from arbitrary OLE stream data."""
    text_parts = []
    current = []
    
    for byte in data:
        if 32 <= byte <= 126 or byte in (9, 10, 13):
            current.append(chr(byte))
        else:
            if len(current) > 20:  # Minimum meaningful text
                text_parts.append(''.join(current))
            current = []
    
    if current and len(current) > 20:
        text_parts.append(''.join(current))
    
    return '\n'.join(text_parts)


# ============================================================================
# EXCEL EXTRACTION
# ============================================================================

def extract_xlsx(file_path: str) -> ExtractionResult:
    """
    Extract data from XLSX with all sheets and formulas.
    
    Features:
        - All worksheets
        - Named ranges
        - Formula results
        - Cell formatting context
    """
    try:
        from openpyxl import load_workbook
        
        wb = load_workbook(file_path, data_only=True)
        text_parts = []
        table_count = 0
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f"\n## Sheet: {sheet_name}\n")
            
            # Get dimensions
            if ws.max_row and ws.max_column:
                rows = []
                
                for row in ws.iter_rows(min_row=1, max_row=min(ws.max_row, 1000)):
                    cells = []
                    for cell in row:
                        value = cell.value
                        if value is not None:
                            cells.append(str(value).strip())
                        else:
                            cells.append('')
                    
                    # Skip completely empty rows
                    if any(cells):
                        rows.append('| ' + ' | '.join(cells) + ' |')
                
                if rows:
                    # Add header separator after first row
                    if len(rows) > 1:
                        col_count = len(rows[0].split('|')) - 2
                        header_sep = '| ' + ' | '.join(['---'] * col_count) + ' |'
                        rows.insert(1, header_sep)
                    
                    text_parts.append('\n'.join(rows))
                    table_count += 1
            else:
                text_parts.append("(Empty sheet)")
        
        wb.close()
        text = '\n'.join(text_parts)
        
        return ExtractionResult(
            success=True,
            text=text,
            method=ExtractionMethod.OPENPYXL.value,
            sheets=len(wb.sheetnames),
            tables=table_count,
            word_count=len(text.split()),
            char_count=len(text)
        )
        
    except ImportError:
        raise ExtractionError("openpyxl library not available", "xlsx")
    except Exception as e:
        raise ExtractionError(f"XLSX extraction failed: {str(e)}", "xlsx")


def extract_xls_legacy(file_path: str) -> ExtractionResult:
    """Extract data from legacy XLS files."""
    try:
        import xlrd
        
        wb = xlrd.open_workbook(file_path)
        text_parts = []
        table_count = 0
        
        for sheet_idx in range(wb.nsheets):
            sheet = wb.sheet_by_index(sheet_idx)
            text_parts.append(f"\n## Sheet: {sheet.name}\n")
            
            rows = []
            for row_idx in range(min(sheet.nrows, 1000)):
                cells = []
                for col_idx in range(sheet.ncols):
                    value = sheet.cell_value(row_idx, col_idx)
                    cells.append(str(value).strip() if value else '')
                
                if any(cells):
                    rows.append('| ' + ' | '.join(cells) + ' |')
            
            if rows:
                if len(rows) > 1:
                    col_count = len(rows[0].split('|')) - 2
                    header_sep = '| ' + ' | '.join(['---'] * col_count) + ' |'
                    rows.insert(1, header_sep)
                
                text_parts.append('\n'.join(rows))
                table_count += 1
        
        text = '\n'.join(text_parts)
        
        return ExtractionResult(
            success=True,
            text=text,
            method=ExtractionMethod.XLRD.value,
            sheets=wb.nsheets,
            tables=table_count,
            word_count=len(text.split()),
            char_count=len(text)
        )
        
    except ImportError:
        raise ExtractionError("xlrd library not available", "xls_legacy")
    except Exception as e:
        raise ExtractionError(f"XLS extraction failed: {str(e)}", "xls_legacy")


# ============================================================================
# POWERPOINT EXTRACTION
# ============================================================================

def extract_pptx(file_path: str) -> ExtractionResult:
    """
    Extract text from PPTX with speaker notes.
    
    Features:
        - All slides
        - Text frames
        - Tables
        - Speaker notes
        - Slide titles
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
        
        prs = Presentation(file_path)
        text_parts = []
        table_count = 0
        slide_count = 0
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_count += 1
            text_parts.append(f"\n## Slide {slide_num}\n")
            
            # Extract title
            if slide.shapes.title and slide.shapes.title.text:
                text_parts.append(f"### {slide.shapes.title.text}\n")
            
            # Extract all text frames
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if shape != slide.shapes.title:
                        text_parts.append(shape.text)
                
                # Extract tables
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append('| ' + ' | '.join(cells) + ' |')
                    
                    if rows and len(rows) > 1:
                        col_count = len(table.columns)
                        header_sep = '| ' + ' | '.join(['---'] * col_count) + ' |'
                        rows.insert(1, header_sep)
                    
                    text_parts.append('\n' + '\n'.join(rows) + '\n')
                    table_count += 1
            
            # Extract speaker notes
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame and notes_frame.text.strip():
                    text_parts.append(f"\n**Speaker Notes:** {notes_frame.text}\n")
        
        text = '\n'.join(text_parts)
        
        return ExtractionResult(
            success=True,
            text=text,
            method=ExtractionMethod.PYTHON_PPTX.value,
            slides=slide_count,
            tables=table_count,
            word_count=len(text.split()),
            char_count=len(text)
        )
        
    except ImportError:
        raise ExtractionError("python-pptx library not available", "pptx")
    except Exception as e:
        raise ExtractionError(f"PPTX extraction failed: {str(e)}", "pptx")


def extract_ppt_legacy(file_path: str) -> ExtractionResult:
    """
    Extract text from legacy .ppt files using OLE parsing.
    
    PowerPoint .ppt files store text in 'PowerPoint Document' stream
    and other OLE containers.
    """
    errors = []
    
    # Strategy 1: Use olefile for proper OLE parsing
    try:
        import olefile
        
        ole = olefile.OleFileIO(file_path)
        text_parts = []
        slide_texts = []
        
        # Get document metadata
        try:
            meta = ole.get_metadata()
            if meta:
                if meta.title:
                    text_parts.append(f"# {meta.title}")
                if meta.subject:
                    text_parts.append(f"**Subject:** {meta.subject}")
                if meta.author:
                    text_parts.append(f"**Author:** {meta.author}")
                if text_parts:
                    text_parts.append("\n---\n")
        except Exception:
            pass
        
        # Extract from PowerPoint Document stream
        if ole.exists('PowerPoint Document'):
            try:
                pp_stream = ole.openstream('PowerPoint Document')
                pp_data = pp_stream.read()
                
                # Parse PowerPoint binary format
                extracted = _extract_powerpoint_text(pp_data)
                if extracted:
                    slide_texts.append(extracted)
            except Exception as e:
                errors.append(f"PowerPoint stream: {e}")
        
        # Try other text-containing streams
        for stream_path in ole.listdir():
            stream_name = '/'.join(stream_path)
            if any(x in stream_name.lower() for x in ['text', 'slide', 'content']):
                try:
                    stream = ole.openstream(stream_path)
                    data = stream.read()
                    extracted = _extract_text_from_ole_stream(data)
                    if extracted and len(extracted.strip()) > 30:
                        slide_texts.append(f"[{stream_name}]\n{extracted}")
                except Exception:
                    continue
        
        ole.close()
        
        if slide_texts:
            text_parts.extend(slide_texts)
            text = '\n\n'.join(text_parts)
            # Clean up
            text = re.sub(r'\s{3,}', '\n\n', text)
            
            return ExtractionResult(
                success=True,
                text=text,
                method="olefile",
                word_count=len(text.split()),
                char_count=len(text),
                metadata={'format': 'ppt_legacy', 'extraction': 'ole_parsing'}
            )
            
    except ImportError:
        errors.append("olefile not available")
    except Exception as e:
        errors.append(f"OLE parsing failed: {e}")
    
    # Strategy 2: Enhanced binary extraction
    try:
        with open(file_path, 'rb') as f:
            content = f.read()
        
        text_parts = []
        
        # Try UTF-16LE (common in PowerPoint)
        try:
            decoded = content.decode('utf-16-le', errors='ignore')
            clean = ''.join(c for c in decoded if c.isprintable() or c in '\n\r\t')
            if len(clean.strip()) > 100:
                text_parts.append(clean)
        except Exception:
            pass
        
        # Extract CP1252/ASCII text
        current = []
        for byte in content:
            if 32 <= byte <= 126 or byte in (9, 10, 13):
                current.append(chr(byte))
            elif 128 <= byte <= 255:
                try:
                    current.append(bytes([byte]).decode('cp1252'))
                except Exception:
                    if len(current) > 15:
                        text_parts.append(''.join(current))
                    current = []
            else:
                if len(current) > 15:
                    text_parts.append(''.join(current))
                current = []
        
        if current and len(current) > 15:
            text_parts.append(''.join(current))
        
        if text_parts:
            # Deduplicate
            seen = set()
            unique = []
            for part in text_parts:
                cleaned = re.sub(r'\s+', ' ', part).strip()
                if cleaned and cleaned not in seen and len(cleaned) > 10:
                    seen.add(cleaned)
                    unique.append(cleaned)
            
            # Organize into slides
            result_parts = []
            slide_num = 1
            current_slide = []
            
            for part in unique:
                if len(part) > 100 or len(current_slide) >= 5:
                    if current_slide:
                        result_parts.append(f"## Slide {slide_num}\n" + '\n'.join(current_slide))
                        slide_num += 1
                    current_slide = [part]
                else:
                    current_slide.append(part)
            
            if current_slide:
                result_parts.append(f"## Slide {slide_num}\n" + '\n'.join(current_slide))
            
            text = '\n\n'.join(result_parts) if result_parts else '\n\n'.join(unique)
            
            if len(text.strip()) > 50:
                return ExtractionResult(
                    success=True,
                    text=f"[Legacy PPT - Binary Extraction]\n\n{text}",
                    method=ExtractionMethod.FALLBACK.value,
                    word_count=len(text.split()),
                    char_count=len(text),
                    metadata={'warning': 'Binary extraction - some content may be missing', 'errors': errors}
                )
        
    except Exception as e:
        errors.append(f"Binary extraction failed: {e}")
    
    raise ExtractionError(
        f"Legacy .ppt extraction failed. Errors: {'; '.join(errors)}. Consider converting to .pptx",
        "ppt_legacy"
    )


def _extract_powerpoint_text(data: bytes) -> str:
    """Extract text from PowerPoint binary document stream."""
    import struct
    
    text_parts = []
    
    # PowerPoint text is often UTF-16LE
    try:
        i = 0
        current = []
        
        while i < len(data) - 1:
            char_code = struct.unpack('<H', data[i:i+2])[0]
            
            if char_code == 0x000D:  # Carriage return - often slide separator
                if current and len(current) > 3:
                    text_parts.append(''.join(current).strip())
                current = []
            elif 32 <= char_code <= 126:
                current.append(chr(char_code))
            elif char_code in (9, 10):
                current.append(' ')
            elif 0x00A0 <= char_code <= 0x00FF:
                current.append(chr(char_code))
            else:
                if len(current) > 5:
                    text_parts.append(''.join(current).strip())
                current = []
            
            i += 2
        
        if current and len(current) > 3:
            text_parts.append(''.join(current).strip())
            
    except Exception:
        pass
    
    # Deduplicate
    seen = set()
    unique = []
    for part in text_parts:
        if part and part not in seen and len(part) > 3:
            seen.add(part)
            unique.append(part)
    
    return '\n'.join(unique)


# ============================================================================
# TEXT FILE EXTRACTION
# ============================================================================

def extract_text_file(file_path: str) -> ExtractionResult:
    """
    Extract text from plain text files with encoding detection.
    
    Supports: .txt, .md, .csv, .json, .log
    """
    try:
        # Try to detect encoding using chardet if available
        try:
            import chardet
        except Exception:
            chardet = None

        raw = None
        if chardet:
            with open(file_path, 'rb') as f:
                raw = f.read()
            detected = chardet.detect(raw)
            encoding = detected.get('encoding') or 'utf-8'
        else:
            encoding = 'utf-8'

    except Exception:
        encoding = 'utf-8'
        raw = None

    encodings_to_try = [encoding, 'utf-8', 'utf-16', 'latin-1', 'cp1252']

    for enc in encodings_to_try:
        try:
            if raw is not None:
                text = raw.decode(enc)
            else:
                with open(file_path, 'r', encoding=enc, errors='ignore') as f:
                    text = f.read()

            return ExtractionResult(
                success=True,
                text=text,
                method=ExtractionMethod.NATIVE.value,
                word_count=len(text.split()),
                char_count=len(text),
                metadata={'encoding': enc}
            )

        except (UnicodeDecodeError, LookupError):
            continue

    raise ExtractionError("Could not decode text file with any encoding", "text")


# ============================================================================
# HTML EXTRACTION
# ============================================================================

def extract_html(file_path: str) -> ExtractionResult:
    """Extract text from HTML with structure preservation."""
    try:
        from bs4 import BeautifulSoup
        
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        soup = BeautifulSoup(content, 'html.parser')
        
        # Remove script and style elements
        for element in soup(['script', 'style', 'meta', 'link']):
            element.decompose()
        
        # Extract title
        text_parts = []
        if soup.title and soup.title.string:
            text_parts.append(f"# {soup.title.string}\n")
        
        # Convert headers
        for i in range(1, 7):
            for header in soup.find_all(f'h{i}'):
                header.replace_with(f"\n{'#' * i} {header.get_text()}\n")
        
        # Convert lists
        for ul in soup.find_all('ul'):
            for li in ul.find_all('li'):
                li.replace_with(f"- {li.get_text()}\n")
        
        for ol in soup.find_all('ol'):
            for idx, li in enumerate(ol.find_all('li'), 1):
                li.replace_with(f"{idx}. {li.get_text()}\n")
        
        # Extract text
        text = soup.get_text(separator='\n', strip=True)
        text_parts.append(text)
        
        final_text = '\n'.join(text_parts)
        
        return ExtractionResult(
            success=True,
            text=final_text,
            method=ExtractionMethod.BEAUTIFULSOUP.value,
            word_count=len(final_text.split()),
            char_count=len(final_text)
        )
        
    except ImportError:
        # Fallback: basic regex extraction
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        # Remove tags
        text = re.sub(r'<script[^>]*>.*?</script>', '', content, flags=re.DOTALL)
        text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL)
        text = re.sub(r'<[^>]+>', ' ', text)
        text = re.sub(r'\s+', ' ', text)
        
        return ExtractionResult(
            success=True,
            text=text.strip(),
            method=ExtractionMethod.NATIVE.value,
            word_count=len(text.split()),
            char_count=len(text)
        )
        
    except Exception as e:
        raise ExtractionError(f"HTML extraction failed: {str(e)}", "html")


# ============================================================================
# XML/XSD EXTRACTION
# ============================================================================

def extract_xml(file_path: str) -> ExtractionResult:
    """
    Extract content from XML/XSD files with structure.
    
    Features:
        - Element hierarchy
        - Attribute extraction
        - Text content
        - Schema documentation
    """
    try:
        import xml.etree.ElementTree as ET
        
        tree = ET.parse(file_path)
        root = tree.getroot()
        
        text_parts = []
        text_parts.append(f"# XML Document: {root.tag}\n")
        
        def extract_element(elem, depth=0):
            """Recursively extract element content."""
            parts = []
            indent = "  " * depth
            
            # Element tag
            tag_name = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
            parts.append(f"{indent}**{tag_name}**")
            
            # Attributes
            if elem.attrib:
                attrs = ', '.join(f'{k}="{v}"' for k, v in elem.attrib.items())
                parts.append(f" ({attrs})")
            
            # Text content
            if elem.text and elem.text.strip():
                parts.append(f": {elem.text.strip()}")
            
            parts.append("\n")
            
            # Children
            for child in elem:
                parts.extend(extract_element(child, depth + 1))
            
            return parts
        
        text_parts.extend(extract_element(root))
        text = ''.join(text_parts)
        
        return ExtractionResult(
            success=True,
            text=text,
            method=ExtractionMethod.XML_PARSER.value,
            word_count=len(text.split()),
            char_count=len(text)
        )
        
    except Exception as e:
        raise ExtractionError(f"XML extraction failed: {str(e)}", "xml")


# ============================================================================
# IMAGE EXTRACTION (OCR-enabled using ImageExtractor class)
# ============================================================================

def extract_image(file_path: str) -> ExtractionResult:
    """
    Extract text from images using OCR (Claude Vision or Textract).
    
    Features:
        - OCR text extraction via Claude Vision (Bedrock)
        - Fallback to AWS Textract if Claude unavailable
        - Image metadata extraction (dimensions, format)
    """
    try:
        # Use the ImageExtractor class for OCR-enabled extraction
        config = ExtractorConfig()
        extractor = ImageExtractor(config)
        
        # Extract with OCR
        extraction_dict = extractor.extract(file_path, document_id='temp')
        
        # Build text from elements
        text_parts = []
        elements = extraction_dict.get('elements', [])
        metadata_dict = extraction_dict.get('metadata', {})
        
        # Extract OCR text from elements
        for elem in elements:
            if elem.get('type') == 'text' and elem.get('content'):
                text_parts.append(elem['content'])
        
        # If no OCR text, add basic image metadata
        if not text_parts:
            try:
                from PIL import Image
                img = Image.open(file_path)
                width, height = img.width, img.height
                fmt = img.format or 'Unknown'
                img.close()
                
                # Create detailed metadata text (ensure MIN_CHUNK_SIZE of 150 chars)
                text_parts.append(f"# Image Document\n\n")
                text_parts.append(f"This is an image file with the following properties:\n\n")
                text_parts.append(f"- **Format:** {fmt} image\n")
                text_parts.append(f"- **Dimensions:** {width} x {height} pixels\n")
                text_parts.append(f"- **File:** {os.path.basename(file_path)}\n")
                text_parts.append(f"- **Size:** {os.path.getsize(file_path):,} bytes\n\n")
                text_parts.append(f"Note: This image was processed but no readable text was detected through OCR (Optical Character Recognition). ")
                text_parts.append(f"The image may contain graphics, photographs, or visual content without text elements.")
                
                metadata_dict.update({
                    'width': width,
                    'height': height,
                    'format': fmt,
                    'has_text': False
                })
            except Exception as e:
                # Fallback to basic file info
                import os
                file_size = os.path.getsize(file_path)
                ext = os.path.splitext(file_path)[1].upper()
                filename = os.path.basename(file_path)
                
                # Ensure sufficient content for chunking (MIN_CHUNK_SIZE = 150)
                text_parts.append(f"# Image Document\n\n")
                text_parts.append(f"This is an image file in {ext} format.\n\n")
                text_parts.append(f"- **File Name:** {filename}\n")
                text_parts.append(f"- **Format:** {ext} image file\n")
                text_parts.append(f"- **Size:** {file_size:,} bytes\n\n")
                text_parts.append(f"Note: This image was uploaded to the system but detailed metadata extraction was unavailable. ")
                text_parts.append(f"The image content has been archived and is available for reference.")
                
                metadata_dict.update({
                    'format': ext.replace('.', ''),
                    'file_size_bytes': file_size,
                    'has_text': False
                })
        
        text = '\n\n'.join(text_parts)
        
        # Determine extraction method
        ocr_performed = metadata_dict.get('ocr_performed', False)
        ocr_text_length = metadata_dict.get('ocr_text_length', 0)
        
        if ocr_performed and ocr_text_length > 0:
            method = ExtractionMethod.IMAGE_METADATA.value  # Could be enhanced to "claude_vision" or "textract"
        else:
            method = ExtractionMethod.IMAGE_METADATA.value
        
        return ExtractionResult(
            success=True,
            text=text,
            method=method,
            images=1,
            word_count=len(text.split()),
            char_count=len(text),
            metadata=metadata_dict
        )

    except Exception as e:
        logger.error(f"Image extraction failed: {str(e)}", exc_info=True)
        raise ExtractionError(f"Image extraction failed: {str(e)}", "image")


# ============================================================================
# MAIN EXTRACTION ROUTER
# ============================================================================

def extract_document(file_path: str, file_type: str) -> ExtractionResult:
    """
    Route extraction to appropriate handler based on file type.
    """
    extractors = {
        'pdf': extract_pdf,
        'docx': extract_docx,
        'doc_legacy': extract_doc_legacy,
        'xlsx': extract_xlsx,
        'xls_legacy': extract_xls_legacy,
        'pptx': extract_pptx,
        'ppt_legacy': extract_ppt_legacy,
        'text': extract_text_file,
        'html': extract_html,
        'xml': extract_xml,
        'image': extract_image,
    }
    
    extractor = extractors.get(file_type)
    
    if not extractor:
        raise ExtractionError(f"No extractor for file type: {file_type}", file_type)
    
    return extractor(file_path)


# ==========================================================================
# STRUCTURED EXTRACTION + ADVANCED CHUNKING
# ==========================================================================

class ExtractorConfig:
    """Runtime configuration for structured extractors."""
    TEMP_DIR = TEMP_DIR
    OCR_ENABLED = OCR_ENABLED
    TEXTRACT_ENABLED = TEXTRACT_ENABLED


class ChunkerConfig:
    """Runtime configuration for ProductionChunker."""
    TARGET_CHUNK_SIZE = TARGET_CHUNK_SIZE
    CHUNK_OVERLAP = CHUNK_OVERLAP_TOKENS
    MAX_CHUNK_SIZE = MAX_CHUNK_SIZE
    MIN_CHUNK_SIZE = MIN_TOKEN_CHUNK_SIZE
    TABLE_MAX_TOKENS = TABLE_MAX_TOKENS
    TABLE_ROWS_PER_CHUNK = TABLE_ROWS_PER_CHUNK
    CODE_TARGET_LINES = CODE_TARGET_LINES
    CODE_OVERLAP_LINES = CODE_OVERLAP_LINES
    ASSEMBLY_PROCEDURE_MARKERS = ['PROC', 'ENDP', 'CSECT', 'DSECT']


def _build_text_sample_from_elements(elements: List[Dict], max_chars: int = 5000) -> str:
    """Build a text sample from structured elements for metadata extraction."""
    parts = []
    total = 0
    for elem in elements:
        content = elem.get('content', '')
        if not content:
            continue
        remaining = max_chars - total
        if remaining <= 0:
            break
        parts.append(content[:remaining])
        total += len(parts[-1])
    return '\n'.join(parts)


def extract_structured_elements(file_path: str, file_type: str, document_id: str) -> Optional[Dict[str, Any]]:
    """Extract structured elements (tables, code, images) when supported."""
    if not ENABLE_ADVANCED_CHUNKING:
        return None

    try:
        os.makedirs(TEMP_DIR, exist_ok=True)
    except Exception:
        pass

    try:
        config = ExtractorConfig()

        if file_type in ['pdf']:
            return PDFExtractor(config).extract(file_path, document_id)

        if file_type in ['docx', 'doc']:
            try:
                return EnhancedDOCXExtractor(config).extract(file_path, document_id)
            except Exception:
                return DOCXExtractor(config).extract(file_path, document_id)

        if file_type in ['xlsx', 'xls']:
            return XLSXExtractor(config).extract(file_path, document_id)

        if file_type in ['pptx', 'ppt']:
            return PPTXExtractor(config).extract(file_path, document_id)

        if file_type in ['txt', 'md', 'csv', 'json', 'xml', 'xsd', 'html', 'htm']:
            return TextExtractor(config).extract(file_path, document_id)

        if file_type in ['png', 'jpg', 'jpeg', 'tiff', 'bmp']:
            return ImageExtractor(config).extract(file_path, document_id)

    except Exception as e:
        logger.warning(f"Structured extraction failed for {file_type}: {e}")

    return None


# ============================================================================
# FILE-TYPE SPECIFIC CHUNKING STRATEGY DISPATCHER
# ============================================================================

# Configuration: Enable/disable hybrid chunking
ENABLE_HYBRID_CHUNKING = os.environ.get('ENABLE_HYBRID_CHUNKING', 'true').lower() == 'true'

# File types that benefit from document structure-aware hybrid chunking
STRUCTURED_DOC_TYPES = {'pdf', 'docx', 'doc', 'pptx', 'ppt'}

# Code file types that need procedure-aware chunking (mainframe + modern)
CODE_FILE_TYPES = {
    # Mainframe
    'asm', 'jcl', 'cbl', 'cob', 'cpy', 'prc', 'clist', 'rexx', 'rex', 'proc',
    # Modern (extend if needed)
    'py', 'java', 'js', 'ts', 'c', 'cpp', 'h', 'hpp', 'cs', 'go', 'rs', 'rb'
}

# Tabular data files
TABULAR_FILE_TYPES = {'xlsx', 'xls', 'csv', 'tsv'}

# Text files that benefit from semantic chunking (use hybrid for all)
TEXT_FILE_TYPES = {'txt', 'md', 'html', 'htm', 'xml', 'json', 'yaml', 'yml', 'rst', 'tex'}

# Image file types (OCR-based chunking)
IMAGE_FILE_TYPES = {'png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'tif', 'webp'}


def get_chunking_strategy(file_type: str) -> str:
    """
    Determine the best chunking strategy based on file type.
    
    Strategy Priority (all use parent-child hierarchy):
        'code': MainframeCodeChunker for procedure-aware chunking
        'tabular': Row-based chunking preserving table structure
        'hybrid': HybridChunker (recursive + semantic + parent-child) - DEFAULT
    
    Returns:
        Strategy name string
    """
    file_type_lower = file_type.lower().lstrip('.')
    
    # Code files need procedure-boundary awareness
    if file_type_lower in CODE_FILE_TYPES:
        return 'code'
    
    # Tabular files need row-based preservation
    if file_type_lower in TABULAR_FILE_TYPES:
        return 'tabular'
    
    # All other files (structured docs, text, etc.) use hybrid
    # This includes: pdf, docx, txt, md, html, xml, json, etc.
    return 'hybrid'


# ============================================================================
# TABULAR DATA CHUNKING (XLSX, CSV, etc.)
# ============================================================================

# Configuration for tabular chunking
TABULAR_ROWS_PER_CHUNK = int(os.environ.get('TABULAR_ROWS_PER_CHUNK', '50'))
TABULAR_MAX_TOKENS = int(os.environ.get('TABULAR_MAX_TOKENS', '800'))


def _chunk_tabular_data(elements: List[Dict], document_id: str, metadata: Dict) -> List[Dict]:
    """
    Chunk tabular data (Excel, CSV) preserving row integrity with headers.
    
    Strategy:
    1. Identify header row (first row or explicitly marked)
    2. Group rows into chunks of TABULAR_ROWS_PER_CHUNK
    3. Prepend header to each chunk for context
    4. Create parent-child hierarchy (sheet → row groups)
    """
    if not elements:
        return []
    
    chunks = []
    chunk_counter = 0
    
    # Group elements by sheet/table
    sheets = {}
    for elem in elements:
        sheet_name = elem.get('metadata', {}).get('sheet_name', 'Sheet1')
        if sheet_name not in sheets:
            sheets[sheet_name] = []
        sheets[sheet_name].append(elem)
    
    for sheet_name, sheet_elements in sheets.items():
        # Extract header (first row or marked as header)
        header_row = None
        data_rows = []
        
        for elem in sheet_elements:
            elem_meta = elem.get('metadata', {})
            content = elem.get('content', '').strip()
            
            if not content:
                continue
            
            is_header = elem_meta.get('is_header', False) or elem_meta.get('row_index', 999) == 0
            
            if is_header and header_row is None:
                header_row = content
            else:
                data_rows.append({
                    'content': content,
                    'row_index': elem_meta.get('row_index', len(data_rows)),
                    'metadata': elem_meta
                })
        
        # Create sheet-level parent chunk
        chunk_counter += 1
        sheet_chunk_id = f"{document_id}-sheet-{chunk_counter:04d}"
        sheet_content = f"Sheet: {sheet_name}\n"
        if header_row:
            sheet_content += f"Headers: {header_row}\n"
        sheet_content += f"Total rows: {len(data_rows)}"
        
        sheet_chunk = {
            'chunk_id': sheet_chunk_id,
            'document_id': document_id,
            'type': 'section',
            'content': sheet_content,
            'heading': sheet_name,
            'heading_path': [sheet_name],
            'is_section': True,
            'is_table': True,
            'child_ids': [],
            'metadata': {
                **metadata,
                'sheet_name': sheet_name,
                'row_count': len(data_rows),
                'has_header': bool(header_row)
            }
        }
        chunks.append(sheet_chunk)
        
        # Create row-group chunks (children)
        for i in range(0, len(data_rows), TABULAR_ROWS_PER_CHUNK):
            row_batch = data_rows[i:i + TABULAR_ROWS_PER_CHUNK]
            
            chunk_counter += 1
            row_chunk_id = f"{document_id}-rows-{chunk_counter:04d}"
            
            # Build content with header context
            content_parts = []
            if header_row:
                content_parts.append(f"[Headers: {header_row}]")
            
            for row in row_batch:
                content_parts.append(row['content'])
            
            row_content = '\n'.join(content_parts)
            
            row_chunk = {
                'chunk_id': row_chunk_id,
                'document_id': document_id,
                'type': 'semantic',
                'content': row_content,
                'parent_chunk_id': sheet_chunk_id,
                'section_id': sheet_chunk_id,
                'heading_path': [sheet_name],
                'is_semantic': True,
                'is_table': True,
                'metadata': {
                    **metadata,
                    'sheet_name': sheet_name,
                    'row_start': i,
                    'row_end': i + len(row_batch),
                    'row_count': len(row_batch),
                    'parent_chunk_id': sheet_chunk_id
                }
            }
            chunks.append(row_chunk)
            sheet_chunk['child_ids'].append(row_chunk_id)
    
    return chunks


def create_chunks_from_elements(
    elements: List[Dict],
    document_id: str,
    metadata: Dict = None,
    document_title: str = None,
    document_author: str = None,
    document_type: str = None
) -> List[Chunk]:
    """
    Create advanced chunks from structured elements using file-type specific strategy.
    
    ALL strategies produce 3-level hierarchies with parent-child relationships:
    - CODE files: MainframeCodeChunker (procedure/function boundaries)
    - TABULAR files: Row-based chunking (preserves table structure)
    - ALL OTHERS: HybridChunker (recursive + semantic + parent-child)
    
    This unified approach ensures consistent chunk structure across all file types.
    """
    if not elements:
        return []
    
    metadata = metadata or {}
    file_type = document_type or metadata.get('file_type', 'txt')
    file_name = metadata.get('file_name', f'{document_id}.{file_type}')
    
    # Normalize content encoding (important for mainframe docs)
    for elem in elements:
        content = elem.get('content', '')
        if content:
            elem['content'] = normalize_mainframe_encoding(content)
    
    # Select chunking strategy
    strategy = get_chunking_strategy(file_type)
    logger.info(f"Chunking strategy: {strategy} for {file_type} ({len(elements)} elements)")
    
    chunk_dicts = []
    
    # =========================================================================
    # STRATEGY 1: CODE FILES - Procedure-aware chunking
    # =========================================================================
    if strategy == 'code':
        try:
            # Combine elements to get full code content
            code_content = '\n'.join([e.get('content', '') for e in elements])
            
            chunk_dicts = mainframe_code_chunker.chunk_code(
                content=code_content,
                document_id=document_id,
                file_extension=file_type,
                file_name=file_name,
                document_metadata=metadata
            )
            logger.info(f"MainframeCodeChunker created {len(chunk_dicts)} procedure-aware chunks")
        except Exception as e:
            logger.warning(f"MainframeCodeChunker failed, falling back to hybrid: {e}")
            strategy = 'hybrid'  # Fallback
    
    # =========================================================================
    # STRATEGY 2: TABULAR FILES - Row-based chunking
    # =========================================================================
    if strategy == 'tabular' and not chunk_dicts:
        try:
            chunk_dicts = _chunk_tabular_data(elements, document_id, metadata)
            logger.info(f"Tabular chunker created {len(chunk_dicts)} row-based chunks")
        except Exception as e:
            logger.warning(f"Tabular chunking failed, falling back to hybrid: {e}")
            strategy = 'hybrid'  # Fallback
    
    # =========================================================================
    # STRATEGY 3: HYBRID (DEFAULT) - Recursive + Semantic + Parent-Child
    # =========================================================================
    if (strategy == 'hybrid' or not chunk_dicts) and ENABLE_HYBRID_CHUNKING:
        try:
            chunk_hierarchy = hybrid_chunker.chunk_document(
                document_id=document_id,
                elements=elements,
                file_type=file_type,
                document_metadata=metadata
            )
            chunk_dicts = chunk_hierarchy.all_chunks()
            logger.info(f"HybridChunker created {len(chunk_dicts)} multi-level chunks ")
        except Exception as e:
            logger.warning(f"HybridChunker failed: {e}")
    
    # =========================================================================
    # FALLBACK: ProductionChunker (simple token-aware)
    # =========================================================================
    if not chunk_dicts:
        try:
            chunker = ProductionChunker(ChunkerConfig())
            chunk_dicts = chunker.chunk(elements, document_id, metadata)
            logger.info(f"ProductionChunker (fallback) created {len(chunk_dicts)} chunks")
        except Exception as e:
            logger.error(f"All chunking strategies failed: {e}")
            return []
    
    # Convert chunk dicts to Chunk objects
    chunks: List[Chunk] = []
    for idx, chunk in enumerate(chunk_dicts):
        chunk_meta = chunk.get('metadata', {})
        content = chunk.get('content', '')
        page_numbers = chunk_meta.get('page_numbers', chunk.get('page_numbers', []))
        
        # Handle different chunk types
        chunk_type = chunk.get('type', 'text')
        is_table = chunk_type == 'table' or chunk.get('is_table', False)
        is_image = chunk_type == 'image' or chunk.get('is_image', False)
        is_code = chunk_type == 'code' or chunk_meta.get('has_code', False)
        
        # Get heading hierarchy
        heading_path = chunk.get('heading_path', chunk_meta.get('heading_path'))
        section_title = chunk.get('heading', chunk_meta.get('section'))
        if heading_path and not section_title:
            section_title = heading_path[-1] if heading_path else None

        chunks.append(Chunk(
            document_id=document_id,
            chunk_id=chunk.get('chunk_id'),
            chunk_number=chunk_meta.get('chunk_index', idx),
            text=content,
            char_count=len(content),
            word_count=len(content.split()),
            page_number=page_numbers[0] if page_numbers else None,
            section_title=section_title,
            heading_hierarchy=heading_path,
            has_table=is_table or chunk_meta.get('has_tables', False),
            has_image=is_image or chunk_meta.get('has_images', False),
            has_list=False,
            has_code=is_code,
            document_title=document_title,
            document_author=document_author,
            document_type=document_type,
            metadata={
                **chunk_meta,
                'chunk_type': chunk_type,
                'parent_chunk_id': chunk.get('parent_chunk_id'),
                'section_id': chunk.get('section_id'),
                'child_ids': chunk.get('child_ids', []),
                'is_section': chunk.get('is_section', False),
                'is_semantic': chunk.get('is_semantic', False),
                'is_child': chunk.get('is_child', False),
                'tokens': chunk.get('tokens', 0)
            }
        ))

    return chunks


# ============================================================================
# ENHANCED CHUNKING WITH CONTEXT
# ============================================================================

def create_chunks(
    text: str, 
    document_id: str, 
    metadata: Dict = None,
    document_title: str = None,
    document_author: str = None,
    document_type: str = None
) -> List[Chunk]:
    """
    Create enhanced chunks with full context for improved RAG retrieval.
    
    Features:
        - Heading hierarchy tracking (breadcrumb navigation)
        - Content type indicators (table, image, list, code)
        - Adjacent chunk previews
        - Document metadata in each chunk
        - Page number tracking
        - MECH-specific metadata extraction (Layer 2)
        - Encoding normalization (Layer 1)
    """
    chunks = []
    
    # Layer 1: Normalize mainframe encodings BEFORE chunking
    text = normalize_mainframe_encoding(text)
    
    # Clean text
    text = re.sub(r'\n{3,}', '\n\n', text.strip())
    
    if not text:
        return chunks
    
    # Build heading index for context tracking
    heading_pattern = r'^(#+)\s+(.+)$'
    heading_positions = []
    
    for match in re.finditer(heading_pattern, text, re.MULTILINE):
        level = len(match.group(1))
        title = match.group(2).strip()
        heading_positions.append({
            'position': match.start(),
            'level': level,
            'title': title
        })
    
    def get_heading_hierarchy(position: int) -> tuple:
        """Get heading breadcrumb and current section at position."""
        hierarchy = []
        current_levels = {}
        
        for h in heading_positions:
            if h['position'] > position:
                break
            current_levels = {k: v for k, v in current_levels.items() if k < h['level']}
            current_levels[h['level']] = h['title']
        
        for level in sorted(current_levels.keys()):
            hierarchy.append(current_levels[level])
        
        section = hierarchy[-1] if hierarchy else None
        return hierarchy, section
    
    # Split by page markers
    page_pattern = r'\[Page (\d+)\]'
    parts = re.split(f'({page_pattern})', text)
    
    current_page = 1
    segments = []
    
    i = 0
    while i < len(parts):
        match = re.match(page_pattern, parts[i] if i < len(parts) else '')
        if match:
            current_page = int(match.group(1))
            i += 1
            if i < len(parts):
                segments.append({'text': parts[i], 'page': current_page})
                i += 1
        else:
            if parts[i].strip():
                segments.append({'text': parts[i], 'page': current_page})
            i += 1
    
    if not segments:
        segments = [{'text': text, 'page': 1}]
    
    # Create chunks with context
    chunk_number = 0
    all_chunk_texts = []
    text_position = 0
    
    for segment in segments:
        segment_text = segment['text']
        page_num = segment['page']
        
        paragraphs = re.split(r'\n\n+', segment_text)
        
        current_chunk_parts = []
        current_length = 0
        
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            
            para_length = len(para)
            
            if current_length + para_length + 2 > CHUNK_SIZE and current_chunk_parts:
                chunk_text = '\n\n'.join(current_chunk_parts)
                
                if len(chunk_text) >= MIN_CHUNK_SIZE:
                    # Detect content types
                    has_table = bool(re.search(r'\|[\s\S]*?\|[\s\S]*?\|', chunk_text))
                    has_image = bool(re.search(r'\[Image[^\]]*\]', chunk_text, re.I))
                    has_list = bool(re.search(r'^[\s]*[-*\d+\.]+\s', chunk_text, re.M))
                    has_code = bool(re.search(r'```', chunk_text))
                    
                    hierarchy, section = get_heading_hierarchy(text_position)
                    
                    chunk = Chunk(
                        document_id=document_id,
                        chunk_id=f"{document_id}_chunk_{chunk_number:04d}",
                        chunk_number=chunk_number,
                        text=chunk_text,
                        char_count=len(chunk_text),
                        word_count=len(chunk_text.split()),
                        page_number=page_num,
                        section_title=section,
                        heading_hierarchy=hierarchy if hierarchy else None,
                        has_table=has_table,
                        has_image=has_image,
                        has_list=has_list,
                        has_code=has_code,
                        document_title=document_title,
                        document_author=document_author,
                        document_type=document_type,
                        metadata=metadata
                    )
                    chunks.append(chunk)
                    all_chunk_texts.append(chunk_text)
                    chunk_number += 1
                
                overlap_chars = chunk_text[-CHUNK_OVERLAP:] if len(chunk_text) > CHUNK_OVERLAP else chunk_text
                current_chunk_parts = [overlap_chars]
                current_length = len(overlap_chars)
            
            current_chunk_parts.append(para)
            current_length += para_length + 2
            text_position += para_length + 2
        
        # Final chunk from segment
        if current_chunk_parts:
            chunk_text = '\n\n'.join(current_chunk_parts)
            
            if len(chunk_text) >= MIN_CHUNK_SIZE:
                has_table = bool(re.search(r'\|[\s\S]*?\|[\s\S]*?\|', chunk_text))
                has_image = bool(re.search(r'\[Image[^\]]*\]', chunk_text, re.I))
                has_list = bool(re.search(r'^[\s]*[-*\d+\.]+\s', chunk_text, re.M))
                has_code = bool(re.search(r'```', chunk_text))
                
                hierarchy, section = get_heading_hierarchy(text_position)
                
                chunk = Chunk(
                    document_id=document_id,
                    chunk_id=f"{document_id}_chunk_{chunk_number:04d}",
                    chunk_number=chunk_number,
                    text=chunk_text,
                    char_count=len(chunk_text),
                    word_count=len(chunk_text.split()),
                    page_number=page_num,
                    section_title=section,
                    heading_hierarchy=hierarchy if hierarchy else None,
                    has_table=has_table,
                    has_image=has_image,
                    has_list=has_list,
                    has_code=has_code,
                    document_title=document_title,
                    document_author=document_author,
                    document_type=document_type,
                    metadata=metadata
                )
                chunks.append(chunk)
                all_chunk_texts.append(chunk_text)
                chunk_number += 1
    
    # Add adjacent chunk previews
    for i, chunk in enumerate(chunks):
        if i > 0:
            prev_text = all_chunk_texts[i - 1]
            chunk.prev_chunk_preview = prev_text[:150] + "..." if len(prev_text) > 150 else prev_text
        if i < len(chunks) - 1:
            next_text = all_chunk_texts[i + 1]
            chunk.next_chunk_preview = next_text[:150] + "..." if len(next_text) > 150 else next_text
    
    # Layer 2: Enrich chunks with MECH-specific metadata
    for chunk in chunks:
        mech_meta = extract_mech_metadata_from_chunk(chunk.text)
        
        # Merge MECH metadata into chunk metadata
        if chunk.metadata is None:
            chunk.metadata = {}
        
        chunk.metadata.update({
            'mech_domain': mech_meta['domain'],
            'bdd_tables': mech_meta['bdd_tables'],
            'field_names': mech_meta['field_names'],
            'shf_segment': mech_meta['shf_segment'],
            'xml_tags': mech_meta['xml_tags'],
            'program_chain': mech_meta['program_chain'],
            'fsf_fields': mech_meta['fsf_fields'],
            'contains_txnlog': mech_meta['contains_txnlog'],
            'contains_rl03': mech_meta['contains_rl03'],
            'contains_xxa': mech_meta['contains_xxa'],
            'contains_hex': mech_meta['contains_hex'],
            'contains_jcl': mech_meta['contains_jcl'],
            'contains_macro': mech_meta['contains_macro'],
            'is_mapping': mech_meta['is_mapping'],
            'is_procedure': mech_meta['is_procedure'],
            'identifier_count': mech_meta.get('identifier_count', 0),
            'identifier_density': mech_meta.get('identifier_density', 0.0),
        })
    
    logger.info(f"Created {len(chunks)} enhanced chunks from document {document_id}")
    return chunks


# ============================================================================
# S3 OPERATIONS
# ============================================================================

def download_from_s3(bucket: str, key: str) -> str:
    """Download file from S3 to temp directory."""
    try:
        s3 = get_client('s3')
        
        # Create temp file with proper extension
        ext = os.path.splitext(key)[1]
        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
            s3.download_fileobj(bucket, key, tmp)
            return tmp.name
            
    except Exception as e:
        raise ExtractionError(f"Failed to download from S3: {str(e)}", "s3")


def save_chunks_to_s3(chunks: List[Chunk], document_id: str, file_name: str = None) -> str:
    """
    Save chunks as JSONL to S3.
    
    PHASE 1.5: Enriches chunks with LLM metadata before saving.
    """
    try:
        s3 = get_client('s3')
        s3_key = f"chunks/{document_id}.jsonl"
        
        # PHASE 1.5: Enrich chunks with metadata
        if ENABLE_LLM_METADATA:
            chunks = enrich_chunks_with_metadata(chunks, file_name)
        
        # Convert chunks to JSONL
        jsonl_lines = []
        for chunk in chunks:
            chunk_dict = asdict(chunk)
            jsonl_lines.append(json.dumps(chunk_dict))
        
        content = '\n'.join(jsonl_lines)
        
        s3.put_object(
            Bucket=BUCKET_NAME,
            Key=s3_key,
            Body=content.encode('utf-8'),
            ContentType='application/jsonl',
            Metadata={
                'document_id': document_id,
                'chunk_count': str(len(chunks)),
                'metadata_enriched': str(ENABLE_LLM_METADATA).lower()
            }
        )
        
        logger.info(f"Saved {len(chunks)} chunks to s3://{BUCKET_NAME}/{s3_key}")
        return s3_key
        
    except Exception as e:
        raise ExtractionError(f"Failed to save chunks to S3: {str(e)}", "s3")


def move_to_processed(source_bucket: str, source_key: str, document_id: str) -> str:
    """Move processed file to processed folder."""
    try:
        s3 = get_client('s3')
        
        # Determine destination key
        file_name = os.path.basename(source_key)
        dest_key = f"processed/{document_id}/{file_name}"
        
        # Copy
        s3.copy_object(
            Bucket=BUCKET_NAME,
            Key=dest_key,
            CopySource={'Bucket': source_bucket, 'Key': source_key}
        )
        
        # Delete original (optional - comment out to keep original)
        # s3.delete_object(Bucket=source_bucket, Key=source_key)
        
        logger.info(f"Moved to s3://{BUCKET_NAME}/{dest_key}")
        return dest_key
        
    except Exception as e:
        logger.warning(f"Failed to move file to processed: {str(e)}")
        return source_key


# ============================================================================
# DYNAMODB OPERATIONS
# ============================================================================

def update_pipeline_status(document_id: str, status: str, message: str = '', 
                           metadata: Dict = None):
    """Update pipeline status in DynamoDB."""
    try:
        table = get_resource('dynamodb').Table(PIPELINE_TABLE)
        
        update_expr = 'SET #s = :status, stage_message = :msg, updated_at = :ts'
        expr_names = {'#s': 'status'}
        expr_values = {
            ':status': status,
            ':msg': message,
            ':ts': datetime.utcnow().isoformat()
        }
        
        if metadata:
            for key, value in metadata.items():
                safe_key = key.replace('-', '_')
                update_expr += f', {safe_key} = :{safe_key}'
                expr_values[f':{safe_key}'] = value
        
        table.update_item(
            Key={'document_id': document_id},
            UpdateExpression=update_expr,
            ExpressionAttributeNames=expr_names,
            ExpressionAttributeValues=expr_values
        )
        
        logger.info(f"Updated pipeline status: {document_id} -> {status}")
        
    except Exception as e:
        logger.error(f"Failed to update pipeline status: {str(e)}")


# ============================================================================
# LAMBDA HANDLER
# ============================================================================

def lambda_handler(event: Dict[str, Any], context: Any) -> Dict[str, Any]:
    """
    Main Lambda handler for document extraction.
    
    Input from Step Functions:
    {
        "document_id": "doc-abc123",
        "s3_bucket": "mechavatar-lab-cac1-s3-inbound",
        "s3_key": "documents/file.pdf",
        "file_type": "pdf",
        "file_name": "file.pdf"
    }
    
    Output:
    {
        "statusCode": 200,
        "document_id": "doc-abc123",
        "chunks_s3_key": "chunks/doc-abc123.jsonl",
        "chunk_count": 42,
        "extraction_method": "pymupdf4llm",
        "metrics": {...}
    }
    """
    import time
    start_time = time.time()
    temp_file = None
    request_id = context.aws_request_id if context else 'local'
    set_correlation_id(request_id)
    log_event('request_start', 'Document extraction invoked', request_id=request_id)
    
    try:
        event = DocumentExtractEvent.model_validate(event).model_dump()

        # Extract parameters
        document_id = event.get('document_id')
        s3_bucket = event.get('s3_bucket', BUCKET_NAME)
        s3_key = event.get('s3_key')
        file_type = event.get('file_type')
        file_name = event.get('file_name', os.path.basename(s3_key))
        
        log_event('parameters_parsed', 'Extraction parameters parsed',
              document_id=document_id, file_type=file_type, s3_key=s3_key)
        
        if not all([document_id, s3_key, file_type]):
            log_event('validation_error', 'Missing required parameters')
            raise ValueError("Missing required parameters: document_id, s3_key, file_type")

        current_status = get_pipeline_status(document_id)
        completed_statuses = {'extracted', 'embeddings_complete', 'indexed', 'partially_indexed', 'completed'}
        if current_status in completed_statuses:
            log_event('idempotent_skip', 'Extraction already completed',
                      document_id=document_id, status=current_status)
            return {
                'statusCode': 200,
                'document_id': document_id,
                'status': current_status,
                'message': 'Extraction already completed'
            }
        
        # Update status
        update_pipeline_status(document_id, 'extracting', f'Downloading {file_name}')
        
        # Download file
        download_start = time.time()
        temp_file = download_from_s3(s3_bucket, s3_key)
        download_time = time.time() - download_start
        logger.info(f"Downloaded to: {temp_file} ({download_time:.2f}s)")
        
        # Update status
        update_pipeline_status(document_id, 'extracting', f'Extracting text from {file_type}')
        
        # Extract text
        extract_start = time.time()
        result = extract_document(temp_file, file_type)
        extract_time = time.time() - extract_start
        
        if not result.success:
            raise ExtractionError(result.error or "Extraction failed", file_type)
        
        logger.info(f"Extraction complete: {result.char_count} chars, method={result.method}, time={extract_time:.2f}s")
        
        # Create enhanced chunks with document context
        update_pipeline_status(document_id, 'chunking', 'Creating semantic chunks')
        
        # Extract document metadata for chunk context
        doc_title = None
        doc_author = None
        if result.metadata:
            doc_title = result.metadata.get('title')
            doc_author = result.metadata.get('author')
        
        # Attempt structured extraction for advanced chunking
        structured = extract_structured_elements(temp_file, file_type, document_id)
        text_sample = result.text[:5000]
        if structured and structured.get('elements'):
            text_sample = _build_text_sample_from_elements(structured['elements'])

        # Extract document-level metadata DYNAMICALLY using LLM (no hardcoded categories)
        doc_metadata = extract_document_metadata(text_sample, file_name)

        chunk_metadata = {
            'file_name': file_name,
            'file_type': file_type,
            'extraction_method': result.method,
            'pages': result.pages,
            'tables': result.tables,
            'images': result.images,
            # DYNAMIC metadata for retrieval (no hardcoded categories)
            'doc_summary': doc_metadata.get('doc_summary', ''),
            'primary_topics': doc_metadata.get('primary_topics', []),
            'key_terms': doc_metadata.get('key_terms', []),
            'doc_type_inferred': doc_metadata.get('doc_type_inferred', 'unknown'),
            'domain_indicators': doc_metadata.get('domain_indicators', [])
        }

        if structured and structured.get('elements'):
            chunks = create_chunks_from_elements(
                structured['elements'],
                document_id,
                metadata=chunk_metadata,
                document_title=doc_title,
                document_author=doc_author,
                document_type=file_type
            )
        else:
            chunks = []

        # Fallback to legacy chunking if advanced chunking is unavailable
        if not chunks:
            chunks = create_chunks(
                result.text,
                document_id,
                metadata=chunk_metadata,
                document_title=doc_title,
                document_author=doc_author,
                document_type=file_type
            )
        
        if not chunks:
            raise ExtractionError("No chunks created - document may be empty", file_type)
        
        # Save chunks to S3 (with metadata enrichment)
        chunks_s3_key = save_chunks_to_s3(chunks, document_id, file_name)
        
        # Move original to processed folder
        move_to_processed(s3_bucket, s3_key, document_id)
        
        # Calculate metrics
        processing_time_ms = int((time.time() - start_time) * 1000)
        
        # Update final status
        update_pipeline_status(
            document_id, 
            'extracted',
            f'Extracted {len(chunks)} chunks',
            metadata={
                'chunk_count': len(chunks),
                'extraction_method': result.method,
                'processing_ms': processing_time_ms
            }
        )
        
        log_event('request_complete', 'Document extraction complete',
                  document_id=document_id, chunk_count=len(chunks), duration_ms=processing_time_ms)
        return {
            'statusCode': 200,
            'document_id': document_id,
            'chunks_s3_key': chunks_s3_key,
            'chunk_count': len(chunks),
            'file_name': file_name,
            'file_type': file_type,
            'extraction': {
                'method': result.method,
                'pages': result.pages,
                'tables': result.tables,
                'images': result.images,
                'sheets': result.sheets,
                'slides': result.slides,
                'word_count': result.word_count,
                'char_count': result.char_count
            },
            'metrics': {
                'processing_ms': processing_time_ms
            }
        }
        
    except ValidationError as e:
        log_event('request_error', 'Invalid input payload', errors=e.errors())
        return {
            'statusCode': 400,
            'error': 'Invalid input payload',
            'error_type': 'validation'
        }

    except (ClientError, BotoCoreError) as e:
        log_event('request_error', f"AWS service error: {str(e)}", document_id=event.get('document_id'))
        return {
            'statusCode': 502,
            'error': str(e),
            'error_type': 'aws_error',
            'document_id': event.get('document_id') if isinstance(event, dict) else None
        }

    except ExtractionError as e:
        log_event('request_error', f"Extraction error: {str(e)}", document_id=event.get('document_id'))
        logger.error(f"Extraction error: {str(e)}")
        
        if 'document_id' in event:
            update_pipeline_status(
                event['document_id'], 
                'failed',
                str(e),
                metadata={'error_type': 'extraction', 'error_method': e.method}
            )
        
        return {
            'statusCode': 500,
            'error': str(e),
            'error_type': 'extraction',
            'document_id': event.get('document_id')
        }
        
    except Exception as e:
        log_event('request_error', f"Unexpected error: {str(e)}", document_id=event.get('document_id'))
        logger.error(f"Unexpected error: {str(e)}", exc_info=True)
        
        if 'document_id' in event:
            update_pipeline_status(event['document_id'], 'failed', str(e))
        
        return {
            'statusCode': 500,
            'error': str(e),
            'error_type': 'unexpected',
            'document_id': event.get('document_id')
        }
        
    finally:
        # Cleanup temp file
        if temp_file and os.path.exists(temp_file):
            try:
                os.unlink(temp_file)
            except Exception:
                pass
